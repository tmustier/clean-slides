"""Inspect-related CLI commands."""

from __future__ import annotations

from collections.abc import Iterable
from typing import Protocol, cast

from pptx.presentation import Presentation as PresentationObj
from pptx.shapes.base import BaseShape
from pptx.slide import Slide
from pptx.util import Emu

from .cli_common import find_shape, get_slide, is_text_shape, json_out, open_presentation
from .constants import Layout


class FileArgs(Protocol):
    file: str


class ShowArgs(FileArgs, Protocol):
    slide: str | None
    shape: str | None


class SlideArgs(FileArgs, Protocol):
    slide: str


class ColorArgs(FileArgs, Protocol):
    hex: str


class XmlArgs(FileArgs, Protocol):
    slide: str
    shape: str


def _classify_shape_type(shape: BaseShape) -> str:
    """Classify shape for display: placeholder, chart, image, table, group, text, connector, decorative."""
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    if shape.is_placeholder:
        return "placeholder"

    st = shape.shape_type
    if st == MSO_SHAPE_TYPE.CHART:
        return "chart"
    if st == MSO_SHAPE_TYPE.PICTURE:
        return "image"
    if st == MSO_SHAPE_TYPE.TABLE:
        return "table"
    if st == MSO_SHAPE_TYPE.GROUP:
        return "group"
    if st in (MSO_SHAPE_TYPE.LINE, MSO_SHAPE_TYPE.FREEFORM):
        return "line"
    if is_text_shape(shape) and shape.text_frame.text.strip():
        return "text"
    return "shape"


def _shape_word_count(shape: BaseShape) -> int:
    """Count words in a shape's text frame."""
    if not is_text_shape(shape):
        return 0
    text = shape.text_frame.text.strip()
    return len(text.split()) if text else 0


def _slide_word_count(slide: Slide) -> int:
    """Total words across all text-bearing shapes on a slide."""
    total = 0
    for shape in slide.shapes:
        total += _shape_word_count(shape)
    return total


def _show_file_summary(prs: PresentationObj, file_path: str) -> None:
    from .inspect_pptx import get_slide_comments, list_slides

    total = len(prs.slides)
    print(f"\n  {file_path}  ({total} slide{'s' if total != 1 else ''})\n")
    for entry in list_slides(prs):
        slide = prs.slides[entry["slide"] - 1]
        layout_name = slide.slide_layout.name
        shapes_n = len(slide.shapes)
        words = _slide_word_count(slide)
        title = entry["title"][:60] or "(no title)"
        comments = get_slide_comments(slide)
        cm_str = f"  {len(comments)}cm" if comments else ""
        print(
            f"  {entry['slide']:3d}  {title:62s}  [{layout_name}]  {shapes_n}sh  {words}w{cm_str}"
        )
    print()


def _show_slide_header(slide: Slide, slide_label: str) -> None:
    title = "(no title)"
    subtitle = ""

    title_shape = slide.shapes.title
    if title_shape is not None:
        title = title_shape.text.replace("\x0b", " | ").replace("\n", " | ")

    for shape in slide.shapes:
        if is_text_shape(shape) and shape.is_placeholder and shape.placeholder_format.idx == 1:
            subtitle = shape.text.replace("\x0b", " | ").replace("\n", " | ")
            break

    layout = slide.slide_layout
    layout_phs = [f"ph{ph.placeholder_format.idx}:{ph.name}" for ph in layout.placeholders]

    print(f"\n  Slide {slide_label}: {title}")
    if subtitle:
        print(f"  {subtitle}")
    print(
        f"  Layout: \"{layout.name}\" → {', '.join(layout_phs) if layout_phs else '(no placeholders)'}"
    )
    print()


def _show_slide_shapes(slide: Slide) -> None:
    from .inspect_pptx import inspect_slide

    shapes = inspect_slide(slide)
    for item in shapes:
        real_shape = list(slide.shapes)[item.index]
        cat = _classify_shape_type(real_shape)

        ph_str = f"ph{item.placeholder_idx}" if item.placeholder_idx is not None else ""
        type_str = ph_str if ph_str else cat

        text = ""
        if item.text_preview:
            preview = item.text_preview[:60]
            if len(item.text_preview) > 60:
                preview += "…"
            text = f"  «{preview}»"

        print(
            f"  [{item.index:2d}] {type_str:8s}  {item.left:5.2f},{item.top:5.2f}  {item.width:5.2f}x{item.height:5.2f}  {item.name:30s}{text}"
        )


def _show_slide_comments(slide: Slide) -> None:
    from .inspect_pptx import get_slide_comments

    comments = get_slide_comments(slide)
    if not comments:
        return

    print(f"\n  Comments ({len(comments)}):")
    for comment in comments:
        preview = comment.text[:120]
        if len(comment.text) > 120:
            preview += "…"
        print(f"    [{comment.author}] {preview}")


def _show_shape_detail_json(slide: Slide, shape_identifier: str) -> None:
    from .inspect_pptx import inspect_chart, inspect_shape

    shape = find_shape(slide, shape_identifier)
    if shape.has_chart if hasattr(shape, "has_chart") else False:
        json_out(inspect_chart(shape))
    else:
        json_out(inspect_shape(shape))


def cmd_show(args: ShowArgs) -> int:
    """Progressive drill-down: file → slide list, slide → shapes, shape → detail."""
    prs = open_presentation(args.file)

    if args.slide is None:
        _show_file_summary(prs, args.file)
        return 0

    slide = get_slide(prs, args.slide)
    if args.shape is None:
        _show_slide_header(slide, args.slide)
        _show_slide_shapes(slide)
        _show_slide_comments(slide)
        print()
        return 0

    _show_shape_detail_json(slide, args.shape)
    return 0


def cmd_list(args: FileArgs) -> int:
    from .inspect_pptx import list_slides

    prs = open_presentation(args.file)
    for entry in list_slides(prs):
        print(f"  {entry['slide']:3d}  {entry['title']}")
    return 0


def cmd_summary(args: SlideArgs) -> int:
    from .inspect_pptx import summarize_slide

    prs = open_presentation(args.file)
    slide = get_slide(prs, args.slide)
    json_out(summarize_slide(slide))
    return 0


def cmd_slide(args: SlideArgs) -> int:
    from .inspect_pptx import inspect_slide

    prs = open_presentation(args.file)
    slide = get_slide(prs, args.slide)
    shapes = inspect_slide(slide)
    for s in shapes:
        ph = f" ph={s.placeholder_idx}" if s.placeholder_idx is not None else ""
        fill = f" fill={s.fill.type}" if s.fill.type != "inherited" else ""
        text = f"  «{s.text_preview}»" if s.text_preview else ""
        print(
            f"  [{s.index:2d}] {s.name:30s}  {s.left:6.2f},{s.top:6.2f}  {s.width:5.2f}x{s.height:5.2f}{ph}{fill}{text}"
        )
    return 0


def cmd_shape(args: XmlArgs) -> int:
    from .inspect_pptx import inspect_shape

    prs = open_presentation(args.file)
    slide = get_slide(prs, args.slide)
    shape = find_shape(slide, args.shape)
    json_out(inspect_shape(shape))
    return 0


def cmd_chart(args: XmlArgs) -> int:
    from .inspect_pptx import inspect_chart

    prs = open_presentation(args.file)
    slide = get_slide(prs, args.slide)
    shape = find_shape(slide, args.shape)
    json_out(inspect_chart(shape))
    return 0


def cmd_layout(args: SlideArgs) -> int:
    from .inspect_pptx import inspect_layout

    prs = open_presentation(args.file)
    slide = get_slide(prs, args.slide)
    json_out(inspect_layout(slide.slide_layout))
    return 0


def cmd_layouts(args: FileArgs) -> int:
    from pptx.enum.shapes import PP_PLACEHOLDER
    from pptx.shapes.placeholder import LayoutPlaceholder

    prs = open_presentation(args.file)

    # Boundary between "structural" placeholders (title, subtitle, tracker) and
    # the slide's main content zone.
    content_y_threshold = int(Layout.CONTENT_START_Y)
    footer_y = int(Layout.FOOTER_LINE_Y)

    for layout in prs.slide_layouts:
        structural: list[str] = []
        # (left_emu, width_in, height_avail_in, name)
        content_phs: list[tuple[int, float, float, str]] = []

        placeholders = cast(Iterable[LayoutPlaceholder], layout.placeholders)
        phs: list[LayoutPlaceholder] = sorted(
            placeholders,
            key=lambda p: (int(p.top), int(p.left)),
        )
        for ph in phs:
            pf = ph.placeholder_format
            w_in = ph.width.inches

            if pf.type == PP_PLACEHOLDER.TITLE:
                structural.append(f"title ({w_in:.1f}in)")
            elif pf.type == PP_PLACEHOLDER.SUBTITLE:
                structural.append(f"subtitle ({w_in:.1f}in)")
            elif pf.type == PP_PLACEHOLDER.PICTURE:
                structural.append(f"image: {ph.name} ({w_in:.1f}×{ph.height.inches:.1f}in)")
            elif int(ph.top) < content_y_threshold:
                # Above content zone — tracker, doc type, etc.
                structural.append(f"{ph.name} ({w_in:.1f}in)")
            else:
                h_avail = Emu(footer_y - int(ph.top)).inches
                content_phs.append((int(ph.left), w_in, h_avail, ph.name))

        # Sort content areas left-to-right, label primary/secondary
        content_phs.sort(key=lambda p: p[0])
        areas: list[str] = []
        for idx, (_x, w, h, _name) in enumerate(content_phs):
            label = "primary" if idx == 0 else "secondary"
            areas.append(f"{label} {w:.1f}×{h:.1f}in")

        print(f"  {layout.name}")
        if structural:
            print(f"    placeholders: {', '.join(structural)}")
        if areas:
            print(f"    content areas: {', '.join(areas)}")
        else:
            print("    content areas: (none)")
    return 0


def cmd_theme(args: FileArgs) -> int:
    from .inspect_pptx import resolve_theme_colors

    prs = open_presentation(args.file)
    colors = resolve_theme_colors(prs)
    for name, hex_val in sorted(colors.items()):
        if not name.startswith("_"):
            print(f"  {name:12s}  {hex_val}")
    return 0


def cmd_color(args: ColorArgs) -> int:
    from .inspect_pptx import identify_color

    prs = open_presentation(args.file)
    result = identify_color(prs, args.hex)
    if result:
        print(result)
    else:
        print(f"No theme match for {args.hex}")
    return 0


def cmd_xml(args: XmlArgs) -> int:
    from .xml_helpers import dump_xml

    prs = open_presentation(args.file)
    slide = get_slide(prs, args.slide)
    shape = find_shape(slide, args.shape)
    print(dump_xml(shape._element))
    return 0
