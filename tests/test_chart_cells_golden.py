from __future__ import annotations

import zipfile
from dataclasses import dataclass
from pathlib import Path
from typing import Any

from lxml import etree
from pptx import Presentation
from pptx.enum.chart import XL_CHART_TYPE

from clean_slides.cli import cmd_generate
from clean_slides.pptx_access import (
    iter_shapes,
    presentation_chart_types,
    shape_has_connector_endpoints,
    shape_is_placeholder,
    shape_text_frame_text,
)


@dataclass
class _GenerateArgs:
    input: list[str]
    template: str | None
    slide_index: int | None
    keep_existing: bool
    detail: bool
    output: str | None
    config: str | None = None


_NS = {
    "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
    "c": "http://schemas.openxmlformats.org/drawingml/2006/chart",
}


def _run_generate(tmp_path: Path, name: str, yaml_text: str) -> Path:
    yaml_path = tmp_path / f"{name}.yaml"
    yaml_path.write_text(yaml_text.strip(), encoding="utf-8")

    output_path = tmp_path / f"{name}.pptx"
    result = cmd_generate(
        _GenerateArgs(
            input=[str(yaml_path)],
            template=None,
            slide_index=None,
            keep_existing=False,
            detail=False,
            output=str(output_path),
        )
    )

    assert result == 0
    assert output_path.exists()
    return output_path


def _chart_xml(path: Path, chart_index: int = 1) -> Any:
    chart_part = f"ppt/charts/chart{chart_index}.xml"
    with zipfile.ZipFile(path) as archive:
        xml_blob = archive.read(chart_part)
    return etree.fromstring(xml_blob)


def _chart_types(path: Path) -> list[int]:
    prs = Presentation(str(path))
    return presentation_chart_types(prs)


def _non_placeholder_texts(path: Path) -> list[str]:
    prs = Presentation(str(path))
    texts: list[str] = []
    for shape in iter_shapes(prs.slides[0]):
        if shape_is_placeholder(shape):
            continue

        raw_text = shape_text_frame_text(shape)
        if raw_text is None:
            continue

        text = raw_text.strip()
        if text:
            texts.append(text)
    return texts


def _connector_count(path: Path) -> int:
    prs = Presentation(str(path))
    count = 0
    for shape in iter_shapes(prs.slides[0]):
        if shape_has_connector_endpoints(shape):
            count += 1
    return count


def test_horizontal_chart_cells_keep_manual_offsets_and_format(tmp_path: Path) -> None:
    """Regression: horizontal chart-cell labels keep gap + custom format."""
    yaml_text = """
title: Horizontal label regression
charts:
  rev:
    dir: horizontal
    values: [100, 200]
    format: "€{}m"
    color: accent1

table:
  rows: 2
  cols: 1
  has_col_header: false
  cells:
    - [rev-1]
    - [rev-2]
"""

    output_path = _run_generate(tmp_path, "horizontal-label-regression", yaml_text)

    assert _chart_types(output_path) == [int(XL_CHART_TYPE.BAR_CLUSTERED)]

    chart_xml = _chart_xml(output_path)
    labels = chart_xml.xpath(".//c:barChart/c:ser/c:dLbls/c:dLbl", namespaces=_NS)
    assert len(labels) == 2

    x_offsets: list[float] = []
    for label in labels:
        assert label.xpath("./c:dLblPos/@val", namespaces=_NS) == ["ctr"]
        assert label.xpath("./c:numFmt/@formatCode", namespaces=_NS) == ['"€"0"m"']
        assert label.xpath("./c:txPr/a:bodyPr/@wrap", namespaces=_NS) == ["none"]

        x_values = label.xpath("./c:layout/c:manualLayout/c:x/@val", namespaces=_NS)
        assert len(x_values) == 1
        x_offset = float(x_values[0])
        assert x_offset > 0
        x_offsets.append(x_offset)

    assert x_offsets[1] > x_offsets[0]


def test_horizontal_chart_cells_zero_decimal_python_format_maps_to_excel_numeric(
    tmp_path: Path,
) -> None:
    """Regression: ``{:.0f}`` stays integer in chart-cell data labels."""
    yaml_text = """
title: Horizontal zero-decimal format regression
charts:
  rev:
    dir: horizontal
    values: [1.2, 2.8]
    format: "{:.0f}"
    color: accent1

table:
  rows: 2
  cols: 1
  has_col_header: false
  cells:
    - [rev-1]
    - [rev-2]
"""

    output_path = _run_generate(tmp_path, "horizontal-zero-decimal-regression", yaml_text)

    chart_xml = _chart_xml(output_path)
    num_formats = chart_xml.xpath(
        ".//c:barChart/c:ser/c:dLbls/c:dLbl/c:numFmt/@formatCode",
        namespaces=_NS,
    )
    assert num_formats == ["0", "0"]


def test_waterfall_chart_cells_keep_connectors_and_formatted_labels(tmp_path: Path) -> None:
    """Regression: waterfall chart-cells keep overlays + connector lines."""
    yaml_text = """
title: Waterfall connector regression
charts:
  wf:
    type: waterfall
    dir: horizontal
    values: [954, 13, -45, 1209]
    totals: [1, 4]
    decreases: [3]
    format: "{:,.0f}"
    connector: true

table:
  rows: 4
  cols: 1
  has_col_header: false
  cells:
    - [wf-1]
    - [wf-2]
    - [wf-3]
    - [wf-4]
"""

    output_path = _run_generate(tmp_path, "waterfall-regression", yaml_text)

    assert _chart_types(output_path) == [int(XL_CHART_TYPE.BAR_STACKED)]

    texts = _non_placeholder_texts(output_path)
    expected_labels = {"954", "13", "-45", "1,209"}
    assert expected_labels.issubset(set(texts))

    # Category labels are intentionally suppressed for chart-cells.
    assert {"1", "2", "3", "4"}.isdisjoint(set(texts))

    assert _connector_count(output_path) >= 3
