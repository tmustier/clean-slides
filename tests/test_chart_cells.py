"""Tests for chart cell parsing, validation, sizing, and rendering."""

# pyright: reportUnknownMemberType=false
# pyright: reportAttributeAccessIssue=false
# pyright: reportUnknownVariableType=false
# pyright: reportPrivateUsage=false

from __future__ import annotations

import unittest
from dataclasses import dataclass
from pathlib import Path

import pytest
from pptx import Presentation
from pptx.util import Emu

from clean_slides.chart_render import (
    ChartGroup,
    _chart_def_to_spec,
    _python_fmt_to_excel_format,
    _waterfall_overlay_label_texts,
)
from clean_slides.cli import cmd_generate, cmd_validate
from clean_slides.constants import Fonts, TableDefaults
from clean_slides.sizing import ColumnSizer, FontConfig, RowSizer
from clean_slides.spec import (
    ChartDef,
    ChartRef,
    ContentArea,
    TableSpec,
    is_chart_ref,
    parse_chart_ref,
)
from clean_slides.text_metrics import TextMetrics

# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------


@dataclass
class GenerateArgs:
    input: list[str]
    template: str | None
    slide_index: int | None
    keep_existing: bool
    detail: bool
    output: str | None
    config: str | None = None


@dataclass
class ValidateArgs:
    input: list[str]
    detail: bool
    json: str | None = None
    config: str | None = None


FONTS = FontConfig(
    body_font=Fonts.BODY,
    body_size_pt=12,
    header_font=Fonts.HEADLINE,
    header_size_pt=12,
)

PAD = int(TableDefaults.CELL_PADDING)


# ===================================================================
# Phase 1: Parsing
# ===================================================================


class TestParseChartRef(unittest.TestCase):
    """Test is_chart_ref() and parse_chart_ref()."""

    def test_simple_ref(self) -> None:
        assert is_chart_ref("revenue-1")
        ref = parse_chart_ref("revenue-1")
        assert ref is not None
        assert ref.name == "revenue"
        assert ref.index == 1

    def test_multi_word_name(self) -> None:
        ref = parse_chart_ref("my_chart-3")
        assert ref is not None
        assert ref.name == "my_chart"
        assert ref.index == 3

    def test_not_a_ref(self) -> None:
        assert not is_chart_ref("plain text")
        assert parse_chart_ref("plain text") is None

    def test_no_index(self) -> None:
        assert not is_chart_ref("revenue-")
        assert not is_chart_ref("revenue")

    def test_negative_index(self) -> None:
        assert not is_chart_ref("revenue--1")

    def test_zero_index(self) -> None:
        # 0 is not a valid 1-based index but parse_chart_ref should still parse it;
        # validation catches it later.
        ref = parse_chart_ref("revenue-0")
        assert ref is not None
        assert ref.index == 0


class TestParseCharts(unittest.TestCase):
    """Test YAML parsing of charts: key and cell grid replacement."""

    def test_from_dict_with_charts(self) -> None:
        data: dict[str, object] = {
            "table": {
                "rows": 2,
                "cols": 2,
                "has_col_header": True,
                "col_headers": ["A", "B"],
                "cells": [
                    ["rev-1", "cost-1"],
                    ["rev-2", "cost-2"],
                ],
            },
            "charts": {
                "rev": {
                    "dir": "vertical",
                    "values": [10, 20],
                },
                "cost": {
                    "dir": "vertical",
                    "values": [5, 8],
                },
            },
        }
        spec = TableSpec.from_dict(data)

        # chart_defs should be populated
        assert spec.chart_defs is not None
        assert "rev" in spec.chart_defs
        assert "cost" in spec.chart_defs

        # Cells should contain ChartRef objects
        assert spec.cells is not None
        assert isinstance(spec.cells[0][0], ChartRef)
        assert spec.cells[0][0].name == "rev"
        assert spec.cells[0][0].index == 1

    def test_from_dict_without_charts(self) -> None:
        data: dict[str, object] = {
            "table": {
                "rows": 2,
                "cols": 1,
                "has_col_header": False,
                "cells": [["hello"], ["world"]],
            },
        }
        spec = TableSpec.from_dict(data)
        assert spec.chart_defs is None or len(spec.chart_defs) == 0

    def test_chart_def_defaults(self) -> None:
        data: dict[str, object] = {
            "table": {
                "rows": 2,
                "cols": 1,
                "has_col_header": False,
                "cells": [["ch-1"], ["text"]],
            },
            "charts": {
                "ch": {
                    "values": [42],
                },
            },
        }
        spec = TableSpec.from_dict(data)
        assert spec.chart_defs is not None
        cd = spec.chart_defs["ch"]
        assert cd.dir == "vertical"  # default
        assert cd.color is None
        assert cd.colors is None
        assert cd.format == "{}"  # default pass-through format
        assert cd.scale_max is None
        assert cd.label_position is None

    def test_chart_def_point_colors(self) -> None:
        data: dict[str, object] = {
            "table": {
                "rows": 3,
                "cols": 1,
                "has_col_header": False,
                "cells": [["wf-1"], ["wf-2"], ["wf-3"]],
            },
            "charts": {
                "wf": {
                    "type": "waterfall",
                    "dir": "horizontal",
                    "values": [954, 13, 1209],
                    "totals": [1, 3],
                    "color": "#4472C4",
                    "total_color": "#0D193B",
                    "colors": [None, "#4472C4", None],
                }
            },
        }

        spec = TableSpec.from_dict(data)
        assert spec.chart_defs is not None
        chart = spec.chart_defs["wf"]
        assert chart.colors == [None, "#4472C4", None]

    def test_chart_def_point_colors_must_match_values_length(self) -> None:
        data: dict[str, object] = {
            "table": {
                "rows": 2,
                "cols": 1,
                "has_col_header": False,
                "cells": [["ch-1"], ["ch-2"]],
            },
            "charts": {
                "ch": {
                    "values": [1, 2],
                    "colors": ["#4472C4"],
                }
            },
        }

        with pytest.raises(ValueError, match="colors must have"):
            TableSpec.from_dict(data)

    def test_grouped_singleton_empty_header_promotes_first_body_label(self) -> None:
        data: dict[str, object] = {
            "table": {
                "cols": 3,
                "has_col_header": False,
                "row_groups": [
                    {
                        "header": "",
                        "rows": [["FY26E EBITDAaL", "wf-1"]],
                    },
                    {
                        "header": "Drivers",
                        "rows": [["CPI", "wf-2"]],
                    },
                ],
            },
            "charts": {
                "wf": {
                    "type": "waterfall",
                    "dir": "horizontal",
                    "values": [954, 13],
                    "totals": [1],
                }
            },
        }

        spec = TableSpec.from_dict(data)
        assert spec.groups is not None
        assert spec.groups[0].header == "FY26E EBITDAaL"
        assert spec.groups[0].promoted is True
        assert spec.cells is not None
        assert spec.cells[0][0] == ""

    def test_grouped_singleton_empty_header_not_promoted_without_chart_ref(self) -> None:
        data: dict[str, object] = {
            "table": {
                "cols": 3,
                "has_col_header": False,
                "row_groups": [
                    {
                        "header": "",
                        "rows": [["Plain Label", "No chart"]],
                    }
                ],
            }
        }

        spec = TableSpec.from_dict(data)
        assert spec.groups is not None
        assert spec.groups[0].header == ""
        assert spec.groups[0].promoted is False
        assert spec.cells is not None
        assert spec.cells[0][0] == "Plain Label"

    def test_row_group_header_string_newline_parenthesis_maps_to_sub(self) -> None:
        data: dict[str, object] = {
            "table": {
                "cols": 3,
                "has_col_header": False,
                "row_groups": [
                    {
                        "header": "Impact of net site additions\n(at constant prices)",
                        "rows": [["MSA", "wf-1"]],
                    }
                ],
            },
            "charts": {
                "wf": {
                    "type": "waterfall",
                    "dir": "horizontal",
                    "values": [13],
                    "totals": [1],
                }
            },
        }

        spec = TableSpec.from_dict(data)
        assert spec.groups is not None
        header = spec.groups[0].header
        assert isinstance(header, dict)
        assert header.get("text") == "Impact of net site additions"
        assert header.get("sub") == "(at constant prices)"

    def test_promoted_header_string_newline_parenthesis_maps_to_sub(self) -> None:
        data: dict[str, object] = {
            "table": {
                "cols": 3,
                "has_col_header": False,
                "row_groups": [
                    {
                        "header": "",
                        "rows": [
                            ["FY33E EBITDAaL at FY26 mgn.\n(i.e. impact of Revenue growth)", "wf-1"]
                        ],
                    }
                ],
            },
            "charts": {
                "wf": {
                    "type": "waterfall",
                    "dir": "horizontal",
                    "values": [1166],
                    "totals": [1],
                }
            },
        }

        spec = TableSpec.from_dict(data)
        assert spec.groups is not None
        assert spec.groups[0].promoted is True
        header = spec.groups[0].header
        assert isinstance(header, dict)
        assert header.get("text") == "FY33E EBITDAaL at FY26 mgn."
        assert header.get("sub") == "(i.e. impact of Revenue growth)"


class TestChartFormatConversion(unittest.TestCase):
    def test_zero_decimal_format_preserved(self) -> None:
        excel = _python_fmt_to_excel_format("{:.0f}", [1.2, 2.8])
        assert excel == "0"


class TestWaterfallSpecAndLabels(unittest.TestCase):
    def test_waterfall_spec_uses_blank_categories_and_zero_based_totals(self) -> None:
        chart = ChartDef(
            name="wf",
            type="waterfall",
            dir="horizontal",
            values=[954, 13, 1209],
            totals=[1, 3],
        )
        group = ChartGroup(
            chart_def=chart,
            refs=[(0, 0, 1), (1, 0, 2), (2, 0, 3)],
            min_row=0,
            max_row=2,
            min_col=0,
            max_col=0,
        )

        spec = _chart_def_to_spec(group)
        assert spec["categories"] == ["", "", ""]

        wf = spec["waterfall"]
        assert wf["total_categories"] == [0, 2]
        assert wf["total_override"] is True
        assert isinstance(wf["connector_inset"], Emu)
        assert int(wf["connector_inset"]) == 3000
        assert wf["connector_overlap"] == 0

    def test_waterfall_overlay_label_texts_respect_python_format(self) -> None:
        meta: dict[str, object] = {
            "overlay": {
                "categories": ["", "", ""],
                "cumulative_totals": [954.0, 967.0, 1166.0],
                "delta_values": [None, 13.0, None],
                "total_categories": {0, 2},
            }
        }

        labels = _waterfall_overlay_label_texts(meta, "{:,.0f}")
        assert labels == ["954", "13", "1,166"]


# ===================================================================
# Phase 2: Validation
# ===================================================================


class TestValidation(unittest.TestCase):
    """Validate chart ref errors caught by cmd_validate."""

    def _write_yaml(self, path: Path, content: str) -> str:
        yaml_path = path / "test.yaml"
        yaml_path.write_text(content.strip())
        return str(yaml_path)

    def _validate(self, tmp_path: Path, yaml: str) -> int:
        yaml_path = self._write_yaml(tmp_path, yaml)
        args = ValidateArgs(input=[yaml_path], detail=False)
        return cmd_validate(args)

    def test_valid_vertical_chart_refs(self) -> None:
        """Vertical chart: refs in same row across columns."""
        import tempfile

        with tempfile.TemporaryDirectory() as td:
            p = Path(td)
            result = self._validate(
                p,
                """
table:
  rows: 2
  cols: 2
  has_col_header: true
  col_headers: ["A", "B"]
  cells:
    - [rev-1, rev-2]

charts:
  rev:
    dir: vertical
    values: [10, 20]
""",
            )
            assert result == 0

    def test_valid_horizontal_chart_refs(self) -> None:
        """Horizontal chart: refs in same column across rows."""
        import tempfile

        with tempfile.TemporaryDirectory() as td:
            p = Path(td)
            result = self._validate(
                p,
                """
table:
  rows: 3
  cols: 1
  has_col_header: true
  col_headers: ["Revenue"]
  cells:
    - [rev-1]
    - [rev-2]

charts:
  rev:
    dir: horizontal
    values: [100, 200]
""",
            )
            assert result == 0

    def test_unknown_chart_name(self) -> None:
        import tempfile

        with tempfile.TemporaryDirectory() as td:
            p = Path(td)
            result = self._validate(
                p,
                """
table:
  rows: 1
  cols: 1
  has_col_header: false
  cells:
    - [unknown-1]

charts:
  rev:
    values: [10]
""",
            )
            assert result != 0

    def test_index_out_of_bounds(self) -> None:
        import tempfile

        with tempfile.TemporaryDirectory() as td:
            p = Path(td)
            result = self._validate(
                p,
                """
table:
  rows: 1
  cols: 1
  has_col_header: false
  cells:
    - [rev-5]

charts:
  rev:
    values: [10]
""",
            )
            assert result != 0

    def test_duplicate_index(self) -> None:
        import tempfile

        with tempfile.TemporaryDirectory() as td:
            p = Path(td)
            result = self._validate(
                p,
                """
table:
  rows: 3
  cols: 1
  has_col_header: true
  col_headers: ["Revenue"]
  cells:
    - [rev-1]
    - [rev-1]

charts:
  rev:
    dir: horizontal
    values: [10, 20]
""",
            )
            assert result != 0


# ===================================================================
# Phase 3: Sizing
# ===================================================================


class TestSizingWithChartRefs(unittest.TestCase):
    """Chart ref cells should not affect text-based sizing."""

    def test_column_sizer_skips_chart_cells(self) -> None:
        spec = TableSpec(
            num_rows=2,
            num_cols=2,
            has_col_header=True,
            has_row_header=False,
            col_headers=["Text Col", "Chart Col"],
            cells=[
                ["Some text", ChartRef(name="rev", index=1)],
                ["More text", ChartRef(name="rev", index=2)],
            ],
            chart_defs={
                "rev": ChartDef(name="rev", type="bar", dir="vertical", values=[10, 20]),
            },
        )
        metrics = TextMetrics()
        area = ContentArea.from_layout("default")
        widths, _warnings = ColumnSizer().size(spec, area.width, metrics, FONTS, pad_top=PAD)

        assert len(widths) == 2
        assert sum(widths) <= area.width
        assert all(w > 0 for w in widths)

    def test_row_sizer_skips_chart_cells(self) -> None:
        spec = TableSpec(
            num_rows=2,
            num_cols=2,
            has_col_header=True,
            has_row_header=False,
            col_headers=["Text Col", "Chart Col"],
            cells=[
                ["Some text", ChartRef(name="rev", index=1)],
                ["More text", ChartRef(name="rev", index=2)],
            ],
            chart_defs={
                "rev": ChartDef(name="rev", type="bar", dir="vertical", values=[10, 20]),
            },
        )
        metrics = TextMetrics()
        area = ContentArea.from_layout("default")
        widths, _cw = ColumnSizer().size(spec, area.width, metrics, FONTS, pad_top=PAD)

        heights, _rw = RowSizer().size(spec, widths, area.height, metrics, FONTS, PAD, PAD)

        total_rows = spec.num_rows + 1  # +1 for header row
        assert len(heights) == total_rows
        assert sum(heights) == area.height
        assert all(h > 0 for h in heights)


# ===================================================================
# Phase 6: Integration — end-to-end generate
# ===================================================================


def _charts_module_available() -> bool:
    """Check if the bundled chart engine can be imported."""
    import importlib

    try:
        importlib.import_module("clean_slides.chart_generator")
    except Exception:
        return False
    return True


_skip_no_charts = unittest.skipUnless(
    _charts_module_available(),
    "Bundled chart engine module not available",
)


@_skip_no_charts
class TestChartCellsIntegration(unittest.TestCase):
    """Generate a slide with chart cells and verify shapes."""

    def test_generate_with_chart_cells(self) -> None:
        import tempfile

        with tempfile.TemporaryDirectory() as td:
            p = Path(td)
            yaml_path = p / "chart.yaml"
            yaml_path.write_text(
                """
title: Chart Test
subtitle: Integration

charts:
  rev:
    dir: horizontal
    values: [100, 200]
    format: "€{}m"
    color: accent1

table:
  rows: 2
  cols: 2
  has_col_header: true
  col_headers: ["Revenue", "Note"]
  cells:
    - [rev-1, "First year"]
    - [rev-2, "Second year"]
""".strip()
            )

            output_path = p / "output.pptx"
            args = GenerateArgs(
                input=[str(yaml_path)],
                template=None,
                slide_index=None,
                keep_existing=False,
                detail=False,
                output=str(output_path),
            )
            result = cmd_generate(args)
            assert result == 0
            assert output_path.exists()

            prs = Presentation(str(output_path))
            slide = prs.slides[0]

            chart_shapes = [s for s in slide.shapes if hasattr(s, "chart")]
            assert len(chart_shapes) >= 1, f"Expected chart shapes, found {len(chart_shapes)}"
            chart = chart_shapes[0].chart
            assert chart.has_legend is False

    def test_generate_vertical_charts(self) -> None:
        import tempfile

        with tempfile.TemporaryDirectory() as td:
            p = Path(td)
            yaml_path = p / "vchart.yaml"
            yaml_path.write_text(
                """
title: Vertical Charts

charts:
  ratio:
    dir: vertical
    values: [4.0, 3.1]
    format: "{}x"
    color: accent2

table:
  rows: 2
  cols: 2
  has_col_header: true
  col_headers: ["FY26E", "FY33E"]
  cells:
    - [ratio-1, ratio-2]
""".strip()
            )

            output_path = p / "output.pptx"
            args = GenerateArgs(
                input=[str(yaml_path)],
                template=None,
                slide_index=None,
                keep_existing=False,
                detail=False,
                output=str(output_path),
            )
            result = cmd_generate(args)
            assert result == 0

            prs = Presentation(str(output_path))
            slide = prs.slides[0]
            chart_shapes = [s for s in slide.shapes if hasattr(s, "chart")]
            assert len(chart_shapes) == 1


if __name__ == "__main__":
    unittest.main()
