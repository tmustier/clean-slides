from __future__ import annotations

from dataclasses import dataclass
from pathlib import Path

from pptx import Presentation
from pptx.util import Emu

from clean_slides.cli import cmd_generate, cmd_insert


@dataclass
class GenerateArgs:
    input: list[str]
    template: str | None
    slide_index: int | None
    keep_existing: bool
    detail: bool
    output: str | None
    config: str | None = None


@dataclass
class InsertArgs:
    file: str
    source: str
    at: int | None
    slides: str | None
    out: str | None


def _write_yaml(path: Path) -> None:
    path.write_text(
        """
layout: default

table:
  rows: 2
  cols: 2
  has_col_header: true
  col_headers: ["Header A", "Header B"]
  cells:
    - ["Alpha", "Beta"]
    - ["Gamma", "Delta"]
""".strip()
    )


def test_insert_smoke(tmp_path: Path) -> None:
    # 1) Generate a source deck with one slide
    yaml_path = tmp_path / "table.yaml"
    _write_yaml(yaml_path)

    src_path = tmp_path / "src.pptx"
    gen_args = GenerateArgs(
        input=[str(yaml_path)],
        template=None,
        slide_index=None,
        keep_existing=False,
        detail=False,
        output=str(src_path),
    )
    assert cmd_generate(gen_args) == 0

    src_prs = Presentation(str(src_path))
    assert len(src_prs.slides) == 1
    src_shapes = len(src_prs.slides[0].shapes)
    assert src_shapes > 0

    # 2) Create a target deck with two blank slides
    target_path = tmp_path / "target.pptx"
    target_prs = Presentation()
    blank_layout = target_prs.slide_layouts[len(target_prs.slide_layouts) - 1]
    target_prs.slides.add_slide(blank_layout)
    target_prs.slides.add_slide(blank_layout)
    target_prs.save(str(target_path))

    # 3) Insert source slide at position 2
    out_path = tmp_path / "out.pptx"
    ins_args = InsertArgs(
        file=str(target_path),
        source=str(src_path),
        at=2,
        slides="1",
        out=str(out_path),
    )
    assert cmd_insert(ins_args) == 0

    out_prs = Presentation(str(out_path))
    assert len(out_prs.slides) == 3

    # Inserted slide is position 2 (index 1)
    assert len(out_prs.slides[1].shapes) == src_shapes


def test_insert_preserves_layout(tmp_path: Path) -> None:
    """Insert should match the source slide's layout name in the destination."""
    # Create a source deck with a named layout
    src_prs = Presentation()
    # Use the first available layout and note its name
    src_layout = src_prs.slide_layouts[0]
    src_layout_name = src_layout.name
    src_slide = src_prs.slides.add_slide(src_layout)

    # Add a shape so the slide isn't empty
    txbox = src_slide.shapes.add_textbox(Emu(100000), Emu(100000), Emu(500000), Emu(200000))
    txbox.text_frame.text = "Test"

    src_path = tmp_path / "src.pptx"
    src_prs.save(str(src_path))

    # Create a destination deck from the same default template (has same layouts)
    dst_prs = Presentation()
    blank_layout = dst_prs.slide_layouts[len(dst_prs.slide_layouts) - 1]
    dst_prs.slides.add_slide(blank_layout)
    dst_path = tmp_path / "dst.pptx"
    dst_prs.save(str(dst_path))

    out_path = tmp_path / "out.pptx"
    ins_args = InsertArgs(
        file=str(dst_path),
        source=str(src_path),
        at=2,
        slides="1",
        out=str(out_path),
    )
    assert cmd_insert(ins_args) == 0

    out_prs = Presentation(str(out_path))
    inserted_slide = out_prs.slides[1]
    assert inserted_slide.slide_layout.name == src_layout_name


def test_insert_copies_hyperlinks(tmp_path: Path) -> None:
    """Insert should copy hyperlink relationships to the destination slide."""
    # Create a source slide with a hyperlink
    src_prs = Presentation()
    src_slide = src_prs.slides.add_slide(src_prs.slide_layouts[0])
    txbox = src_slide.shapes.add_textbox(Emu(100000), Emu(100000), Emu(500000), Emu(200000))
    p = txbox.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = "Click here"
    r.hyperlink.address = "https://example.com"

    src_path = tmp_path / "src.pptx"
    src_prs.save(str(src_path))

    # Create destination deck
    dst_prs = Presentation()
    dst_prs.slides.add_slide(dst_prs.slide_layouts[0])
    dst_path = tmp_path / "dst.pptx"
    dst_prs.save(str(dst_path))

    out_path = tmp_path / "out.pptx"
    ins_args = InsertArgs(
        file=str(dst_path),
        source=str(src_path),
        at=2,
        slides="1",
        out=str(out_path),
    )
    assert cmd_insert(ins_args) == 0

    # Verify the hyperlink relationship exists on the inserted slide
    out_prs = Presentation(str(out_path))
    inserted_slide = out_prs.slides[1]

    hyperlink_rels = [
        rel
        for rel in inserted_slide.part.rels.values()
        if rel.is_external
        and "hyperlink" in rel.reltype
    ]
    assert len(hyperlink_rels) == 1
    assert hyperlink_rels[0].target_ref == "https://example.com"
