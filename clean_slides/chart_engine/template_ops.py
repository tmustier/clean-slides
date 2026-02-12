"""Template and OPC-level chart replacement helpers."""

from __future__ import annotations

import hashlib
import posixpath
import xml.etree.ElementTree as ET
import zipfile
from collections.abc import Sequence
from pathlib import Path
from typing import Any, NamedTuple, cast

from pptx import Presentation

from ..pptx_access import slide_charts


def read_pptx_part(path: Path, part_name: str) -> bytes:
    """Read a specific part from a PPTX package."""
    try:
        with zipfile.ZipFile(path) as pptx:
            return pptx.read(part_name)
    except KeyError as exc:
        raise ValueError(f"Missing {part_name} in {path}") from exc


def read_content_type_defaults(path: Path) -> dict[str, str]:
    """Read default extension → content-type mappings from ``[Content_Types].xml``."""
    content_xml = read_pptx_part(path, "[Content_Types].xml")
    root = ET.fromstring(content_xml)
    ns = "{http://schemas.openxmlformats.org/package/2006/content-types}"

    defaults: dict[str, str] = {}
    for item in root.findall(f"{ns}Default"):
        ext = item.get("Extension")
        content_type = item.get("ContentType")
        if ext and content_type:
            defaults[ext] = content_type
    return defaults


def normalize_relationship_target(base_part: str, target: str) -> str:
    """Resolve a relationship target to a normalized package path."""
    base_dir = posixpath.dirname(base_part)
    return posixpath.normpath(posixpath.join(base_dir, target))


def replace_chart_with_template(
    output_path: Path,
    chart_part: str,
    template_path: Path,
    template_slide_index: int = 0,
    template_chart_index: int = 0,
) -> None:
    """Replace generated chart XML with template chart XML + dependencies."""
    template_prs = Presentation(str(template_path))
    if template_slide_index < 0 or template_slide_index >= len(template_prs.slides):
        raise ValueError(
            f"chart_template_slide {template_slide_index + 1} is out of range "
            f"(1-{len(template_prs.slides)})"
        )

    template_slide = template_prs.slides[template_slide_index]
    template_charts = slide_charts(template_slide)
    if template_chart_index < 0 or template_chart_index >= len(template_charts):
        raise ValueError(
            f"chart_template_chart_index {template_chart_index} is out of range "
            f"(0-{max(0, len(template_charts) - 1)})"
        )

    template_chart = template_charts[template_chart_index]
    template_chart_obj = cast(Any, template_chart)
    template_chart_part = str(template_chart_obj.part.partname).lstrip("/")
    template_chart_rels_part = f"ppt/charts/_rels/{Path(template_chart_part).name}.rels"

    template_chart_xml = read_pptx_part(template_path, template_chart_part)
    template_chart_rels = read_pptx_part(template_path, template_chart_rels_part)

    rel_ns = {"rel": "http://schemas.openxmlformats.org/package/2006/relationships"}
    rels_root = ET.fromstring(template_chart_rels)
    embedding_parts: dict[str, bytes] = {}

    for rel in rels_root.findall("rel:Relationship", rel_ns):
        if (
            rel.get("Type")
            != "http://schemas.openxmlformats.org/officeDocument/2006/relationships/package"
        ):
            continue

        target = rel.get("Target")
        if not target:
            continue

        part = normalize_relationship_target(template_chart_part, target)
        embedding_parts[part] = read_pptx_part(template_path, part)

    content_defaults = read_content_type_defaults(template_path)
    content_xml = read_pptx_part(output_path, "[Content_Types].xml")
    content_root = ET.fromstring(content_xml)
    ns = "{http://schemas.openxmlformats.org/package/2006/content-types}"

    existing_exts = {
        item.get("Extension")
        for item in content_root.findall(f"{ns}Default")
        if item.get("Extension")
    }

    for part in embedding_parts:
        ext = Path(part).suffix.lstrip(".")
        if not ext or ext in existing_exts:
            continue

        content_type = content_defaults.get(ext)
        if not content_type:
            continue

        node = ET.Element(f"{ns}Default")
        node.set("Extension", ext)
        node.set("ContentType", content_type)
        content_root.append(node)
        existing_exts.add(ext)

    new_content_xml = ET.tostring(content_root, xml_declaration=True, encoding="UTF-8")

    normalized_chart_part = chart_part.lstrip("/")
    chart_rels_part = f"ppt/charts/_rels/{Path(normalized_chart_part).name}.rels"
    temp_path = output_path.with_suffix(output_path.suffix + ".tmp")

    with zipfile.ZipFile(output_path) as zin, zipfile.ZipFile(temp_path, "w") as zout:
        for item in zin.infolist():
            if item.filename == normalized_chart_part:
                zout.writestr(item, template_chart_xml)
            elif item.filename == chart_rels_part:
                zout.writestr(item, template_chart_rels)
            elif item.filename == "[Content_Types].xml":
                zout.writestr(item, new_content_xml)
            elif item.filename in embedding_parts:
                continue
            else:
                zout.writestr(item, zin.read(item.filename))

        for part, data in embedding_parts.items():
            zout.writestr(part, data)

    temp_path.replace(output_path)


def slide_master_signature(path: Path) -> str:
    """Compute a stable hash of all slide-master XML parts."""
    with zipfile.ZipFile(path) as pptx:
        master_files = sorted(
            name
            for name in pptx.namelist()
            if name.startswith("ppt/slideMasters/slideMaster") and name.endswith(".xml")
        )
        if not master_files:
            raise ValueError(f"No slide masters found in {path}")

        digest = hashlib.sha256()
        for name in master_files:
            digest.update(name.encode("utf-8"))
            digest.update(pptx.read(name))

    return digest.hexdigest()


def theme_name(path: Path) -> str | None:
    """Return theme name from ``ppt/theme/theme1.xml`` if present."""
    try:
        theme_xml = read_pptx_part(path, "ppt/theme/theme1.xml")
    except ValueError:
        return None

    root = ET.fromstring(theme_xml)
    return root.attrib.get("name")


class ChartTemplateReplacement(NamedTuple):
    """Pending replacement action for chart template copy mode."""

    chart_part: str
    template_path: Path
    template_slide_index: int
    template_chart_index: int


def apply_chart_template_replacements(
    output_path: Path,
    replacements: Sequence[ChartTemplateReplacement],
) -> None:
    """Apply all queued chart template replacements."""
    for replacement in replacements:
        replace_chart_with_template(
            output_path,
            replacement.chart_part,
            replacement.template_path,
            template_slide_index=replacement.template_slide_index,
            template_chart_index=replacement.template_chart_index,
        )
