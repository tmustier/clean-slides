"""
Chart cell renderer — embeds native chart shapes inside table cells.

Delegates to the bundled chart engine module (loaded via charts.py) for
all chart creation and styling. This module handles only the table-specific
logic: grouping adjacent ChartRef cells, computing bounding boxes from the
table layout, and translating ChartDef into the JSON spec the generator expects.
"""

from __future__ import annotations

from collections.abc import Iterable
from dataclasses import dataclass
from typing import Protocol, cast

from pptx.oxml.xmlchemy import OxmlElement
from pptx.slide import Slide
from pptx.util import Emu, Pt

from .chart_engine.spec_utils import object_list, optional_str_list, str_key_dict
from .charts import ChartEngine
from .pptx_access import (
    chart_first_plot,
    chart_series,
    chart_xml_element,
    paragraph_font,
    plot_data_labels,
    point_fill_fore_color,
    point_fill_solid,
    point_line_fill_background,
    series_points,
    set_chart_has_legend,
    set_font_size,
    set_plot_has_data_labels,
    set_text_frame_text,
    shape_chart,
    shape_has_text_frame,
    shape_text_frame,
    slide_add_chart,
    slide_size_emu,
    text_frame_paragraphs,
    text_frame_text,
)
from .spec import Box, ChartDef, ChartRef, ContentArea, TableLayout, TableSpec

# ---------------------------------------------------------------------------
# Chart group — a set of adjacent ChartRef cells that merge into one shape
# ---------------------------------------------------------------------------


@dataclass
class ChartGroup:
    """A contiguous block of cells referencing the same chart."""

    chart_def: ChartDef
    refs: list[tuple[int, int, int]]  # (body_row, body_col, bar_index)
    min_row: int
    max_row: int
    min_col: int
    max_col: int


class _XmlElementLike(Protocol):
    def iter(self, tag: str) -> Iterable[object]: ...

    def find(self, path: str) -> object | None: ...

    def append(self, element: object) -> None: ...

    def insert(self, index: int, element: object) -> None: ...

    def addprevious(self, element: object) -> None: ...


class _XmlSettable(Protocol):
    def set(self, key: str, value: str) -> None: ...


def _iterable_objects(value: object) -> list[object]:
    """Coerce list/tuple/set payloads into ``list[object]``."""
    if isinstance(value, set):
        return list(cast(set[object], value))
    return object_list(value)


def _int_set(value: object) -> set[int]:
    """Coerce list/tuple/set payload to ``set[int]``."""
    result: set[int] = set()
    for item in _iterable_objects(value):
        if isinstance(item, int):
            result.add(item)
    return result


def _collect_chart_groups(spec: TableSpec) -> list[ChartGroup]:
    """Scan the cell grid and group adjacent ChartRef cells by chart name."""
    if not spec.cells or not spec.chart_defs:
        return []

    refs_by_name: dict[str, list[tuple[int, int, int]]] = {}
    for ri, row in enumerate(spec.cells):
        for ci, value in enumerate(row):
            if isinstance(value, ChartRef) and value.name in spec.chart_defs:
                refs_by_name.setdefault(value.name, []).append((ri, ci, value.index))

    groups: list[ChartGroup] = []
    for chart_name, ref_list in refs_by_name.items():
        chart_def = spec.chart_defs[chart_name]
        rows = [r for r, _, _ in ref_list]
        cols = [c for _, c, _ in ref_list]
        groups.append(
            ChartGroup(
                chart_def=chart_def,
                refs=ref_list,
                min_row=min(rows),
                max_row=max(rows),
                min_col=min(cols),
                max_col=max(cols),
            )
        )

    return groups


def _chart_bounding_box(
    group: ChartGroup,
    layout: TableLayout,
    area: ContentArea,
    row_offset: int,
    col_offset: int,
) -> Box:
    """Compute the bounding box (x, y, w, h) in EMU for a chart group."""
    grid_row_start = group.min_row + row_offset
    grid_row_end = group.max_row + row_offset
    grid_col_start = group.min_col + col_offset
    grid_col_end = group.max_col + col_offset

    x = area.x + sum(layout.col_widths[:grid_col_start])
    w = sum(layout.col_widths[grid_col_start : grid_col_end + 1])

    y = area.y + sum(layout.row_heights[:grid_row_start])
    h = sum(layout.row_heights[grid_row_start : grid_row_end + 1])

    return (x, y, w, h)


# ---------------------------------------------------------------------------
# ChartDef → JSON spec conversion
# ---------------------------------------------------------------------------

_CELL_CHART_GAP_WIDTH = 50

# Label headroom: text height + gap to bar top, expressed as a
# multiplier on font size (in pt → EMU via 12700).
#   ~1.0× for the glyph, ~0.5× padding below the label to the bar.
_LABEL_HEADROOM_FACTOR: float = 1.5
_EMU_PER_PT: int = 12700
# Clamp so headroom never eats more than half the chart.
_MAX_HEADROOM_FRAC: float = 0.45


def _horizontal_plot_layout(
    chart_width_emu: int,
    label_font_size_pt: int,
    has_labels: bool,
    values: list[float],
    fmt: str,
) -> dict[str, float]:
    """Compute plot layout for a horizontal (bar) chart.

    Reserves width on the right for outside_end data labels, based on
    the widest formatted value at the given font size.
    """
    if not has_labels or label_font_size_pt <= 0 or chart_width_emu <= 0:
        return {"x": 0.0, "y": 0.0, "w": 1.0, "h": 1.0}

    # Estimate widest label string
    widest = 0.0
    for v in values:
        try:
            label = fmt.format(v)
        except (ValueError, IndexError, KeyError):
            label = str(v)
        # Rough char-width estimate: 0.6 × font size per character
        label_emu = len(label) * label_font_size_pt * 0.6 * _EMU_PER_PT
        widest = max(widest, label_emu)

    # Add padding (1 char-width gap between bar end and label)
    widest += label_font_size_pt * 0.6 * _EMU_PER_PT

    w_frac = max(1.0 - widest / chart_width_emu, 0.50)
    return {"x": 0.0, "y": 0.0, "w": round(w_frac, 3), "h": 1.0}


def _vertical_plot_layout(
    chart_height_emu: int, label_font_size_pt: int, has_labels: bool
) -> dict[str, float]:
    """Compute plot layout for a vertical (column) chart.

    Reserves exactly enough headroom for outside_end labels above the
    tallest bar: label text height + padding gap, as a fraction of the
    chart shape height.
    """
    if not has_labels or label_font_size_pt <= 0 or chart_height_emu <= 0:
        return {"x": 0.0, "y": 0.0, "w": 1.0, "h": 1.0}

    headroom_emu = int(label_font_size_pt * _LABEL_HEADROOM_FACTOR * _EMU_PER_PT)
    y_frac = min(headroom_emu / chart_height_emu, _MAX_HEADROOM_FRAC)
    return {"x": 0.0, "y": round(y_frac, 3), "w": 1.0, "h": round(1.0 - y_frac, 3)}


def _sorted_values(group: ChartGroup) -> list[float]:
    """Return chart values in bar order (sorted by ref index)."""
    sorted_refs = sorted(group.refs, key=lambda r: r[2])
    return [group.chart_def.values[ref[2] - 1] for ref in sorted_refs]


def _chart_def_to_spec(group: ChartGroup, label_font_size: int = 8) -> dict[str, object]:
    """Convert a ChartDef + group refs into a JSON spec dict.

    Dispatches to bar or waterfall spec builder based on chart type.
    """
    if group.chart_def.type == "waterfall":
        return _waterfall_chart_spec(group, label_font_size)
    return _bar_chart_spec(group, label_font_size)


def _bar_chart_spec(group: ChartGroup, label_font_size: int = 8) -> dict[str, object]:
    """Build a JSON spec dict for build_bar_payload (clustered bar chart)."""
    chart_def = group.chart_def
    values = _sorted_values(group)
    categories = [str(i + 1) for i in range(len(values))]

    series_entry: dict[str, object] = {
        "name": "Values",
        "values": values,
    }
    if chart_def.color:
        series_entry["color"] = chart_def.color

    # Data label format conversion: Python "€{}m" → Excel "0" style
    data_labels: dict[str, object] = {}
    fmt = chart_def.format
    if fmt and fmt != "{}":
        data_labels["format"] = _python_fmt_to_excel_format(fmt, values)
    else:
        data_labels["format"] = "General"

    label_pos = chart_def.label_position
    if label_pos == "on":
        data_labels["position"] = "center"
    elif label_pos == "none":
        pass  # no data labels
    else:
        data_labels["position"] = "outside_end"

    data_labels["font_size"] = label_font_size

    bar_spec: dict[str, object] = {
        "orientation": chart_def.dir,
        "gap_width": _CELL_CHART_GAP_WIDTH,
        "overlap": 0,
        "plot_layout": {},  # computed dynamically in render_chart_cells
        "series_border_color": "none",
        "axis_line_color": "none",
    }
    if chart_def.scale_max is not None:
        bar_spec["axis_max"] = chart_def.scale_max
    if chart_def.scale_min is not None:
        bar_spec["axis_min"] = chart_def.scale_min

    spec: dict[str, object] = {
        "type": "clustered",
        "categories": categories,
        "series": [series_entry],
        "show_data_labels": label_pos != "none",
        "data_labels": data_labels,
        "show_legend": False,
        "orientation": chart_def.dir,
        "bar": bar_spec,
    }

    return spec


def _waterfall_chart_spec(group: ChartGroup, label_font_size: int = 8) -> dict[str, object]:
    """Build a JSON spec dict for build_waterfall_payload."""
    chart_def = group.chart_def
    values = _sorted_values(group)

    # Keep category labels visually empty in the chart itself.
    # We still need one category per value for geometry/stacking.
    categories = ["" for _ in values]

    series_entry: dict[str, object] = {
        "name": "Values",
        "values": values,
    }
    if chart_def.color:
        series_entry["color"] = chart_def.color

    # Totals/decreases are 1-based in YAML; convert to 0-based indices.
    total_categories: list[int] = []
    if chart_def.totals:
        total_categories = [max(i - 1, 0) for i in chart_def.totals]

    decrease_categories: list[int] = []
    if chart_def.decreases:
        decrease_categories = [max(i - 1, 0) for i in chart_def.decreases]

    # Data label format
    fmt = chart_def.format
    excel_fmt: str | None = None
    if fmt and fmt != "{}":
        excel_fmt = _python_fmt_to_excel_format(fmt, values)

    label_pos = chart_def.label_position

    wf_config: dict[str, object] = {
        "orientation": chart_def.dir,
        "total_categories": total_categories,
        "decrease_categories": decrease_categories,
        "gap_width": _CELL_CHART_GAP_WIDTH,
        "overlap": 100,
        "plot_layout": {},  # computed dynamically in render_chart_cells
        "axis_line_color": "none",  # hide category axis line in cells
        "total_override": True,  # honor explicit total values from YAML
    }

    if chart_def.total_color:
        wf_config["total_color"] = chart_def.total_color

    if not chart_def.connector:
        wf_config["connector_style"] = "none"
    else:
        wf_config["connector_style"] = "gap"
        # In table cells, nudge connectors half a stroke inside the bar tip
        # for cleaner visual alignment across renderers.
        wf_config["connector_inset"] = Emu(3000)
        wf_config["connector_overlap"] = 0

    if excel_fmt:
        wf_config["data_label_format"] = excel_fmt

    spec: dict[str, object] = {
        "type": "waterfall",
        "categories": categories,
        "series": [series_entry],
        "show_data_labels": label_pos != "none",
        "show_legend": False,
        "orientation": chart_def.dir,
        "waterfall": wf_config,
    }

    return spec


def _python_fmt_to_excel_format(fmt: str, values: list[float]) -> str:
    """Convert a Python-style format like '€{}m' to an Excel number format.

    Examples:
        '€{}m'    → '"€"0"m"'
        '{}x'     → '0.0"x"'
        '{:.1f}x' → '0.0"x"'
        '{:.2f}%' → '0.00"%"'
    """
    if fmt == "{}":
        return "General"

    import re

    # Match {}, {:.<N>f}, {:.0f}, etc.
    m = re.match(r"^(.*?)\{(?::([^}]*))?\}(.*)$", fmt)
    if m is None:
        return "General"

    prefix = m.group(1)
    format_spec = m.group(2) or ""
    suffix = m.group(3)

    # Determine decimal places from format spec (e.g. ".1f" → 1)
    decimal_places: int | None = None
    spec_match = re.match(r"\.(\d+)f", format_spec)
    if spec_match:
        decimal_places = int(spec_match.group(1))

    if decimal_places is not None:
        num_fmt = "0" + ("." + "0" * decimal_places if decimal_places > 0 else "")
    else:
        has_decimals = any(v != int(v) for v in values)
        num_fmt = "0.0" if has_decimals else "0"

    result = ""
    if prefix:
        result += f'"{prefix}"'
    result += num_fmt
    if suffix:
        result += f'"{suffix}"'

    return result


def _format_chart_label_value(value: float | None, fmt: str) -> str | None:
    """Format a chart value using the user-provided Python format string."""
    if value is None:
        return None

    try:
        if fmt == "{}":
            rounded = round(value)
            if abs(value - rounded) < 1e-9:
                return str(int(rounded))
            return str(value)
        return str(fmt.format(value))
    except (ValueError, IndexError, KeyError):
        rounded = round(value)
        if abs(value - rounded) < 1e-9:
            return str(int(rounded))
        return str(value)


def _waterfall_overlay_label_texts(meta: dict[str, object], fmt: str) -> list[str]:
    """Compute overlay label texts for a waterfall chart.

    Mirrors the generator's label-value selection logic but applies the
    clean-slides format string (e.g. ``{:,.0f}``) so labels can include
    thousands separators and custom affixes.
    """
    overlay_obj = str_key_dict(meta.get("overlay"))
    if not overlay_obj:
        return []

    categories_obj = overlay_obj.get("categories")
    cumulative_obj = overlay_obj.get("cumulative_totals")
    if not isinstance(categories_obj, list) or not isinstance(cumulative_obj, list):
        return []

    categories = object_list(cast(object, categories_obj))
    cumulative_values = object_list(cast(object, cumulative_obj))
    delta_values_raw = object_list(overlay_obj.get("delta_values"))
    total_categories = _int_set(overlay_obj.get("total_categories"))

    categories_count = len(categories)
    cumulative_totals: list[float | None] = [
        float(value) if isinstance(value, (int, float)) else None for value in cumulative_values
    ]
    delta_values: list[float | None] = [
        float(value) if isinstance(value, (int, float)) else None for value in delta_values_raw
    ]

    texts: list[str] = []
    for idx in range(categories_count):
        total_value = cumulative_totals[idx] if idx < len(cumulative_totals) else None
        if total_value is None:
            continue

        if idx == 0 or idx in total_categories:
            label_value = total_value
        else:
            label_value = delta_values[idx] if idx < len(delta_values) else None

        label_text = _format_chart_label_value(label_value, fmt)
        if label_text is not None:
            texts.append(label_text)

    return texts


def _hide_waterfall_overlay_category_labels(meta: dict[str, object]) -> None:
    """Suppress category labels rendered by add_waterfall_overlays()."""
    overlay_obj = str_key_dict(meta.get("overlay"))
    if not overlay_obj:
        return

    categories_obj = overlay_obj.get("categories")
    if not isinstance(categories_obj, list):
        return

    overlay_obj["categories"] = ["" for _item in object_list(cast(object, categories_obj))]
    meta["overlay"] = overlay_obj


def _rewrite_overlay_value_label_texts(
    slide: Slide,
    start_shape_index: int,
    label_texts: list[str],
    font_size_pt: int,
) -> None:
    """Rewrite newly added overlay text shapes with formatted label text."""
    if not label_texts:
        return

    label_idx = 0
    total_shapes = len(slide.shapes)
    for shape_idx in range(start_shape_index, total_shapes):
        if label_idx >= len(label_texts):
            break

        shape = slide.shapes[shape_idx]
        if not shape_has_text_frame(shape):
            continue

        text_frame = shape_text_frame(shape)
        if text_frame is None:
            continue

        existing_text = text_frame_text(text_frame)
        if existing_text is None or not existing_text.strip():
            continue

        set_text_frame_text(text_frame, label_texts[label_idx])

        paragraphs = text_frame_paragraphs(text_frame)
        if paragraphs:
            first_paragraph = paragraphs[0]
            font = paragraph_font(first_paragraph)
            if font is not None:
                set_font_size(font, Pt(font_size_pt))

        label_idx += 1


def _set_label_nowrap(chart: object) -> None:
    """Set ``wrap="none"`` on every data-label ``bodyPr`` in the chart.

    Without this, PowerPoint (and LibreOffice) may wrap short labels
    like ``4.0x`` across multiple lines when the default label text box
    is narrower than the formatted string.
    """

    ns_a = "http://schemas.openxmlformats.org/drawingml/2006/main"
    ns_c = "http://schemas.openxmlformats.org/drawingml/2006/chart"
    chart_el_obj = chart_xml_element(chart)
    if chart_el_obj is None:
        return
    chart_el = cast(_XmlElementLike, chart_el_obj)

    for dlbls_obj in chart_el.iter(f"{{{ns_c}}}dLbls"):
        dlbls = cast(_XmlElementLike, dlbls_obj)
        tx_pr_obj = dlbls.find(f"{{{ns_c}}}txPr")
        if tx_pr_obj is None:
            continue
        tx_pr = cast(_XmlElementLike, tx_pr_obj)
        body_pr_obj = tx_pr.find(f"{{{ns_a}}}bodyPr")
        if body_pr_obj is None:
            continue
        body_pr = cast(_XmlSettable, body_pr_obj)
        body_pr.set("wrap", "none")


# Visual gap between bar tip and label left edge, as a fraction of chart width.
_HORIZONTAL_LABEL_VISUAL_GAP: float = 0.02


def _set_horizontal_label_offsets(
    chart: object,
    values: list[float],
    axis_max: float,
    plot_w: float,
    chart_width_emu: int,
    num_format: str = "General",
    label_format: str = "{}",
    font_size_pt: int = 12,
) -> None:
    """Position horizontal bar labels to the right of each bar with a gap.

    PowerPoint's ``outEnd`` labels can overlap bar tips.  The reference
    deck avoids this by using ``ctr`` (center) position on each label,
    then applying a manual x-offset that pushes the label past the bar.

    Offset from center = (value / axis_max) × plot_w / 2 + gap

    This places every label at the bar's right edge + a fixed gap,
    regardless of bar length.
    """
    from lxml import etree

    chart_el_obj = chart_xml_element(chart)
    if chart_el_obj is None:
        return
    chart_el = cast(_XmlElementLike, chart_el_obj)
    ns_c = "http://schemas.openxmlformats.org/drawingml/2006/chart"

    # The OOXML structure is:
    #   barChart > dLbls  (chart-level defaults: showVal etc.)
    #   barChart > ser > dLbls > dLbl*  (per-point overrides with offsets)
    # python-pptx creates dLbls at the barChart level.  We need a ser-level
    # dLbls container to hold per-point dLbl elements with manual offsets.

    bar_chart_obj = chart_el.find(f".//{{{ns_c}}}barChart")
    if bar_chart_obj is None:
        return
    bar_chart = cast(_XmlElementLike, bar_chart_obj)

    ser_obj = bar_chart.find(f"{{{ns_c}}}ser")
    if ser_obj is None:
        return
    ser = cast(_XmlElementLike, ser_obj)

    # Create ser-level dLbls container (or reuse if it already exists)
    ser_dlbls_obj = ser.find(f"{{{ns_c}}}dLbls")
    if ser_dlbls_obj is None:
        ser_dlbls_el = OxmlElement("c:dLbls")
        # Insert before cat/val
        insert_before_obj: object | None = None
        for tag_suffix in ("cat", "val", "shape", "extLst"):
            insert_before_obj = ser.find(f"{{{ns_c}}}{tag_suffix}")
            if insert_before_obj is not None:
                break
        if insert_before_obj is not None:
            insert_before = cast(_XmlElementLike, insert_before_obj)
            insert_before.addprevious(ser_dlbls_el)
        else:
            ser.append(ser_dlbls_el)
        ser_dlbls_obj = ser.find(f"{{{ns_c}}}dLbls")
        if ser_dlbls_obj is None:
            return
    ser_dlbls = cast(_XmlElementLike, ser_dlbls_obj)

    # Character width estimate: ~0.6 × font size per character
    char_width_emu = font_size_pt * 0.6 * _EMU_PER_PT
    chart_w = max(chart_width_emu, 1)

    for i, v in enumerate(values):
        ratio = v / axis_max if axis_max > 0 else 0

        # Estimate this label's width in chart-fraction units
        try:
            label_text = label_format.format(v)
        except (ValueError, IndexError, KeyError):
            label_text = str(v)
        label_w_frac = len(label_text) * char_width_emu / chart_w
        # Offset = half_bar (from ctr to tip) + visual_gap + half_label
        x_offset = ratio * plot_w / 2 + _HORIZONTAL_LABEL_VISUAL_GAP + label_w_frac / 2

        ns_a = "http://schemas.openxmlformats.org/drawingml/2006/main"
        dlbl = OxmlElement("c:dLbl")

        idx_el = OxmlElement("c:idx")
        idx_el.set("val", str(i))
        dlbl.append(idx_el)

        # Manual layout with x-offset
        layout_el = OxmlElement("c:layout")
        manual = OxmlElement("c:manualLayout")
        x_el = OxmlElement("c:x")
        x_el.set("val", str(x_offset))
        manual.append(x_el)
        y_el = OxmlElement("c:y")
        y_el.set("val", "0")
        manual.append(y_el)
        layout_el.append(manual)
        dlbl.append(layout_el)

        # Number format
        num_fmt = OxmlElement("c:numFmt")
        num_fmt.set("formatCode", num_format)
        num_fmt.set("sourceLinked", "0")
        dlbl.append(num_fmt)

        # Transparent background (no fill, no outline)
        spPr = OxmlElement("c:spPr")
        etree.SubElement(spPr, f"{{{ns_a}}}noFill")
        ln = etree.SubElement(spPr, f"{{{ns_a}}}ln")
        etree.SubElement(ln, f"{{{ns_a}}}noFill")
        dlbl.append(spPr)

        # Text properties: no wrap, font size
        txPr = OxmlElement("c:txPr")
        bodyPr = etree.SubElement(txPr, f"{{{ns_a}}}bodyPr")
        bodyPr.set("wrap", "none")
        etree.SubElement(txPr, f"{{{ns_a}}}lstStyle")
        p = etree.SubElement(txPr, f"{{{ns_a}}}p")
        pPr = etree.SubElement(p, f"{{{ns_a}}}pPr")
        defRPr = etree.SubElement(pPr, f"{{{ns_a}}}defRPr")
        defRPr.set("sz", str(font_size_pt * 100))  # hundredths of a pt
        dlbl.append(txPr)

        pos_el = OxmlElement("c:dLblPos")
        pos_el.set("val", "ctr")
        dlbl.append(pos_el)

        for tag_name, val in [
            ("showLegendKey", "0"),
            ("showVal", "1"),
            ("showCatName", "0"),
            ("showSerName", "0"),
            ("showPercent", "0"),
            ("showBubbleSize", "0"),
        ]:
            el = OxmlElement(f"c:{tag_name}")
            el.set("val", val)
            dlbl.append(el)

        ser_dlbls.insert(i, dlbl)

    # Add shared settings to ser-level dLbls (after per-point labels)
    for tag_name, val in [
        ("showLegendKey", "0"),
        ("showVal", "1"),
        ("showCatName", "0"),
        ("showSerName", "0"),
        ("showPercent", "0"),
        ("showBubbleSize", "0"),
    ]:
        shared_el = OxmlElement(f"c:{tag_name}")
        shared_el.set("val", val)
        ser_dlbls.append(shared_el)


def _delete_auto_title(chart: object) -> None:
    """Suppress the auto-generated chart title (series name watermark)."""
    ns_c = "http://schemas.openxmlformats.org/drawingml/2006/chart"
    chart_el_obj = chart_xml_element(chart)
    if chart_el_obj is None:
        return

    chart_el = cast(_XmlElementLike, chart_el_obj)
    auto_title_obj = chart_el.find(".//{" + ns_c + "}autoTitleDeleted")
    if auto_title_obj is None:
        return

    auto_title = cast(_XmlSettable, auto_title_obj)
    auto_title.set("val", "1")


def _apply_point_colors(
    chart: object,
    charts_module: ChartEngine,
    point_colors: list[str | None],
    series_idx: int,
) -> None:
    """Apply optional per-point fill colors to one chart series."""
    if not point_colors:
        return

    series_list = chart_series(chart)
    if series_idx < 0 or series_idx >= len(series_list):
        return

    series = series_list[series_idx]
    points = series_points(series)
    for point_idx, color in enumerate(point_colors):
        if not color or point_idx >= len(points):
            continue

        point = points[point_idx]
        point_fill_solid(point)
        fore_color = point_fill_fore_color(point)
        if fore_color is None:
            continue

        applied = bool(charts_module.apply_color(fore_color, color))
        if applied:
            point_line_fill_background(point)


# ---------------------------------------------------------------------------
# Public API
# ---------------------------------------------------------------------------


def render_chart_cells(
    slide: Slide,
    spec: TableSpec,
    layout: TableLayout,
    area: ContentArea,
    charts_module: ChartEngine,
    label_font_size_pt: int = 8,
) -> None:
    """Render all chart cell groups as native chart shapes on the slide.

    Uses the bundled chart engine module (loaded via charts.py) for
    chart creation and styling.
    """
    if not spec.chart_defs or not spec.cells:
        return

    groups = _collect_chart_groups(spec)
    if not groups:
        return

    row_offset = spec.row_offset
    col_offset = spec.col_offset

    for group in groups:
        box = _chart_bounding_box(group, layout, area, row_offset, col_offset)
        x, y, w, h = box

        # Inset by padding so the chart doesn't touch cell edges
        pad = layout.pad_top
        x += pad
        y += pad
        w = max(w - 2 * pad, 1)
        h = max(h - 2 * pad, 1)

        if group.chart_def.type == "waterfall":
            _render_waterfall_group(
                slide,
                group,
                charts_module,
                x,
                y,
                w,
                h,
                label_font_size_pt,
            )
        else:
            _render_bar_group(
                slide,
                group,
                charts_module,
                x,
                y,
                w,
                h,
                label_font_size_pt,
            )


def _render_bar_group(
    slide: Slide,
    group: ChartGroup,
    charts_module: ChartEngine,
    x: int,
    y: int,
    w: int,
    h: int,
    label_font_size_pt: int,
) -> None:
    """Render a clustered bar chart group."""
    chart_spec = _chart_def_to_spec(group, label_font_size=label_font_size_pt)

    # Compute direction-aware plot layout
    has_labels = bool(chart_spec.get("show_data_labels", False))
    if group.chart_def.dir == "horizontal":
        plot_layout = _horizontal_plot_layout(
            w,
            label_font_size_pt,
            has_labels,
            group.chart_def.values,
            group.chart_def.format,
        )
    else:
        plot_layout = _vertical_plot_layout(h, label_font_size_pt, has_labels)

    bar_cfg = str_key_dict(chart_spec.get("bar"))
    bar_cfg["plot_layout"] = plot_layout
    chart_spec["bar"] = bar_cfg

    chart_type, chart_data, style = charts_module.build_bar_payload(chart_spec)

    chart_frame = slide_add_chart(
        slide,
        chart_type,
        Emu(x),
        Emu(y),
        Emu(w),
        Emu(h),
        chart_data,
    )
    if chart_frame is None:
        return

    chart = shape_chart(chart_frame)
    if chart is None:
        return

    set_chart_has_legend(chart, False)
    _delete_auto_title(chart)

    series_colors = optional_str_list(style.get("series_colors", []))
    charts_module.apply_series_colors(chart, series_colors)

    if group.chart_def.colors:
        _apply_point_colors(chart, charts_module, group.chart_def.colors, series_idx=0)

    style_bar = str_key_dict(style.get("bar"))
    if style_bar:
        charts_module.apply_bar_chart_style(chart, style_bar)

    if has_labels:
        data_cfg = str_key_dict(chart_spec.get("data_labels"))
        plot = chart_first_plot(chart)
        if plot is None:
            return

        set_plot_has_data_labels(plot, True)
        labels = plot_data_labels(plot)
        if labels is None:
            return

        charts_module.apply_data_label_style(labels, data_cfg)
        _set_label_nowrap(chart)

        if group.chart_def.dir == "horizontal":
            axis_max_obj = bar_cfg.get("axis_max")
            axis_max = (
                float(axis_max_obj)
                if isinstance(axis_max_obj, (int, float))
                else max(group.chart_def.values)
            )
            excel_fmt = str(data_cfg.get("format", "General"))
            _set_horizontal_label_offsets(
                chart,
                group.chart_def.values,
                axis_max,
                plot_layout["w"],
                chart_width_emu=w,
                num_format=excel_fmt,
                label_format=group.chart_def.format,
                font_size_pt=label_font_size_pt,
            )


def _render_waterfall_group(
    slide: Slide,
    group: ChartGroup,
    charts_module: ChartEngine,
    x: int,
    y: int,
    w: int,
    h: int,
    label_font_size_pt: int,
) -> None:
    """Render a waterfall chart group."""
    chart_spec = _chart_def_to_spec(group, label_font_size=label_font_size_pt)

    # Compute plot layout — waterfall labels use center positioning,
    # so reserve space at the right for the largest bar label.
    has_labels = bool(chart_spec.get("show_data_labels", False))
    values = _sorted_values(group)
    if group.chart_def.dir == "horizontal":
        plot_layout = _horizontal_plot_layout(
            w,
            label_font_size_pt,
            has_labels,
            values,
            group.chart_def.format,
        )
    else:
        plot_layout = _vertical_plot_layout(h, label_font_size_pt, has_labels)

    waterfall_cfg = str_key_dict(chart_spec.get("waterfall"))
    waterfall_cfg["plot_layout"] = plot_layout
    chart_spec["waterfall"] = waterfall_cfg

    chart_type, chart_data, style = charts_module.build_waterfall_payload(chart_spec)

    chart_frame = slide_add_chart(
        slide,
        chart_type,
        Emu(x),
        Emu(y),
        Emu(w),
        Emu(h),
        chart_data,
    )
    if chart_frame is None:
        return

    chart = shape_chart(chart_frame)
    if chart is None:
        return

    set_chart_has_legend(chart, False)
    _delete_auto_title(chart)

    # Apply series colors (offset series + value series)
    series_colors = optional_str_list(style.get("series_colors", []))
    charts_module.apply_series_colors(chart, series_colors)

    # Apply waterfall-specific styling (offset no-fill, total point colors)
    wf_meta = str_key_dict(style.get("waterfall", {}))
    if wf_meta:
        charts_module.apply_waterfall_style(chart, wf_meta)
        charts_module.apply_waterfall_chart_style(chart, wf_meta)

    if group.chart_def.colors:
        offset_idx_obj = wf_meta.get("offset_series_idx")
        offset_idx = offset_idx_obj if isinstance(offset_idx_obj, int) else 0
        series_list = chart_series(chart)
        target_idx: int | None = None
        for idx in range(len(series_list)):
            if idx != offset_idx:
                target_idx = idx
                break
        if target_idx is not None:
            _apply_point_colors(chart, charts_module, group.chart_def.colors, series_idx=target_idx)

    # Add connector lines and overlay labels as separate shapes on the slide.
    # This handles both incremental and total bar labels, plus connectors.
    # We do NOT use built-in chart data labels — overlays give full control.
    if has_labels and wf_meta:
        # Cell tables already provide row labels, so hide overlay category labels.
        _hide_waterfall_overlay_category_labels(wf_meta)

        # Pre-compute desired label texts (supports format strings like {:,.0f}).
        label_texts = _waterfall_overlay_label_texts(wf_meta, group.chart_def.format)

        slide_size = slide_size_emu(slide)
        if slide_size is None:
            return

        before_shape_count = len(slide.shapes)
        charts_module.add_waterfall_overlays(
            slide,
            (x, y, w, h),
            wf_meta,
            slide_size=slide_size,
        )
        _rewrite_overlay_value_label_texts(
            slide,
            before_shape_count,
            label_texts,
            label_font_size_pt,
        )
