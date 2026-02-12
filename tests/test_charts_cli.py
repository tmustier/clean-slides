from __future__ import annotations

import copy
import json
from dataclasses import dataclass
from pathlib import Path
from typing import ClassVar, cast
from unittest.mock import patch

from pptx import Presentation
from pptx.enum.chart import XL_CHART_TYPE

from clean_slides.chart_engine.payloads import build_bar_payload
from clean_slides.charts import generate_charts_from_json
from clean_slides.cli_render import cmd_charts


@dataclass
class ChartsArgs:
    input: str
    output: str
    template: str | None = None
    layout: str | None = None
    expected_template: str | None = None


def _write_json(path: Path, payload: object) -> None:
    path.write_text(json.dumps(payload), encoding="utf-8")


def _to_str_dict(value: object) -> dict[str, object]:
    if not isinstance(value, dict):
        return {}

    result: dict[str, object] = {}
    for key, item in cast(dict[object, object], value).items():
        if isinstance(key, str):
            result[key] = item
    return result


def _chart_types(path: Path) -> list[int]:
    prs = Presentation(str(path))
    chart_types: list[int] = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if not bool(getattr(shape, "has_chart", False)):
                continue
            chart = getattr(shape, "chart", None)
            if chart is None:
                continue
            chart_type = getattr(chart, "chart_type", None)
            if chart_type is None:
                continue
            chart_types.append(int(chart_type))
    return chart_types


def _run_cmd_charts(tmp_path: Path, name: str, spec: dict[str, object]) -> Path:
    input_path = tmp_path / f"{name}.json"
    output_path = tmp_path / f"{name}.pptx"
    _write_json(input_path, spec)

    result = cmd_charts(
        ChartsArgs(
            input=str(input_path),
            output=str(output_path),
        )
    )

    assert result == 0
    assert output_path.exists()
    return output_path


def _simple_chart_spec(value: int) -> dict[str, object]:
    return {
        "type": "clustered",
        "categories": ["A"],
        "series": [{"name": "S1", "values": [value]}],
    }


class _FakePresentation:
    created_templates: ClassVar[list[str | None]] = []

    def __init__(self, template: str | None = None) -> None:
        self.template = template
        _FakePresentation.created_templates.append(template)

    def save(self, path: str | Path) -> None:
        Path(path).write_bytes(b"fake-pptx")


class _FakeEngine:
    Presentation = _FakePresentation

    def __init__(self, raw_spec: object) -> None:
        self.raw_spec = raw_spec
        self.built_specs: list[dict[str, object]] = []
        self.template_paths: list[Path | None] = []
        self.save_flags: list[bool] = []
        self.defer_template_copy_flags: list[bool] = []

    def load_spec(self, path: Path) -> object:
        return self.raw_spec

    def normalize_chart_specs(
        self, raw: object
    ) -> tuple[list[dict[str, object]], dict[str, object]]:
        if isinstance(raw, dict):
            raw_dict = _to_str_dict(cast(object, raw))
            if "charts" in raw_dict:
                charts_raw = raw_dict.get("charts")
                charts: list[dict[str, object]] = []
                if isinstance(charts_raw, list):
                    for item in cast(list[object], charts_raw):
                        chart_spec = _to_str_dict(item)
                        if chart_spec:
                            charts.append(copy.deepcopy(chart_spec))
                deck_meta = {key: value for key, value in raw_dict.items() if key != "charts"}
                return charts, deck_meta

            return [copy.deepcopy(raw_dict)], {}

        if isinstance(raw, list):
            charts: list[dict[str, object]] = []
            for item in cast(list[object], raw):
                chart_spec = _to_str_dict(item)
                if chart_spec:
                    charts.append(copy.deepcopy(chart_spec))
            return charts, {}

        return [], {}

    def resolve_expected_template(
        self,
        spec: dict[str, object],
        spec_path: Path,
        expected_template: str | Path | None,
    ) -> Path | None:
        return None

    def ensure_expected_template(self, expected: Path, actual: Path | None) -> None:
        raise AssertionError("ensure_expected_template should not be called in this test")

    def build_chart(
        self,
        prs: object,
        spec: dict[str, object],
        output_path: Path,
        template_path: Path | None = None,
        layout_name: str | None = None,
        save: bool = True,
        defer_template_copy: bool = False,
    ) -> list[object]:
        self.built_specs.append(copy.deepcopy(spec))
        self.template_paths.append(template_path)
        self.save_flags.append(save)
        self.defer_template_copy_flags.append(defer_template_copy)
        return []

    def apply_chart_template_replacements(
        self, output_path: Path, replacements: list[object]
    ) -> None:
        return None


def test_cmd_charts_clustered_smoke(tmp_path: Path) -> None:
    output_path = _run_cmd_charts(
        tmp_path,
        "clustered",
        {
            "type": "clustered",
            "categories": ["A", "B", "C"],
            "series": [{"name": "S1", "values": [10, 20, 30], "color": "accent1"}],
            "show_data_labels": True,
            "data_labels": {"font_size": 14, "format": "0"},
        },
    )

    assert _chart_types(output_path) == [int(XL_CHART_TYPE.COLUMN_CLUSTERED)]


def test_cmd_charts_horizontal_smoke(tmp_path: Path) -> None:
    output_path = _run_cmd_charts(
        tmp_path,
        "horizontal",
        {
            "type": "clustered",
            "categories": ["A", "B"],
            "series": [{"name": "S1", "values": [100, 200], "color": "accent2"}],
            "show_data_labels": True,
            "data_labels": {"font_size": 14, "format": "0"},
            "bar": {"orientation": "horizontal"},
        },
    )

    assert _chart_types(output_path) == [int(XL_CHART_TYPE.BAR_CLUSTERED)]


def test_cmd_charts_waterfall_smoke(tmp_path: Path) -> None:
    output_path = _run_cmd_charts(
        tmp_path,
        "waterfall",
        {
            "type": "waterfall",
            "categories": ["Start", "Growth", "Costs", "End"],
            "series": [
                {"name": "Values", "values": [100, 40, -20, 120], "color": "accent1"},
            ],
            "show_data_labels": True,
            "add_overlay_labels": True,
            "waterfall": {
                "total_categories": ["Start", "End"],
                "decrease_categories": ["Costs"],
                "connector_style": "gap",
            },
        },
    )

    assert _chart_types(output_path) == [int(XL_CHART_TYPE.COLUMN_STACKED)]


def test_generate_charts_from_json_sets_base_dir_for_each_chart(tmp_path: Path) -> None:
    spec_path = tmp_path / "specs" / "charts.json"
    spec_path.parent.mkdir(parents=True, exist_ok=True)
    output_path = tmp_path / "out.pptx"

    raw_spec = {
        "charts": [
            _simple_chart_spec(1),
            _simple_chart_spec(2),
        ],
    }
    _write_json(spec_path, raw_spec)

    fake_engine = _FakeEngine(raw_spec)
    with patch("clean_slides.charts.load_chart_engine", return_value=fake_engine):
        generate_charts_from_json(spec_path, output_path)

    expected_base_dir = str(spec_path.parent.resolve())
    assert len(fake_engine.built_specs) == 2
    assert all(spec.get("_base_dir") == expected_base_dir for spec in fake_engine.built_specs)
    assert fake_engine.built_specs[1].get("append_slide") is True
    assert output_path.exists()


def test_generate_charts_from_json_defers_template_copy_in_builder_calls(tmp_path: Path) -> None:
    spec_path = tmp_path / "specs" / "charts.json"
    spec_path.parent.mkdir(parents=True, exist_ok=True)
    output_path = tmp_path / "out-defer.pptx"

    raw_spec = {
        "charts": [
            _simple_chart_spec(1),
            _simple_chart_spec(2),
        ],
    }
    _write_json(spec_path, raw_spec)

    fake_engine = _FakeEngine(raw_spec)
    with patch("clean_slides.charts.load_chart_engine", return_value=fake_engine):
        generate_charts_from_json(spec_path, output_path)

    assert fake_engine.save_flags == [False, False]
    assert fake_engine.defer_template_copy_flags == [True, True]


def test_generate_charts_from_json_resolves_relative_deck_template(tmp_path: Path) -> None:
    spec_path = tmp_path / "specs" / "charts.json"
    spec_path.parent.mkdir(parents=True, exist_ok=True)
    output_path = tmp_path / "out-template.pptx"

    raw_spec = {
        "template": "templates/base.pptx",
        "charts": [_simple_chart_spec(1)],
    }
    _write_json(spec_path, raw_spec)

    _FakePresentation.created_templates = []
    fake_engine = _FakeEngine(raw_spec)
    with patch("clean_slides.charts.load_chart_engine", return_value=fake_engine):
        generate_charts_from_json(spec_path, output_path)

    expected_template = (spec_path.parent / "templates" / "base.pptx").resolve()
    assert _FakePresentation.created_templates == [str(expected_template)]
    assert fake_engine.template_paths == [expected_template]


def test_build_bar_payload_resolves_chart_template_against_base_dir(tmp_path: Path) -> None:
    base_dir = tmp_path / "spec-dir"
    base_dir.mkdir(parents=True, exist_ok=True)

    _chart_type, _chart_data, style = build_bar_payload(
        {
            "_base_dir": str(base_dir),
            "type": "clustered",
            "categories": ["A"],
            "series": [{"name": "S1", "values": [1]}],
            "bar": {
                "chart_template": "templates/style.pptx",
            },
        }
    )

    expected_template = (base_dir / "templates" / "style.pptx").resolve()
    bar_style = cast(dict[str, object], style["bar"])
    assert bar_style["chart_template"] == str(expected_template)


def test_build_bar_payload_keeps_empty_chart_template_disabled(tmp_path: Path) -> None:
    base_dir = tmp_path / "spec-dir"
    base_dir.mkdir(parents=True, exist_ok=True)

    _chart_type, _chart_data, style = build_bar_payload(
        {
            "_base_dir": str(base_dir),
            "type": "clustered",
            "categories": ["A"],
            "series": [{"name": "S1", "values": [1]}],
            "bar": {
                "chart_template": "",
            },
        }
    )

    bar_style = cast(dict[str, object], style["bar"])
    assert bar_style["chart_template"] == ""
