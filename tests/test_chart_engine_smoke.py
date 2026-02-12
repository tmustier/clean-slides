from __future__ import annotations

import zipfile
from pathlib import Path

from pptx import Presentation
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.shapes import MSO_SHAPE_TYPE

from clean_slides.chart_engine.builder import build_chart
from clean_slides.pptx_access import (
    iter_shapes,
    iter_slides,
    presentation_chart_types,
    shape_text_frame_text,
)


def _chart_types(path: Path) -> list[int]:
    prs = Presentation(str(path))
    return presentation_chart_types(prs)


def _shape_texts(path: Path) -> list[str]:
    prs = Presentation(str(path))
    texts: list[str] = []
    for slide in iter_slides(prs):
        for shape in iter_shapes(slide):
            text_value = shape_text_frame_text(shape)
            if text_value is None:
                continue
            text = text_value.strip()
            if text:
                texts.append(text)
    return texts


def _pptx_part(path: Path, part_name: str) -> bytes:
    with zipfile.ZipFile(path) as archive:
        return archive.read(part_name)


def _blank_auto_shape_count(path: Path) -> int:
    prs = Presentation(str(path))
    count = 0
    for shape in iter_shapes(prs.slides[0]):
        shape_type = getattr(shape, "shape_type", None)
        if shape_type != MSO_SHAPE_TYPE.AUTO_SHAPE:
            continue

        text = shape_text_frame_text(shape)
        if text is not None and text.strip():
            continue

        count += 1
    return count


def _waterfall_overlay_spec(connector_style: str, connector_dash_style: str) -> dict[str, object]:
    return {
        "type": "waterfall",
        "categories": ["Start", "A", "B", "C", "End"],
        "series": [
            {"name": "Values", "values": [100, 25, -15, 10, 120], "color": "accent1"},
        ],
        "show_data_labels": True,
        "add_overlay_labels": True,
        "waterfall": {
            "orientation": "horizontal",
            "total_categories": ["Start", "End"],
            "decrease_categories": ["B"],
            "connector_value": "totals",
            "connector_style": connector_style,
            "connector_dash_style": connector_dash_style,
        },
    }


def test_build_chart_bar_with_overlays_smoke(tmp_path: Path) -> None:
    output_path = tmp_path / "bar-smoke.pptx"

    spec: dict[str, object] = {
        "type": "clustered",
        "categories": ["A", "B"],
        "series": [{"name": "S1", "values": [10, 20], "color": "accent1"}],
        "show_data_labels": True,
        "add_overlay_labels": True,
        "bar": {
            "orientation": "horizontal",
            "show_totals": True,
            "segment_labels": [{"show": True, "series_indices": [0]}],
        },
    }

    prs = Presentation()
    replacements = build_chart(prs, spec, output_path)

    assert replacements == []
    assert output_path.exists()
    assert _chart_types(output_path) == [int(XL_CHART_TYPE.BAR_CLUSTERED)]

    texts = _shape_texts(output_path)
    assert "A" in texts
    assert "B" in texts


def test_build_chart_waterfall_with_overlays_smoke(tmp_path: Path) -> None:
    output_path = tmp_path / "waterfall-smoke.pptx"

    spec: dict[str, object] = {
        "type": "waterfall",
        "categories": ["Start", "Growth", "Costs", "End"],
        "series": [{"name": "Values", "values": [100, 40, -20, 120], "color": "accent1"}],
        "show_data_labels": True,
        "add_overlay_labels": True,
        "waterfall": {
            "total_categories": ["Start", "End"],
            "decrease_categories": ["Costs"],
            "connector_style": "gap",
        },
    }

    prs = Presentation()
    replacements = build_chart(prs, spec, output_path)

    assert replacements == []
    assert output_path.exists()
    assert _chart_types(output_path) == [int(XL_CHART_TYPE.COLUMN_STACKED)]

    texts = _shape_texts(output_path)
    assert "Start" in texts
    assert "End" in texts


def test_build_chart_defers_chart_template_copy_replacement(tmp_path: Path) -> None:
    template_path = tmp_path / "template-without-charts.pptx"
    Presentation().save(str(template_path))

    output_path = tmp_path / "deferred-template-copy.pptx"
    spec: dict[str, object] = {
        "type": "clustered",
        "categories": ["A"],
        "series": [{"name": "S1", "values": [1], "color": "accent1"}],
        "bar": {
            "chart_template": str(template_path),
            "chart_template_copy": True,
        },
    }

    prs = Presentation()
    replacements = build_chart(
        prs,
        spec,
        output_path,
        save=False,
        defer_template_copy=True,
    )

    assert len(replacements) == 1
    replacement = replacements[0]
    assert replacement.template_path == template_path
    assert replacement.chart_part.startswith("ppt/charts/chart")


def test_build_chart_applies_chart_template_copy_immediately_when_saving(tmp_path: Path) -> None:
    template_path = tmp_path / "template-chart.pptx"
    template_spec: dict[str, object] = {
        "type": "clustered",
        "categories": ["A", "B"],
        "series": [{"name": "S1", "values": [1, 2], "color": "accent1"}],
        "show_data_labels": True,
        "data_labels": {"format": "0.00", "font_size": 18},
    }
    build_chart(Presentation(), template_spec, template_path)

    output_path = tmp_path / "immediate-template-copy.pptx"
    target_spec: dict[str, object] = {
        "type": "clustered",
        "categories": ["A", "B"],
        "series": [{"name": "S1", "values": [9, 8], "color": "accent2"}],
        "show_data_labels": True,
        "data_labels": {"format": "0", "font_size": 10},
        "bar": {
            "chart_template": str(template_path),
            "chart_template_copy": True,
        },
    }

    replacements = build_chart(
        Presentation(),
        target_spec,
        output_path,
        save=True,
        defer_template_copy=False,
    )

    assert replacements == []

    template_chart_xml = _pptx_part(template_path, "ppt/charts/chart1.xml")
    output_chart_xml = _pptx_part(output_path, "ppt/charts/chart1.xml")
    assert output_chart_xml == template_chart_xml

    template_chart_rels = _pptx_part(template_path, "ppt/charts/_rels/chart1.xml.rels")
    output_chart_rels = _pptx_part(output_path, "ppt/charts/_rels/chart1.xml.rels")
    assert output_chart_rels == template_chart_rels


def test_build_chart_waterfall_connector_styles_change_overlay_segment_counts(
    tmp_path: Path,
) -> None:
    gap_solid_path = tmp_path / "waterfall-gap-solid.pptx"
    step_solid_path = tmp_path / "waterfall-step-solid.pptx"
    gap_dash_path = tmp_path / "waterfall-gap-dash.pptx"
    gap_dot_path = tmp_path / "waterfall-gap-dot.pptx"

    build_chart(Presentation(), _waterfall_overlay_spec("gap", "solid"), gap_solid_path)
    build_chart(Presentation(), _waterfall_overlay_spec("step", "solid"), step_solid_path)
    build_chart(Presentation(), _waterfall_overlay_spec("gap", "long_dash"), gap_dash_path)
    build_chart(Presentation(), _waterfall_overlay_spec("gap", "dot"), gap_dot_path)

    gap_solid_blank_shapes = _blank_auto_shape_count(gap_solid_path)
    step_solid_blank_shapes = _blank_auto_shape_count(step_solid_path)
    gap_dash_blank_shapes = _blank_auto_shape_count(gap_dash_path)
    gap_dot_blank_shapes = _blank_auto_shape_count(gap_dot_path)

    assert _chart_types(gap_solid_path) == [int(XL_CHART_TYPE.BAR_STACKED)]

    assert step_solid_blank_shapes > gap_solid_blank_shapes
    assert gap_dash_blank_shapes > gap_solid_blank_shapes
    assert gap_dot_blank_shapes > gap_dash_blank_shapes
