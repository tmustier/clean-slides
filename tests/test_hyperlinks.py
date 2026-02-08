"""Tests for clickable hyperlink generation from markdown link syntax."""

from __future__ import annotations

import os
import tempfile
from pathlib import Path

from lxml import etree
from pptx import Presentation

from clean_slides.cli import cmd_generate

NS_A = "http://schemas.openxmlformats.org/drawingml/2006/main"
NS_R = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
HYPERLINK_REL = "http://schemas.openxmlformats.org/officeDocument/2006/relationships/hyperlink"


def _generate_pptx(yaml_content: str) -> Path:
    """Write a YAML spec to a temp file, generate a PPTX, return its path."""
    with tempfile.TemporaryDirectory() as td:
        yaml_path = Path(td) / "test.yaml"
        yaml_path.write_text(yaml_content)
        out_path = Path(td) / "out.pptx"

        # Build a minimal argparse Namespace
        class Args:
            input = [str(yaml_path)]
            output = str(out_path)
            template = None
            config = None
            slide_index = None
            keep_existing = False
            detail = False

        cmd_generate(Args())  # type: ignore[arg-type]
        assert out_path.exists(), "PPTX not generated"
        # Copy to a non-temp location so it survives the context manager
        import shutil

        fd, tmp = tempfile.mkstemp(suffix=".pptx")
        os.close(fd)
        persist = Path(tmp)
        shutil.copy2(out_path, persist)
        return persist


def test_markdown_link_creates_hyperlink() -> None:
    """[text](url) in cell content should produce an <a:hlinkClick> element."""
    yaml = """\
title: Link Test
table:
  rows: 2
  cols: 1
  has_col_header: true
  col_headers: ["Info"]
  cells:
    - ["[Click me](https://example.com)"]
"""
    pptx_path = _generate_pptx(yaml)
    try:
        prs = Presentation(str(pptx_path))
        slide = prs.slides[0]

        # Check relationship was registered
        hyperlink_rels = {
            rel.rId: rel.target_ref
            for rel in slide.part.rels.values()
            if rel.reltype == HYPERLINK_REL
        }
        assert len(hyperlink_rels) >= 1, "No hyperlink relationships found"
        assert "https://example.com" in hyperlink_rels.values()

        # Check XML has hlinkClick with correct text
        found = False
        for shape in slide.shapes:
            for r_el in shape._element.iter(f"{{{NS_A}}}r"):
                hlink = r_el.find(f".//{{{NS_A}}}hlinkClick")
                t_el = r_el.find(f"{{{NS_A}}}t")
                if hlink is not None and t_el is not None and t_el.text == "Click me":
                    rid = hlink.get(f"{{{NS_R}}}id")
                    assert rid is not None
                    assert hyperlink_rels.get(rid) == "https://example.com"
                    found = True
        assert found, "hlinkClick element with 'Click me' text not found"
    finally:
        pptx_path.unlink(missing_ok=True)


def test_hyperlink_element_order_in_rpr() -> None:
    """hlinkClick must come after latin in rPr (OOXML sequence)."""
    yaml = """\
title: Order Test
table:
  rows: 2
  cols: 1
  has_col_header: true
  col_headers: ["Info"]
  cells:
    - ["[Link](https://example.com)"]
"""
    pptx_path = _generate_pptx(yaml)
    try:
        prs = Presentation(str(pptx_path))
        slide = prs.slides[0]

        for shape in slide.shapes:
            for rPr in shape._element.iter(f"{{{NS_A}}}rPr"):
                children = [etree.QName(c).localname for c in rPr]
                if "hlinkClick" in children and "latin" in children:
                    latin_idx = children.index("latin")
                    hlink_idx = children.index("hlinkClick")
                    assert (
                        hlink_idx > latin_idx
                    ), f"hlinkClick (pos {hlink_idx}) must come after latin (pos {latin_idx})"
    finally:
        pptx_path.unlink(missing_ok=True)


def test_plain_text_has_no_hyperlink() -> None:
    """Plain text without markdown links should not create hyperlink rels."""
    yaml = """\
title: No Link Test
table:
  rows: 2
  cols: 1
  has_col_header: true
  col_headers: ["Info"]
  cells:
    - ["Just plain text here"]
"""
    pptx_path = _generate_pptx(yaml)
    try:
        prs = Presentation(str(pptx_path))
        slide = prs.slides[0]

        hyperlink_rels = [rel for rel in slide.part.rels.values() if rel.reltype == HYPERLINK_REL]
        assert len(hyperlink_rels) == 0, "Unexpected hyperlink relationships found"
    finally:
        pptx_path.unlink(missing_ok=True)
