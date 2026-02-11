"""Shared CLI helpers for presentation loading, lookup, and JSON output."""

from __future__ import annotations

import json
import sys
from pathlib import Path
from typing import Protocol, Union

from pptx import Presentation
from pptx.presentation import Presentation as PresentationObj
from pptx.shapes.autoshape import Shape
from pptx.shapes.base import BaseShape
from pptx.slide import Slide, SlideLayout
from typing_extensions import TypeGuard


class ToDict(Protocol):
    def to_dict(self) -> object: ...


def has_to_dict(obj: object) -> TypeGuard[ToDict]:
    return hasattr(obj, "to_dict")


def is_object_list(value: object) -> TypeGuard[list[object]]:
    return isinstance(value, list)


def is_str_object_dict(value: object) -> TypeGuard[dict[str, object]]:
    return isinstance(value, dict)


def is_text_shape(shape: BaseShape) -> TypeGuard[Shape]:
    return isinstance(shape, Shape) and shape.has_text_frame


def open_presentation(path: Union[str, Path]) -> PresentationObj:
    """Open a PPTX file."""
    return Presentation(str(path))


def get_slide(prs: PresentationObj, num: Union[str, int]) -> Slide:
    """Get slide by 1-indexed number."""
    idx = int(num) - 1
    slides: list[Slide] = list(prs.slides)
    if idx < 0 or idx >= len(slides):
        print(f"Error: slide {num} out of range (1-{len(slides)})", file=sys.stderr)
        sys.exit(1)
    return slides[idx]


def find_shape(slide: Slide, identifier: str) -> BaseShape:
    """Find shape by index (int) or name (substring match)."""
    try:
        idx = int(identifier)
        shapes: list[BaseShape] = list(slide.shapes)
        if 0 <= idx < len(shapes):
            return shapes[idx]
        for shape in shapes:
            if shape.shape_id == idx:
                return shape
    except ValueError:
        pass

    ident_lower = identifier.lower()
    for shape in slide.shapes:
        if ident_lower in shape.name.lower():
            return shape

    print(f"Error: shape '{identifier}' not found", file=sys.stderr)
    sys.exit(1)


def _layout_by_index(prs: PresentationObj, identifier: str) -> SlideLayout | None:
    try:
        idx = int(identifier)
    except ValueError:
        return None

    layouts: list[SlideLayout] = list(prs.slide_layouts)
    if 0 <= idx < len(layouts):
        return layouts[idx]
    return None


def _layout_by_name(prs: PresentationObj, identifier: str) -> SlideLayout | None:
    ident_lower = identifier.lower()
    for layout in prs.slide_layouts:
        if ident_lower in layout.name.lower():
            return layout
    return None


def _fallback_layout(layouts: list[SlideLayout]) -> SlideLayout:
    # Prefer a sensible, content-friendly layout.
    # NOTE: templates vary widely; avoid hardcoding indices.
    for key in ("default", "blank"):
        exact = next((layout for layout in layouts if layout.name.lower() == key), None)
        if exact is not None:
            return exact

    for key in ("default", "blank"):
        partial = next((layout for layout in layouts if key in layout.name.lower()), None)
        if partial is not None:
            return partial

    return layouts[0]


def find_layout(prs: PresentationObj, identifier: str, fallback: bool = False) -> SlideLayout:
    """Find layout by index (int) or name (substring match).

    If fallback=True, returns a best-effort layout instead of exiting.
    """
    by_index = _layout_by_index(prs, identifier)
    if by_index is not None:
        return by_index

    by_name = _layout_by_name(prs, identifier)
    if by_name is not None:
        return by_name

    if fallback:
        return _fallback_layout(list(prs.slide_layouts))

    print(f"Error: layout '{identifier}' not found", file=sys.stderr)
    sys.exit(1)


def try_find_layout(prs: PresentationObj, name: str) -> SlideLayout | None:
    """Best-effort lookup by layout name (case-insensitive).

    Returns None if no match.
    """
    name_lower = name.lower()
    layouts: list[SlideLayout] = list(prs.slide_layouts)

    for layout in layouts:
        if layout.name and layout.name.lower() == name_lower:
            return layout

    for layout in layouts:
        if layout.name and name_lower in layout.name.lower():
            return layout

    return None


def json_out(obj: object) -> None:
    """Serialize to JSON, handling dataclass .to_dict()."""
    payload: object

    if has_to_dict(obj):
        payload = obj.to_dict()
    elif is_object_list(obj):
        items: list[object] = []
        for x_obj in obj:
            if has_to_dict(x_obj):
                items.append(x_obj.to_dict())
            else:
                items.append(x_obj)
        payload = items
    else:
        payload = obj

    print(json.dumps(payload, indent=2, ensure_ascii=False, default=str))


def shape_text(shape: BaseShape) -> str | None:
    """Current text from a shape."""
    if not is_text_shape(shape):
        return None
    return shape.text_frame.text.replace("\x0b", "\\n").replace("\n", "\\n")
