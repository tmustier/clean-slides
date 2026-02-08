"""Safe writing to placeholders and text frames.

Core principle: only set properties you explicitly pass. Everything else
inherits from the placeholder → layout → master → theme chain.

Progressive disclosure:
    write_to_placeholder(ph, "Simple text")
    write_to_placeholder(ph, [("Bold ", {"bold": True}), ("normal", {})])
    write_to_placeholder(ph, [("Colored", {"color": "accent1", "size": 25})])
"""

from __future__ import annotations

from collections import Counter
from typing import Any, List, Mapping, Optional, Protocol, Sequence, Tuple, TypedDict, Union

from lxml import etree
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.enum.text import MSO_VERTICAL_ANCHOR, PP_ALIGN, PP_PARAGRAPH_ALIGNMENT
from pptx.oxml.ns import qn
from pptx.oxml.text import CT_RegularTextRun, CT_TextParagraph
from pptx.shapes.autoshape import Shape
from pptx.text.text import Font, TextFrame
from pptx.util import Inches, Length, Pt
from typing_extensions import TypeGuard

# ── Theme color string → scheme value mapping ─────────────────────────

_SCHEME_MAP: dict[str, str] = {
    "accent1": "accent1",
    "accent2": "accent2",
    "accent3": "accent3",
    "accent4": "accent4",
    "accent5": "accent5",
    "accent6": "accent6",
    "dk1": "dk1",
    "dk2": "dk2",
    "lt1": "lt1",
    "lt2": "lt2",
    "tx1": "dk1",
    "tx2": "dk2",
    "bg1": "lt1",
    "bg2": "lt2",
    "hlink": "hlink",
    "folHlink": "folHlink",
}

_THEME_COLOR_MAP: dict[str, MSO_THEME_COLOR] = {
    "accent1": MSO_THEME_COLOR.ACCENT_1,
    "accent2": MSO_THEME_COLOR.ACCENT_2,
    "accent3": MSO_THEME_COLOR.ACCENT_3,
    "accent4": MSO_THEME_COLOR.ACCENT_4,
    "accent5": MSO_THEME_COLOR.ACCENT_5,
    "accent6": MSO_THEME_COLOR.ACCENT_6,
    "dk1": MSO_THEME_COLOR.DARK_1,
    "dk2": MSO_THEME_COLOR.DARK_2,
    "lt1": MSO_THEME_COLOR.LIGHT_1,
    "lt2": MSO_THEME_COLOR.LIGHT_2,
    "tx1": MSO_THEME_COLOR.DARK_1,
    "tx2": MSO_THEME_COLOR.DARK_2,
    "bg1": MSO_THEME_COLOR.LIGHT_1,
    "bg2": MSO_THEME_COLOR.LIGHT_2,
    "hlink": MSO_THEME_COLOR.HYPERLINK,
    "folHlink": MSO_THEME_COLOR.FOLLOWED_HYPERLINK,
}


def _is_dict(value: object) -> TypeGuard[dict[str, Any]]:
    return isinstance(value, dict)


def _is_sequence(value: object) -> TypeGuard[Sequence[object]]:
    return isinstance(value, (list, tuple))


def _rgb_from_hex(value: str) -> RGBColor:
    rgb_hex_str = value.lstrip("#")
    if len(rgb_hex_str) != 6:
        raise ValueError(f"Invalid RGB hex string: {value!r}")
    r = int(rgb_hex_str[0:2], 16)
    g = int(rgb_hex_str[2:4], 16)
    b = int(rgb_hex_str[4:6], 16)
    return RGBColor(r, g, b)


def _to_float(value: object) -> float:
    if isinstance(value, bool):
        return float(int(value))
    if isinstance(value, (int, float)):
        return float(value)
    return float(str(value))


class RunOverrides(TypedDict, total=False):
    color: object  # str | RGBColor | MSO_THEME_COLOR
    bold: bool
    italic: bool
    underline: bool
    size: float  # points
    font: str
    superscript: bool
    subscript: bool


RunContent = object


class ParagraphSpec(TypedDict, total=False):
    runs: RunContent
    level: int
    alignment: Union[str, PP_PARAGRAPH_ALIGNMENT]
    spacing_before: float  # points
    spacing_after: float  # points
    line_spacing: float  # percent, e.g. 114
    bullet: Union[bool, str]


class RunLike(Protocol):
    @property
    def text(self) -> str: ...

    @text.setter
    def text(self, text: str) -> None: ...

    @property
    def font(self) -> Font: ...


class ParagraphLike(Protocol):
    def add_run(self) -> RunLike: ...

    def add_line_break(self) -> None: ...

    @property
    def runs(self) -> tuple[RunLike, ...]: ...

    @property
    def level(self) -> int: ...

    @level.setter
    def level(self, level: int) -> None: ...

    @property
    def alignment(self) -> Optional[PP_PARAGRAPH_ALIGNMENT]: ...

    @alignment.setter
    def alignment(self, value: Optional[PP_PARAGRAPH_ALIGNMENT]) -> None: ...

    @property
    def space_before(self) -> Optional[Length]: ...

    @space_before.setter
    def space_before(self, value: Optional[Length]) -> None: ...

    @property
    def space_after(self) -> Optional[Length]: ...

    @space_after.setter
    def space_after(self, value: Optional[Length]) -> None: ...

    @property
    def line_spacing(self) -> Optional[Union[int, float, Length]]: ...

    @line_spacing.setter
    def line_spacing(self, value: Optional[Union[int, float, Length]]) -> None: ...


# ── Run-level overrides ───────────────────────────────────────────────


def _parse_run_overrides(raw: object) -> RunOverrides:
    if not _is_dict(raw):
        return {}

    overrides: RunOverrides = {}

    if "bold" in raw:
        overrides["bold"] = bool(raw["bold"])
    if "italic" in raw:
        overrides["italic"] = bool(raw["italic"])
    if "underline" in raw:
        overrides["underline"] = bool(raw["underline"])

    if "size" in raw:
        overrides["size"] = _to_float(raw["size"])
    if "font" in raw and raw["font"] is not None:
        overrides["font"] = str(raw["font"])

    if "color" in raw:
        overrides["color"] = raw["color"]

    if "superscript" in raw:
        overrides["superscript"] = bool(raw["superscript"])
    if "subscript" in raw:
        overrides["subscript"] = bool(raw["subscript"])

    return overrides


def _set_run_baseline(run: RunLike, baseline: str) -> None:
    """Set baseline attribute on the run XML (used for super/subscript)."""
    r_obj = getattr(run, "_r", None)
    if not isinstance(r_obj, CT_RegularTextRun):
        return
    rPr = r_obj.get_or_add_rPr()
    rPr.set("baseline", baseline)


def _apply_run_overrides(run: RunLike, overrides: RunOverrides) -> None:
    """Apply formatting overrides to a run.

    Override keys:
        color       "accent1" | "#FF0000" | RGBColor | MSO_THEME_COLOR
        bold        bool
        italic      bool
        underline   bool
        size        float (points)
        font        str (font name)
        superscript bool
        subscript   bool
    """
    if not overrides:
        return

    if "bold" in overrides:
        run.font.bold = overrides["bold"]

    if "italic" in overrides:
        run.font.italic = overrides["italic"]

    if "underline" in overrides:
        run.font.underline = overrides["underline"]

    if "size" in overrides:
        run.font.size = Pt(float(overrides["size"]))

    if "font" in overrides:
        run.font.name = overrides["font"]

    if "color" in overrides:
        _apply_color(run, overrides["color"])

    if overrides.get("superscript"):
        _set_run_baseline(run, "30000")

    if overrides.get("subscript"):
        _set_run_baseline(run, "-25000")


def _apply_color(run: RunLike, color: object) -> None:
    """Apply color to a run.

    Accepts:
        "accent1"       → theme color
        "#FF0000"       → RGB hex
        RGBColor(...)   → direct
        MSO_THEME_COLOR → enum
    """

    if isinstance(color, MSO_THEME_COLOR):
        run.font.color.theme_color = color
        return

    if isinstance(color, RGBColor):
        run.font.color.rgb = color
        return

    if not isinstance(color, str):
        return

    if color.startswith("#"):
        run.font.color.rgb = _rgb_from_hex(color)
        return

    # Theme scheme name
    scheme = _SCHEME_MAP.get(color)
    if scheme is not None:
        theme_color = _THEME_COLOR_MAP.get(color)
        if theme_color is not None:
            run.font.color.theme_color = theme_color
        return

    # Try as hex without '#'
    run.font.color.rgb = _rgb_from_hex(color)


def _normalize_content(content: Any) -> List[Tuple[str, RunOverrides]]:
    """Normalize content to list of (text, overrides) tuples.

    "hello"                              → [("hello", {})]
    [("bold", {"bold": True})]           → as-is
    [["bold", {"bold": True}]]           → same (from JSON)
    [["just text"]]                      → [("just text", {})]
    ["just text"]                        → [("just text", {})]
    [{"text": "bold", "bold": True}]     → inspect output format
    """

    if isinstance(content, str):
        return [(content, {})]

    if not _is_sequence(content):
        return [(str(content), {})]

    result: List[Tuple[str, RunOverrides]] = []
    for item in content:
        if isinstance(item, str):
            result.append((item, {}))
            continue

        # Dict format: {"text": "...", "bold": true, ...} — matches inspect output
        if _is_dict(item) and "text" in item:
            text = str(item["text"])
            # Everything except "text" is an override
            opts_raw = {k: v for k, v in item.items() if k != "text"}
            result.append((text, _parse_run_overrides(opts_raw)))
            continue

        # Array format: ["text", {"bold": true}] — legacy format
        if _is_sequence(item):
            text = str(item[0]) if len(item) > 0 else ""
            opts_raw: object = item[1] if len(item) > 1 else {}
            result.append((text, _parse_run_overrides(opts_raw)))
            continue

        result.append((str(item), {}))

    return result


# ── Public API ─────────────────────────────────────────────────────────


def write_to_placeholder(placeholder: Shape, content: object, **frame_overrides: object) -> None:
    """Write content to a placeholder, preserving inherited styling."""

    tf: TextFrame = placeholder.text_frame
    tf.clear()
    p: ParagraphLike = tf.paragraphs[0]

    for text, overrides in _normalize_content(content):
        run: RunLike = p.add_run()
        run.text = text
        _apply_run_overrides(run, overrides)

    _apply_frame_overrides(tf, frame_overrides)


def set_text(shape: Shape, content: object, **frame_overrides: object) -> None:
    """Write content to any shape's text frame (textbox, rectangle, etc.)."""

    tf: TextFrame = shape.text_frame
    tf.clear()
    p: ParagraphLike = tf.paragraphs[0]

    for text, overrides in _normalize_content(content):
        run: RunLike = p.add_run()
        run.text = text
        _apply_run_overrides(run, overrides)

    _apply_frame_overrides(tf, frame_overrides)


def _snapshot_defaults(tf: TextFrame) -> RunOverrides:
    """Capture the dominant run formatting from an existing text frame."""

    sizes: Counter[int] = Counter()
    fonts: Counter[str] = Counter()

    for p in tf.paragraphs:
        for run in p.runs:
            size = run.font.size
            if size is not None:
                sizes[int(size)] += 1

            name = run.font.name
            if name is not None:
                fonts[name] += 1

    defaults: RunOverrides = {}
    if sizes:
        # Convert EMU to points for the override key
        most_common_size_emu = sizes.most_common(1)[0][0]
        defaults["size"] = most_common_size_emu / 12_700

    if fonts:
        defaults["font"] = fonts.most_common(1)[0][0]

    return defaults


def _apply_defaults(overrides: RunOverrides, defaults: RunOverrides) -> RunOverrides:
    """Merge defaults into overrides — user overrides win."""
    merged: RunOverrides = {}
    merged.update(defaults)
    merged.update(overrides)
    return merged


def snapshot_defaults(tf: TextFrame) -> RunOverrides:
    """Capture the dominant run formatting from an existing text frame."""
    return _snapshot_defaults(tf)


def apply_defaults(overrides: RunOverrides, defaults: RunOverrides) -> RunOverrides:
    """Merge defaults into overrides — user overrides win."""
    return _apply_defaults(overrides, defaults)


def write_paragraphs(shape: Shape, paragraphs: Sequence[ParagraphSpec]) -> None:
    """Write multiple paragraphs to a shape's text frame.

    Captures font defaults (size, name) from existing content before clearing,
    so new runs match the original formatting unless explicitly overridden.
    """

    tf: TextFrame = shape.text_frame
    defaults = _snapshot_defaults(tf)
    tf.clear()

    for i, para_spec in enumerate(paragraphs):
        p: ParagraphLike = tf.paragraphs[0] if i == 0 else tf.add_paragraph()

        # Level
        level_raw = para_spec.get("level", 0)
        p.level = int(level_raw)

        # Alignment
        alignment_raw = para_spec.get("alignment")
        if alignment_raw:
            if isinstance(alignment_raw, str):
                align_map: dict[str, PP_PARAGRAPH_ALIGNMENT] = {
                    "left": PP_ALIGN.LEFT,
                    "center": PP_ALIGN.CENTER,
                    "right": PP_ALIGN.RIGHT,
                    "justify": PP_ALIGN.JUSTIFY,
                }
                align = align_map.get(alignment_raw)
                if align is not None:
                    p.alignment = align
            else:
                # After the string check, this is already narrowed to the enum.
                p.alignment = alignment_raw

        # Spacing
        sb = para_spec.get("spacing_before")
        if sb is not None:
            p.space_before = Pt(_to_float(sb))

        sa = para_spec.get("spacing_after")
        if sa is not None:
            p.space_after = Pt(_to_float(sa))

        ls = para_spec.get("line_spacing")
        if ls is not None:
            # Input is percent, e.g. 114 → 1.14
            p.line_spacing = _to_float(ls) / 100.0

        # Bullet control (requires XML access)
        bullet = para_spec.get("bullet")
        if bullet is not None:
            p_elm = getattr(p, "_element", None)
            if isinstance(p_elm, CT_TextParagraph):
                pPr = p_elm.get_or_add_pPr()

                # Remove any existing bullet elements
                for tag in ("a:buNone", "a:buChar", "a:buAutoNum"):
                    for el in pPr.findall(qn(tag)):
                        pPr.remove(el)

                if bullet is False:
                    etree.SubElement(pPr, qn("a:buNone"))
                elif isinstance(bullet, str):
                    etree.SubElement(pPr, qn("a:buChar")).set("char", bullet)

        # Runs — apply snapshotted defaults
        runs_spec: object = para_spec.get("runs", "")
        normalized = _normalize_content(runs_spec)
        for text, overrides in normalized:
            run: RunLike = p.add_run()
            run.text = text
            _apply_run_overrides(run, _apply_defaults(overrides, defaults))


def add_line_break(paragraph: ParagraphLike) -> ParagraphLike:
    """Insert <a:br/> into a paragraph."""
    paragraph.add_line_break()
    return paragraph


def add_run(paragraph: ParagraphLike, text: str, **overrides: object) -> RunLike:
    """Add a run to an existing paragraph with optional overrides."""

    run: RunLike = paragraph.add_run()
    run.text = text
    _apply_run_overrides(run, _parse_run_overrides(overrides))
    return run


# ── Frame-level overrides ─────────────────────────────────────────────


def _apply_frame_overrides(tf: TextFrame, overrides: Mapping[str, object]) -> None:
    """Apply text frame-level overrides."""
    if not overrides:
        return

    if "anchor" in overrides:
        anchor_raw = overrides.get("anchor")
        if isinstance(anchor_raw, MSO_VERTICAL_ANCHOR):
            tf.vertical_anchor = anchor_raw
        elif isinstance(anchor_raw, str):
            anchor_map: dict[str, MSO_VERTICAL_ANCHOR] = {
                "t": MSO_VERTICAL_ANCHOR.TOP,
                "ctr": MSO_VERTICAL_ANCHOR.MIDDLE,
                "b": MSO_VERTICAL_ANCHOR.BOTTOM,
            }
            anchor = anchor_map.get(anchor_raw)
            if anchor is not None:
                tf.vertical_anchor = anchor

    if "word_wrap" in overrides:
        ww = overrides.get("word_wrap")
        if isinstance(ww, bool) or ww is None:
            tf.word_wrap = ww

    if "margins" in overrides:
        margins_raw = overrides.get("margins")
        if _is_dict(margins_raw):
            if "left" in margins_raw:
                tf.margin_left = Inches(_to_float(margins_raw["left"]))
            if "right" in margins_raw:
                tf.margin_right = Inches(_to_float(margins_raw["right"]))
            if "top" in margins_raw:
                tf.margin_top = Inches(_to_float(margins_raw["top"]))
            if "bottom" in margins_raw:
                tf.margin_bottom = Inches(_to_float(margins_raw["bottom"]))

    if "alignment" in overrides:
        alignment_raw = overrides.get("alignment")
        if alignment_raw:
            if isinstance(alignment_raw, str):
                align_map: dict[str, PP_PARAGRAPH_ALIGNMENT] = {
                    "left": PP_ALIGN.LEFT,
                    "center": PP_ALIGN.CENTER,
                    "right": PP_ALIGN.RIGHT,
                    "justify": PP_ALIGN.JUSTIFY,
                }
                align = align_map.get(alignment_raw)
                if align is None:
                    return
                for p in tf.paragraphs:
                    p.alignment = align
            elif isinstance(alignment_raw, PP_PARAGRAPH_ALIGNMENT):
                for p in tf.paragraphs:
                    p.alignment = alignment_raw
