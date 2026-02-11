"""
Unified CLI for PowerPoint inspection, editing, generation, and rendering.

    pptx <command> [args]

Inspect (progressive drill-down):
    pptx show          <file>                    # slide index with layout + word count
    pptx show          <file> <slide>            # all shapes sorted by position
    pptx show          <file> <slide> <shape>    # full detail (text/formatting/chart data)

Other inspect:
    pptx theme         <file>                    # scheme colors
    pptx xml           <file> <slide> <shape>    # raw XML

Edit:
    pptx edit          <file> <slide> <shape> <text> [--out PATH]
    pptx batch         <file> <edits.json | -> [--out PATH]

Slide management:
    pptx add-slide     <file> <layout> [--at N] [--out PATH]
    pptx delete-slide  <file> <slide> --confirm [--out PATH]
    pptx delete-shape  <file> <slide> <shape> [--out PATH]
    pptx insert        <deck.pptx> <source.pptx> [--at N] [--slides 1,3-5] [--out PATH]

Render:
    pptx render        <file> <slides> [--out DIR] [--dpi N] [--engine E]
    pptx crop          <png> <L> <T> <R> <B> [--out PATH]

Charts (from JSON):
    pptx charts        <spec.json> <output.pptx> [--template PATH] [--layout NAME]

Setup:
    pptx init          [-t template.pptx]        # create .clean-slides/ project dir
    pptx init-config   <template.pptx>           # generate config from a template

Generate (from YAML):
    pptx generate      <yaml...> [-o out.pptx] [-t template.pptx]
    pptx validate      <yaml...>
    pptx verify        <yaml...>

Auto-discovery: when no -t/-c flags are given, looks for .clean-slides/template.pptx
and .clean-slides/config.yaml walking up from the current directory.

Shape identified by index (int) or name (substring match).
Slide numbers are 1-indexed for inspect/edit, 0-indexed for generate --slide-index.
"""

from __future__ import annotations

import argparse
import json
import sys
from collections.abc import Iterable
from copy import deepcopy
from pathlib import Path
from typing import Any, Protocol, Union, cast

from pptx import Presentation
from pptx.presentation import Presentation as PresentationObj
from pptx.shapes.autoshape import Shape
from pptx.shapes.base import BaseShape
from pptx.slide import Slide, SlideLayout
from pptx.util import Emu
from typing_extensions import TypeGuard

from .cli_text import text_preview as _text_preview
from .cli_text import write_to_shape as _write_to_shape
from .constants import EMU_PER_INCH, Fonts, FontSizes, Layout, TableDefaults
from .metadata import fill_slide_metadata
from .placeholder import fill_placeholders
from .renderer import TableRenderer
from .screenshot import ScreenshotGenerator, crop_region
from .solver import ConstraintSolver
from .spec import ContentArea, TableLayout, TableSpec
from .spec_pipeline import YamlDict, load_yaml, parse_spec, preview_spec, validate_spec
from .template_config import TEMPLATE_CONFIG, set_template_config
from .text_metrics import EMU_PER_PT, TextMetrics

# ============================================================================
# SHARED HELPERS
# ============================================================================


class _ToDict(Protocol):
    def to_dict(self) -> object: ...


def _has_to_dict(obj: object) -> TypeGuard[_ToDict]:
    return hasattr(obj, "to_dict")


def _is_object_list(value: object) -> TypeGuard[list[object]]:
    return isinstance(value, list)


def _is_str_object_dict(value: object) -> TypeGuard[dict[str, object]]:
    return isinstance(value, dict)


def _is_text_shape(shape: BaseShape) -> TypeGuard[Shape]:
    return isinstance(shape, Shape) and shape.has_text_frame


class _FileArgs(Protocol):
    file: str


class _ShowArgs(_FileArgs, Protocol):
    slide: str | None
    shape: str | None


class _SlideArgs(_FileArgs, Protocol):
    slide: str


class _ColorArgs(_FileArgs, Protocol):
    hex: str


class _XmlArgs(_FileArgs, Protocol):
    slide: str
    shape: str


class _EditArgs(_FileArgs, Protocol):
    slide: str
    shape: str
    text: str
    out: str | None


class _BatchArgs(_FileArgs, Protocol):
    edits: str
    out: str | None


class _AddSlideArgs(_FileArgs, Protocol):
    layout: str
    at: int | None
    out: str | None


class _DeleteSlideArgs(_FileArgs, Protocol):
    slide: str
    confirm: bool
    out: str | None


class _DeleteShapeArgs(_FileArgs, Protocol):
    slide: str
    shape: str
    out: str | None


class _InsertArgs(_FileArgs, Protocol):
    source: str
    at: int | None
    slides: str | None
    out: str | None


class _RenderArgs(_FileArgs, Protocol):
    slides: str
    dpi: int
    out: str | None
    engine: str | None


class _CropArgs(Protocol):
    png: str
    left: float
    top: float
    right: float
    bottom: float
    out: str | None


class _ChartsArgs(Protocol):
    input: str
    output: str
    template: str | None
    layout: str | None
    expected_template: str | None


class _InputArgs(Protocol):
    input: list[str]


class _GenerateArgs(_InputArgs, Protocol):
    output: str | None
    template: str | None
    config: str | None
    slide_index: int | None
    keep_existing: bool
    detail: bool


class _VerifyArgs(_InputArgs, Protocol):
    detail: bool
    json: str | None
    config: str | None


class _ScreenshotArgs(_InputArgs, Protocol):
    output_dir: str | None
    slide: int
    soffice: str | None


class _ValidateArgs(_InputArgs, Protocol):
    config: str | None


class _InitArgs(Protocol):
    template: str | None
    output: str | None


class _InitConfigArgs(Protocol):
    file: str
    output: str | None


# ── Project-level auto-discovery ───────────────────────────────────────

PROJECT_DIR_NAME = ".clean-slides"
_CONFIG_NAME = "config.yaml"
_TEMPLATE_NAME = "template.pptx"


def _discover_project_dir() -> Path | None:
    """Walk from CWD to filesystem root looking for a `.clean-slides/` directory."""
    cur = Path.cwd().resolve()
    for parent in [cur, *cur.parents]:
        candidate = parent / PROJECT_DIR_NAME
        if candidate.is_dir():
            return candidate
    return None


def _discover_config() -> Path | None:
    """Return the project config path if auto-discovered."""
    proj = _discover_project_dir()
    if proj is not None:
        cfg = proj / _CONFIG_NAME
        if cfg.is_file():
            return cfg
    return None


def _discover_template() -> Path | None:
    """Return the project template path if auto-discovered."""
    proj = _discover_project_dir()
    if proj is not None:
        tpl = proj / _TEMPLATE_NAME
        if tpl.is_file():
            return tpl
    return None


def _apply_config(config_path: str | None) -> None:
    """Load template config — explicit path, auto-discovered, or built-in defaults."""
    if config_path is not None:
        set_template_config(Path(config_path))
        return
    # Auto-discover from .clean-slides/
    discovered = _discover_config()
    if discovered is not None:
        set_template_config(discovered)


def _open(path: Union[str, Path]) -> PresentationObj:
    """Open a PPTX file."""
    return Presentation(str(path))


def _get_slide(prs: PresentationObj, num: Union[str, int]) -> Slide:
    """Get slide by 1-indexed number."""
    idx = int(num) - 1
    slides: list[Slide] = list(prs.slides)
    if idx < 0 or idx >= len(slides):
        print(f"Error: slide {num} out of range (1-{len(slides)})", file=sys.stderr)
        sys.exit(1)
    return slides[idx]


def _find_shape(slide: Slide, identifier: str) -> BaseShape:
    """Find shape by index (int) or name (substring match)."""
    try:
        idx = int(identifier)
        shapes: list[BaseShape] = list(slide.shapes)
        if 0 <= idx < len(shapes):
            return shapes[idx]
        for s in shapes:
            if s.shape_id == idx:
                return s
    except ValueError:
        pass

    ident_lower = identifier.lower()
    for s in slide.shapes:
        if ident_lower in s.name.lower():
            return s

    print(f"Error: shape '{identifier}' not found", file=sys.stderr)
    sys.exit(1)


def _layout_by_index(prs: PresentationObj, identifier: str) -> SlideLayout | None:
    try:
        idx = int(identifier)
    except ValueError:
        return None

    layouts: list[SlideLayout] = list(prs.slide_layouts)
    if 0 <= idx < len(layouts):
        return layouts[idx]
    return None


def _layout_by_name(prs: PresentationObj, identifier: str) -> SlideLayout | None:
    ident_lower = identifier.lower()
    for layout in prs.slide_layouts:
        if ident_lower in layout.name.lower():
            return layout
    return None


def _fallback_layout(layouts: list[SlideLayout]) -> SlideLayout:
    # Prefer a sensible, content-friendly layout.
    # NOTE: templates vary widely; avoid hardcoding indices.
    for key in ("default", "blank"):
        exact = next((layout for layout in layouts if layout.name.lower() == key), None)
        if exact is not None:
            return exact

    for key in ("default", "blank"):
        partial = next((layout for layout in layouts if key in layout.name.lower()), None)
        if partial is not None:
            return partial

    return layouts[0]


def _find_layout(prs: PresentationObj, identifier: str, fallback: bool = False) -> SlideLayout:
    """Find layout by index (int) or name (substring match).

    If fallback=True, returns a best-effort layout instead of exiting.
    """
    by_index = _layout_by_index(prs, identifier)
    if by_index is not None:
        return by_index

    by_name = _layout_by_name(prs, identifier)
    if by_name is not None:
        return by_name

    if fallback:
        return _fallback_layout(list(prs.slide_layouts))

    print(f"Error: layout '{identifier}' not found", file=sys.stderr)
    sys.exit(1)


def _try_find_layout(prs: PresentationObj, name: str) -> SlideLayout | None:
    """Best-effort lookup by layout name (case-insensitive).

    Returns None if no match.
    """
    name_lower = name.lower()
    layouts: list[SlideLayout] = list(prs.slide_layouts)

    for layout in layouts:
        if layout.name and layout.name.lower() == name_lower:
            return layout

    for layout in layouts:
        if layout.name and name_lower in layout.name.lower():
            return layout

    return None


def _content_area_from_layout(slide_layout: SlideLayout) -> ContentArea | None:
    """Extract the primary content area from a slide layout.

    Finds the best OBJECT / BODY placeholder to use as the table content
    area.  When multiple content placeholders exist (e.g. "Two Content"),
    picks the largest; ties are broken by topmost then leftmost position
    so the primary (upper-left) area wins consistently.

    Returns ``None`` when no suitable placeholder exists.
    """
    from .spec import ContentArea

    # Placeholder type names that represent content areas (OBJECT, BODY,
    # TABLE, CHART, etc.).  We exclude TITLE, SUBTITLE, DATE, FOOTER,
    # SLIDE_NUMBER, HEADER.
    _SKIP_TYPES = {"TITLE", "CENTER_TITLE", "SUBTITLE", "DATE", "FOOTER", "SLIDE_NUMBER", "HEADER"}

    footer_y = int(Layout.FOOTER_LINE_Y)

    # Collect candidate content placeholders.
    candidates: list[tuple[int, int, int, ContentArea]] = []
    for ph in slide_layout.placeholders:
        pf = ph.placeholder_format
        type_name = str(pf.type).split("(")[0].strip() if pf.type is not None else ""
        if type_name in _SKIP_TYPES:
            continue
        ph_area = int(ph.width) * int(ph.height)
        ph_bottom = int(ph.top) + int(ph.height)
        area = ContentArea(
            x=int(ph.left),
            y=int(ph.top),
            width=int(ph.width),
            height=int(min(ph_bottom, footer_y) - int(ph.top)),
        )
        candidates.append((ph_area, int(ph.top), int(ph.left), area))

    if not candidates:
        return None

    # Pick the primary content area: topmost then leftmost among the
    # largest placeholders (within 10% of the max area).
    max_area = max(c[0] for c in candidates)
    threshold = int(max_area * 0.9)
    large = [(top, left, ca) for (a, top, left, ca) in candidates if a >= threshold]
    large.sort()  # topmost, then leftmost
    return large[0][2]


def _sidebar_content_area(slide_layout: SlideLayout) -> ContentArea | None:
    """Extract the secondary (sidebar) content area from a split slide layout.

    Returns the ContentArea for the right-side placeholder in layouts like
    2/3, 3/4, 1/2. Returns None when the layout has no secondary area.
    """
    from .spec import ContentArea

    content_y_threshold = 1600000
    footer_y = int(Layout.FOOTER_LINE_Y)

    # Collect (left, top, width) for content-region placeholders
    candidates: list[tuple[int, int, int]] = []
    for ph in slide_layout.placeholders:  # type: ignore[union-attr]
        top: int = int(ph.top)  # type: ignore[arg-type]
        if top < content_y_threshold:
            continue
        candidates.append((int(ph.left), top, int(ph.width)))  # type: ignore[arg-type]

    if len(candidates) < 2:
        return None

    candidates.sort()
    x, y, w = candidates[1]
    return ContentArea(x=x, y=y, width=w, height=int(footer_y - y))


def _json_out(obj: object) -> None:
    """Serialize to JSON, handling dataclass .to_dict()."""

    payload: object

    if _has_to_dict(obj):
        payload = obj.to_dict()
    elif _is_object_list(obj):
        items: list[object] = []
        for x_obj in obj:
            if _has_to_dict(x_obj):
                items.append(x_obj.to_dict())
            else:
                items.append(x_obj)
        payload = items
    else:
        payload = obj

    print(json.dumps(payload, indent=2, ensure_ascii=False, default=str))


def _shape_text(shape: BaseShape) -> str | None:
    """Current text from a shape."""
    if not _is_text_shape(shape):
        return None
    return shape.text_frame.text.replace("\x0b", "\\n").replace("\n", "\\n")


# ============================================================================
# TEXT PARSING (for edit/batch commands)
# ============================================================================


# Text parsing and writing helpers moved to clean_slides.cli_text

def _get_text_limit(key: str, default: int) -> int:
    try:
        limits = TEMPLATE_CONFIG.section("text_limits")
    except KeyError:
        return default

    raw = limits.get(key)
    if raw is None:
        return default

    try:
        return int(raw)
    except (TypeError, ValueError):
        return default


_SIDEBAR_MIN_PT = 8  # never shrink below this


def _warn_sidebar_overflow(
    paragraphs: list[Any],
    area: ContentArea,
    metrics: TextMetrics,
) -> None:
    """Warn when sidebar content exceeds available height."""
    total = _sidebar_height(paragraphs, area, metrics)
    if total > area.height:
        overflow_in = (total - area.height) / EMU_PER_INCH
        print(
            f"  WARNING: sidebar content overflows by ~{overflow_in:.1f}in. "
            f"Shorten text, reduce paragraphs, or set sidebar_shrink: true.",
            file=sys.stderr,
        )


def _sidebar_height(
    paragraphs: list[Any],
    area: ContentArea,
    metrics: TextMetrics,
) -> int:
    """Estimate total sidebar height in EMU."""
    PARA_GAP_EMU = int(4 * EMU_PER_PT)
    total = 0
    for para in paragraphs:
        font: str = str(para.font or Fonts.BODY)
        size_pt: int = int(para.size_pt or FontSizes.DEFAULT)
        total += metrics.text_height(str(para.text), area.width, font, size_pt) + PARA_GAP_EMU
    return total


def _shrink_sidebar_to_fit(
    paragraphs: list[Any],
    area: ContentArea,
    metrics: TextMetrics,
) -> None:
    """Proportionally shrink sidebar font sizes until content fits the area.

    Modifies paragraph objects in-place. Warns if content still overflows
    at the minimum font size.
    """
    total = _sidebar_height(paragraphs, area, metrics)
    if total <= area.height:
        return  # fits already

    # Compute scale factor and apply proportionally
    scale = area.height / total
    original_sizes: list[int] = []
    for para in paragraphs:
        orig = int(para.size_pt or FontSizes.DEFAULT)
        original_sizes.append(orig)
        shrunk = max(_SIDEBAR_MIN_PT, int(orig * scale))
        para.size_pt = shrunk

    # Re-check (rounding / min clamp may still overflow)
    total = _sidebar_height(paragraphs, area, metrics)
    if total > area.height:
        overflow_in = (total - area.height) / EMU_PER_INCH
        print(
            f"  WARNING: sidebar content still overflows by ~{overflow_in:.1f}in after "
            f"shrinking fonts (min {_SIDEBAR_MIN_PT}pt). Shorten text or reduce paragraphs.",
            file=sys.stderr,
        )
    else:
        reduced = [
            f"{orig}→{int(para.size_pt or orig)}pt"
            for para, orig in zip(paragraphs, original_sizes)
            if int(para.size_pt or orig) != orig
        ]
        if reduced:
            print(
                f"  sidebar: shrunk fonts to fit ({reduced[0].split('→')[1].rstrip('pt')}pt body)"
            )


def _warn_placeholder_text_limits(slide: Slide, shape: Shape) -> None:
    """Warn (to stderr) when placeholder text likely wraps beyond configured max lines."""
    if not shape.is_placeholder:
        return

    ph_idx = shape.placeholder_format.idx
    if ph_idx not in {0, 1}:
        return

    text = str(shape.text_frame.text or "").strip()
    if not text:
        return

    is_title_slide = "title" in (slide.slide_layout.name or "").lower()

    if ph_idx == 0:
        max_lines = _get_text_limit("title_max_lines", 2)
        font = Fonts.HEADLINE
        size_pt = int(FontSizes.TITLE)
        label = "title"
    else:
        if is_title_slide:
            max_lines = _get_text_limit("title_slide_subtitle_max_lines", 1)
        else:
            max_lines = _get_text_limit("subtitle_max_lines", 1)
        font = Fonts.HEADLINE
        size_pt = int(FontSizes.SUBTITLE)
        label = "subtitle"

    width_emu = int(shape.width)
    metrics = TextMetrics()
    needed = metrics.lines_needed(text, width_emu, font, size_pt)

    if needed > max_lines:
        print(
            f"Warning: {label} text likely wraps to ~{needed} lines (max {max_lines}) in placeholder '{shape.name}'. "
            "Consider shortening.",
            file=sys.stderr,
        )


# ============================================================================
# INSPECT COMMANDS
# ============================================================================


def _classify_shape_type(shape: BaseShape) -> str:
    """Classify shape for display: placeholder, chart, image, table, group, text, connector, decorative."""
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    if shape.is_placeholder:
        return "placeholder"

    st = shape.shape_type
    if st == MSO_SHAPE_TYPE.CHART:
        return "chart"
    if st == MSO_SHAPE_TYPE.PICTURE:
        return "image"
    if st == MSO_SHAPE_TYPE.TABLE:
        return "table"
    if st == MSO_SHAPE_TYPE.GROUP:
        return "group"
    if st in (MSO_SHAPE_TYPE.LINE, MSO_SHAPE_TYPE.FREEFORM):
        return "line"
    if _is_text_shape(shape) and shape.text_frame.text.strip():
        return "text"
    return "shape"


def _shape_word_count(shape: BaseShape) -> int:
    """Count words in a shape's text frame."""
    if not _is_text_shape(shape):
        return 0
    text = shape.text_frame.text.strip()
    return len(text.split()) if text else 0


def _slide_word_count(slide: Slide) -> int:
    """Total words across all text-bearing shapes on a slide."""
    total = 0
    for shape in slide.shapes:
        total += _shape_word_count(shape)
    return total


def _show_file_summary(prs: PresentationObj, file_path: str) -> None:
    from .inspect_pptx import get_slide_comments, list_slides

    total = len(prs.slides)
    print(f"\n  {file_path}  ({total} slide{'s' if total != 1 else ''})\n")
    for entry in list_slides(prs):
        slide = prs.slides[entry["slide"] - 1]
        layout_name = slide.slide_layout.name
        shapes_n = len(slide.shapes)
        words = _slide_word_count(slide)
        title = entry["title"][:60] or "(no title)"
        comments = get_slide_comments(slide)
        cm_str = f"  {len(comments)}cm" if comments else ""
        print(f"  {entry['slide']:3d}  {title:62s}  [{layout_name}]  {shapes_n}sh  {words}w{cm_str}")
    print()


def _show_slide_header(slide: Slide, slide_label: str) -> None:
    title = "(no title)"
    subtitle = ""

    title_shape = slide.shapes.title
    if title_shape is not None:
        title = title_shape.text.replace("\x0b", " | ").replace("\n", " | ")

    for shape in slide.shapes:
        if _is_text_shape(shape) and shape.is_placeholder and shape.placeholder_format.idx == 1:
            subtitle = shape.text.replace("\x0b", " | ").replace("\n", " | ")
            break

    layout = slide.slide_layout
    layout_phs = [f"ph{ph.placeholder_format.idx}:{ph.name}" for ph in layout.placeholders]

    print(f"\n  Slide {slide_label}: {title}")
    if subtitle:
        print(f"  {subtitle}")
    print(f"  Layout: \"{layout.name}\" → {', '.join(layout_phs) if layout_phs else '(no placeholders)'}")
    print()


def _show_slide_shapes(slide: Slide) -> None:
    from .inspect_pptx import inspect_slide

    shapes = inspect_slide(slide)
    for item in shapes:
        real_shape = list(slide.shapes)[item.index]
        cat = _classify_shape_type(real_shape)

        ph_str = f"ph{item.placeholder_idx}" if item.placeholder_idx is not None else ""
        type_str = ph_str if ph_str else cat

        text = ""
        if item.text_preview:
            preview = item.text_preview[:60]
            if len(item.text_preview) > 60:
                preview += "…"
            text = f"  «{preview}»"

        print(
            f"  [{item.index:2d}] {type_str:8s}  {item.left:5.2f},{item.top:5.2f}  {item.width:5.2f}x{item.height:5.2f}  {item.name:30s}{text}"
        )


def _show_slide_comments(slide: Slide) -> None:
    from .inspect_pptx import get_slide_comments

    comments = get_slide_comments(slide)
    if not comments:
        return

    print(f"\n  Comments ({len(comments)}):")
    for comment in comments:
        preview = comment.text[:120]
        if len(comment.text) > 120:
            preview += "…"
        print(f"    [{comment.author}] {preview}")


def _show_shape_detail_json(slide: Slide, shape_identifier: str) -> None:
    from .inspect_pptx import inspect_chart, inspect_shape

    shape = _find_shape(slide, shape_identifier)
    if shape.has_chart if hasattr(shape, "has_chart") else False:
        _json_out(inspect_chart(shape))
    else:
        _json_out(inspect_shape(shape))


def cmd_show(args: _ShowArgs) -> int:
    """Progressive drill-down: file → slide list, slide → shapes, shape → detail."""
    prs = _open(args.file)

    if args.slide is None:
        _show_file_summary(prs, args.file)
        return 0

    slide = _get_slide(prs, args.slide)
    if args.shape is None:
        _show_slide_header(slide, args.slide)
        _show_slide_shapes(slide)
        _show_slide_comments(slide)
        print()
        return 0

    _show_shape_detail_json(slide, args.shape)
    return 0


def cmd_list(args: _FileArgs) -> int:
    from .inspect_pptx import list_slides

    prs = _open(args.file)
    for entry in list_slides(prs):
        print(f"  {entry['slide']:3d}  {entry['title']}")
    return 0


def cmd_summary(args: _SlideArgs) -> int:
    from .inspect_pptx import summarize_slide

    prs = _open(args.file)
    slide = _get_slide(prs, args.slide)
    _json_out(summarize_slide(slide))
    return 0


def cmd_slide(args: _SlideArgs) -> int:
    from .inspect_pptx import inspect_slide

    prs = _open(args.file)
    slide = _get_slide(prs, args.slide)
    shapes = inspect_slide(slide)
    for s in shapes:
        ph = f" ph={s.placeholder_idx}" if s.placeholder_idx is not None else ""
        fill = f" fill={s.fill.type}" if s.fill.type != "inherited" else ""
        text = f"  «{s.text_preview}»" if s.text_preview else ""
        print(
            f"  [{s.index:2d}] {s.name:30s}  {s.left:6.2f},{s.top:6.2f}  {s.width:5.2f}x{s.height:5.2f}{ph}{fill}{text}"
        )
    return 0


def cmd_shape(args: _XmlArgs) -> int:
    from .inspect_pptx import inspect_shape

    prs = _open(args.file)
    slide = _get_slide(prs, args.slide)
    shape = _find_shape(slide, args.shape)
    _json_out(inspect_shape(shape))
    return 0


def cmd_chart(args: _XmlArgs) -> int:
    from .inspect_pptx import inspect_chart

    prs = _open(args.file)
    slide = _get_slide(prs, args.slide)
    shape = _find_shape(slide, args.shape)
    _json_out(inspect_chart(shape))
    return 0


def cmd_layout(args: _SlideArgs) -> int:
    from .inspect_pptx import inspect_layout

    prs = _open(args.file)
    slide = _get_slide(prs, args.slide)
    _json_out(inspect_layout(slide.slide_layout))
    return 0


def cmd_layouts(args: _FileArgs) -> int:
    from pptx.enum.shapes import PP_PLACEHOLDER
    from pptx.shapes.placeholder import LayoutPlaceholder

    prs = _open(args.file)

    # Boundary between "structural" placeholders (title, subtitle, tracker) and
    # the slide's main content zone.
    content_y_threshold = int(Layout.CONTENT_START_Y)
    footer_y = int(Layout.FOOTER_LINE_Y)

    for layout in prs.slide_layouts:
        structural: list[str] = []
        # (left_emu, width_in, height_avail_in, name)
        content_phs: list[tuple[int, float, float, str]] = []

        placeholders = cast(Iterable[LayoutPlaceholder], layout.placeholders)
        phs: list[LayoutPlaceholder] = sorted(
            placeholders,
            key=lambda p: (int(p.top), int(p.left)),
        )
        for ph in phs:
            pf = ph.placeholder_format
            w_in = ph.width.inches

            if pf.type == PP_PLACEHOLDER.TITLE:
                structural.append(f"title ({w_in:.1f}in)")
            elif pf.type == PP_PLACEHOLDER.SUBTITLE:
                structural.append(f"subtitle ({w_in:.1f}in)")
            elif pf.type == PP_PLACEHOLDER.PICTURE:
                structural.append(f"image: {ph.name} ({w_in:.1f}×{ph.height.inches:.1f}in)")
            elif int(ph.top) < content_y_threshold:
                # Above content zone — tracker, doc type, etc.
                structural.append(f"{ph.name} ({w_in:.1f}in)")
            else:
                h_avail = Emu(footer_y - int(ph.top)).inches
                content_phs.append((int(ph.left), w_in, h_avail, ph.name))

        # Sort content areas left-to-right, label primary/secondary
        content_phs.sort(key=lambda p: p[0])
        areas: list[str] = []
        for idx, (_x, w, h, _name) in enumerate(content_phs):
            label = "primary" if idx == 0 else "secondary"
            areas.append(f"{label} {w:.1f}×{h:.1f}in")

        print(f"  {layout.name}")
        if structural:
            print(f"    placeholders: {', '.join(structural)}")
        if areas:
            print(f"    content areas: {', '.join(areas)}")
        else:
            print("    content areas: (none)")
    return 0


def cmd_theme(args: _FileArgs) -> int:
    from .inspect_pptx import resolve_theme_colors

    prs = _open(args.file)
    colors = resolve_theme_colors(prs)
    for name, hex_val in sorted(colors.items()):
        if not name.startswith("_"):
            print(f"  {name:12s}  {hex_val}")
    return 0


def cmd_color(args: _ColorArgs) -> int:
    from .inspect_pptx import identify_color

    prs = _open(args.file)
    result = identify_color(prs, args.hex)
    if result:
        print(result)
    else:
        print(f"No theme match for {args.hex}")
    return 0


def cmd_xml(args: _XmlArgs) -> int:
    from .xml_helpers import dump_xml

    prs = _open(args.file)
    slide = _get_slide(prs, args.slide)
    shape = _find_shape(slide, args.shape)
    print(dump_xml(shape._element))
    return 0


# ============================================================================
# EDIT COMMANDS
# ============================================================================


def cmd_edit(args: _EditArgs) -> int:
    prs = _open(args.file)
    slide = _get_slide(prs, args.slide)
    shape_raw = _find_shape(slide, args.shape)

    if not _is_text_shape(shape_raw):
        print(f"Error: shape '{args.shape}' has no text frame", file=sys.stderr)
        return 1

    shape = shape_raw

    before = _shape_text(shape)
    out_path = args.out or args.file

    print(f"Before: {before}")
    print(f"After:  {_text_preview(args.text)}")

    _write_to_shape(shape, args.text)
    _warn_placeholder_text_limits(slide, shape)

    prs.save(out_path)
    print(f"Saved → {out_path}")
    return 0


def cmd_batch(args: _BatchArgs) -> int:
    prs = _open(args.file)
    out_path = args.out or args.file

    if args.edits == "-":
        edits = json.load(sys.stdin)
    else:
        with open(args.edits) as f:
            edits = json.load(f)

    for idx, edit in enumerate(edits):
        slide_num = edit["slide"]
        shape_id = str(edit["shape"])
        text_arg = edit["text"]

        slide = _get_slide(prs, str(slide_num))
        shape_raw = _find_shape(slide, shape_id)

        if not _is_text_shape(shape_raw):
            print(f"  [{idx+1}] SKIP {shape_id} — no text frame", file=sys.stderr)
            continue

        shape = shape_raw

        before = _shape_text(shape)
        _write_to_shape(shape, text_arg)

        print(f"  [{idx+1}] slide {slide_num} / {shape.name}")
        print(f"      Before: {before}")
        print(f"      After:  {_text_preview(text_arg)}")

    prs.save(out_path)
    print(f"Saved → {out_path}  ({len(edits)} edits)")
    return 0


# ============================================================================
# SLIDE MANAGEMENT COMMANDS
# ============================================================================


def _parse_slide_selection(selection: str | None, total: int) -> list[int]:
    """Parse a slide selection string (e.g. "1,3-5") into 1-indexed slide numbers."""
    if not selection:
        return list(range(1, total + 1))

    result: list[int] = []
    seen: set[int] = set()

    for part in selection.split(","):
        part = part.strip()
        if not part:
            continue
        if "-" in part:
            a, b = part.split("-", 1)
            start = int(a)
            end = int(b)
            if start > end:
                start, end = end, start
            nums = range(start, end + 1)
        else:
            nums = [int(part)]

        for n in nums:
            if n < 1 or n > total:
                raise ValueError(f"slide {n} out of range (1-{total})")
            if n not in seen:
                seen.add(n)
                result.append(n)

    return result


def _move_last_slide_to(prs: PresentationObj, at_pos: int) -> None:
    """Move the last slide in *prs* to position *at_pos* (1-indexed)."""
    sldIdLst = prs.slides.element
    sld_ids = sldIdLst.sldId_lst
    if not sld_ids:
        raise ValueError("presentation has no slide ID list")

    last = sld_ids[-1]
    sldIdLst.remove(last)

    insert_idx = at_pos - 1
    if insert_idx < len(sldIdLst.sldId_lst):
        sldIdLst.insert(insert_idx, last)
    else:
        sldIdLst.append(last)


_NS_R = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
_HYPERLINK_RELTYPE = "http://schemas.openxmlformats.org/officeDocument/2006/relationships/hyperlink"


def _collect_relationship_ids(slide: Slide) -> set[str]:
    """Return all r:id values referenced by shapes on this slide."""
    ids: set[str] = set()
    spTree = slide.shapes.element
    for el in spTree.iter():
        for key, val in el.attrib.items():
            if str(key).startswith(f"{{{_NS_R}}}"):
                ids.add(str(val))
    return ids


def _check_and_copy_relationships(src: Slide, dst: Slide) -> None:
    """Copy external hyperlink rels from src to dst, reject embedded content.

    Hyperlinks are external relationships (just a URL string) and can be
    safely re-created on the destination slide. Embedded content (images,
    charts, media) requires OPC-level blob copying which is not yet supported.
    """
    ref_ids = _collect_relationship_ids(src)
    if not ref_ids:
        return

    # Build a rId remapping: src rId → dst rId
    remap: dict[str, str] = {}

    for rid in ref_ids:
        try:
            rel = src.part.rels[rid]
        except KeyError:
            continue

        if rel.is_external and rel.reltype == _HYPERLINK_RELTYPE:
            # Add the hyperlink as a new external relationship on dst
            new_rid = dst.part.rels.get_or_add_ext_rel(_HYPERLINK_RELTYPE, rel.target_ref)
            remap[rid] = new_rid
        else:
            raise ValueError(
                f"source slide contains embedded relationship rId={rid} "
                f"(type={rel.reltype}); pptx insert supports text and hyperlinks only"
            )

    # Remap r:id attributes in the destination shapes
    if remap:
        dst_spTree = dst.shapes.element
        for el in dst_spTree.iter():
            for key in list(el.attrib):
                if str(key).startswith(f"{{{_NS_R}}}"):
                    old_id = str(el.attrib[key])
                    if old_id in remap:
                        el.attrib[key] = remap[old_id]


def _replace_slide_shapes(dst: Slide, src: Slide) -> None:
    """Replace dst slide shapes with a deep-copy of src slide shapes."""
    dst_spTree = dst.shapes.element
    src_spTree = src.shapes.element

    # Keep spTree group headers (nvGrpSpPr + grpSpPr). Remove everything else.
    for child in list(dst_spTree)[2:]:
        dst_spTree.remove(child)

    for child in list(src_spTree)[2:]:
        dst_spTree.append(deepcopy(child))


def cmd_add_slide(args: _AddSlideArgs) -> int:
    prs = _open(args.file)
    layout = _find_layout(prs, args.layout)
    total_before = len(prs.slides)
    out_path = args.out or args.file

    slide = prs.slides.add_slide(layout)

    if args.at is not None:
        at_pos = args.at
        if at_pos < 1 or at_pos > total_before + 1:
            print(f"Error: --at {at_pos} out of range (1-{total_before + 1})", file=sys.stderr)
            return 1

        # Reorder slide by moving its <p:sldId> entry in the presentation XML.
        sldIdLst = prs.slides.element
        sld_ids = sldIdLst.sldId_lst
        if not sld_ids:
            print("Error: presentation has no slide ID list", file=sys.stderr)
            return 1

        new_sldId = sld_ids[-1]
        sldIdLst.remove(new_sldId)

        insert_idx = at_pos - 1
        if insert_idx < len(sldIdLst.sldId_lst):
            sldIdLst.insert(insert_idx, new_sldId)
        else:
            sldIdLst.append(new_sldId)

        final_pos = at_pos
    else:
        final_pos = total_before + 1

    phs: list[str] = [f"{ph.placeholder_format.idx}:{ph.name}" for ph in slide.placeholders]
    print(f'Added slide {final_pos} from layout "{layout.name}"')
    print(f"  Placeholders: {', '.join(phs) if phs else '(none)'}")
    print(f"  Total slides: {total_before + 1}")

    prs.save(out_path)
    print(f"Saved → {out_path}")
    return 0


def cmd_delete_slide(args: _DeleteSlideArgs) -> int:
    if not args.confirm:
        print("Error: delete-slide requires --confirm flag", file=sys.stderr)
        return 1

    prs = _open(args.file)
    total = len(prs.slides)
    out_path = args.out or args.file

    if total <= 1:
        print("Error: cannot delete the only slide in the presentation", file=sys.stderr)
        return 1

    slide = _get_slide(prs, args.slide)

    title = "(no title)"
    title_shape = slide.shapes.title
    if title_shape is not None:
        title = title_shape.text.replace("\x0b", " | ").replace("\n", " | ")[:60]

    print(f'Deleting slide {args.slide}: "{title}" ({len(slide.shapes)} shapes)')

    # Find relationship id (rId) for this slide part
    slide_part = slide.part
    rId: str | None = None
    for rel_key in prs.part.rels:
        rel = prs.part.rels[rel_key]
        if rel.target_part is slide_part:
            rId = rel_key
            break

    if rId is None:
        print("Error: could not find slide relationship", file=sys.stderr)
        return 1

    # Remove the <p:sldId> element referencing this slide
    sldIdLst = prs.slides.element
    for sldId in list(sldIdLst.sldId_lst):
        if sldId.rId == rId:
            sldIdLst.remove(sldId)
            break

    prs.part.rels.pop(rId)

    print(f"  Remaining slides: {total - 1}")
    prs.save(out_path)
    print(f"Saved → {out_path}")
    return 0


def cmd_delete_shape(args: _DeleteShapeArgs) -> int:
    prs = _open(args.file)
    slide = _get_slide(prs, args.slide)
    shape = _find_shape(slide, args.shape)
    out_path = args.out or args.file

    text_preview = ""
    if _is_text_shape(shape):
        text_preview = shape.text_frame.text.replace("\x0b", " ").replace("\n", " ")[:50]

    print(f'Deleting shape: "{shape.name}"')
    if text_preview:
        print(f"  Text: {text_preview}")

    slide.shapes.element.remove(shape.element)

    prs.save(out_path)
    print(f"Saved → {out_path}")
    return 0


def cmd_insert(args: _InsertArgs) -> int:
    """Insert slides from another PPTX into this presentation.

    Intended workflow:
        pptx generate spec.yaml -t template.pptx -o /tmp/table.pptx
        pptx insert  deck.pptx /tmp/table.pptx --at 5

    Notes:
    - Supports text and hyperlink slides. Embedded content (images/charts/media)
      is not yet supported.
    - Matches the source slide's layout name in the destination deck, falling
      back to Default (with a warning) when no match is found.  Native
      PowerPoint would import the source layout/master into the destination
      instead; we don't do that because it requires deep OPC-level copying
      of theme, fonts, and background assets.
    """

    prs = _open(args.file)
    src_prs = _open(args.source)
    out_path = args.out or args.file

    total_before = len(prs.slides)
    src_total = len(src_prs.slides)

    try:
        selected = _parse_slide_selection(args.slides, src_total)
    except ValueError as e:
        print(f"Error: {e}", file=sys.stderr)
        return 1

    if args.at is None:
        at_pos = total_before + 1
    else:
        at_pos = args.at
        if at_pos < 1 or at_pos > total_before + 1:
            print(f"Error: --at {at_pos} out of range (1-{total_before + 1})", file=sys.stderr)
            return 1

    insert_pos = at_pos
    for slide_num in selected:
        src_slide = _get_slide(src_prs, slide_num)

        # Match the source slide's layout by name in the destination deck.
        # Fall back to "Default" if no match, then best-effort.
        #
        # NOTE: Native PowerPoint imports the source layout (and its master)
        # into the destination when no matching layout exists, creating
        # numbered duplicates like "1_LayoutName".  We don't do that —
        # importing a full layout/master with its theme, fonts, and
        # backgrounds is a deep OPC-level operation.  Instead we fall back
        # to "Default" and warn.  Shapes are still copied faithfully.
        dst_layout: SlideLayout | None = None
        src_layout_name = getattr(src_slide.slide_layout, "name", "")
        if src_layout_name:
            dst_layout = _try_find_layout(prs, src_layout_name)
        if dst_layout is None and src_layout_name:
            print(
                f"  WARNING: slide {slide_num} layout '{src_layout_name}' not found "
                f"in destination; falling back to Default. "
                f"(Native PowerPoint would import the source layout instead.)"
            )
        if dst_layout is None:
            dst_layout = _try_find_layout(prs, "default")
        if dst_layout is None:
            dst_layout = _find_layout(prs, "default", fallback=True)
            print(
                f"  WARNING: slide {slide_num} — 'Default' layout not found either; "
                f"using '{dst_layout.name}'."
            )

        dst_slide = prs.slides.add_slide(dst_layout)
        _replace_slide_shapes(dst_slide, src_slide)

        # Copy hyperlink relationships from source to destination (remapping rIds).
        # Raises ValueError for embedded content (images/charts) which isn't supported yet.
        try:
            _check_and_copy_relationships(src_slide, dst_slide)
        except ValueError as e:
            print(f"Error: slide {slide_num}: {e}", file=sys.stderr)
            return 1

        # Reorder to requested insertion position.
        _move_last_slide_to(prs, insert_pos)
        insert_pos += 1

    prs.save(out_path)
    print(
        f"Inserted {len(selected)} slide(s) from {args.source} into {args.file} at position {at_pos} → {out_path}"
    )
    return 0


# ============================================================================
# RENDER COMMANDS
# ============================================================================


def cmd_render(args: _RenderArgs) -> int:
    from .screenshot import render_slides

    prs = _open(args.file)
    total = len(prs.slides)
    try:
        slide_nums = _parse_slide_selection(args.slides, total)
    except ValueError as e:
        print(f"Error: {e}", file=sys.stderr)
        return 1

    paths = render_slides(
        args.file,
        slide_nums,
        dpi=args.dpi,
        output_dir=args.out,
        engine=args.engine,
    )
    for p in paths:
        print(p)
    return 0


def cmd_crop(args: _CropArgs) -> int:
    result = crop_region(
        args.png,
        args.left,
        args.top,
        args.right,
        args.bottom,
        output_path=args.out,
    )
    print(result)
    return 0


def cmd_charts(args: _ChartsArgs) -> int:
    from .charts import generate_charts_from_json

    input_path = Path(args.input)
    output_path = Path(args.output)
    template_path = Path(args.template) if args.template else None

    try:
        generate_charts_from_json(
            input_path,
            output_path,
            template=template_path,
            layout=args.layout,
            expected_template=args.expected_template,
        )
    except (FileNotFoundError, ImportError) as exc:
        print(f"Error: {exc}", file=sys.stderr)
        return 1

    print(f"Saved → {output_path}")
    return 0


# ============================================================================
# GENERATE COMMANDS (YAML → PPTX table pipeline)
# ============================================================================



# YAML parsing/validation helpers now live in clean_slides.spec_pipeline

def _expand_inputs(inputs: list[str]) -> list[Path]:
    files: list[Path] = []
    for pattern in inputs:
        path = Path(pattern)
        if path.is_file():
            files.append(path)
        else:
            files.extend(Path(".").glob(pattern))
    return files


def _clear_content_area(slide: Slide, area: ContentArea) -> None:
    for shape in list(slide.shapes):
        if _boxes_intersect(
            (int(shape.left), int(shape.top), int(shape.width), int(shape.height)),
            (area.x, area.y, area.width, area.height),
        ):
            slide.shapes.element.remove(shape.element)


_Box = tuple[int, int, int, int]


def _boxes_intersect(box_a: _Box, box_b: _Box) -> bool:
    ax, ay, aw, ah = box_a
    bx, by, bw, bh = box_b
    return ax < bx + bw and ax + aw > bx and ay < by + bh and ay + ah > by


def _infer_layout_from_slide(slide: Slide) -> str:
    for shape in slide.shapes:
        if shape.top < Layout.CONTENT_START_Y:
            return "content"
    return "default"


def _delete_all_slides(prs: PresentationObj) -> None:
    """Remove all existing slides from a presentation."""
    # Must iterate in reverse to avoid index shifting
    # Access internal _sldIdLst to properly delete slides
    slides = prs.slides
    for i in range(len(slides) - 1, -1, -1):
        rId: str = slides._sldIdLst[i].rId  # type: ignore[union-attr]
        prs.part.drop_rel(rId)  # type: ignore[union-attr]
        del slides._sldIdLst[i]  # type: ignore[union-attr]


def _clear_body_placeholders(slide: Slide) -> None:
    """Remove unfilled body/content placeholders from a slide."""
    # Only keep these specific placeholder types (by idx):
    # 0 = title, 1 = subtitle
    # Tracker placeholder is identified by name, not idx (idx varies by layout)
    KEEP_PLACEHOLDER_INDICES = {0, 1}

    shapes_to_remove: list[BaseShape] = []
    for shape in slide.shapes:
        if not shape.is_placeholder:
            continue
        ph_idx = shape.placeholder_format.idx

        # Always keep title and subtitle
        if ph_idx in KEEP_PLACEHOLDER_INDICES:
            continue

        # Keep tracker placeholder (identified by name pattern)
        name_lower = shape.name.lower()
        if "tracker" in name_lower or "on-page" in name_lower:
            continue

        # Check if placeholder has meaningful content
        if hasattr(shape, "text_frame"):
            text: str = str(shape.text_frame.text).strip()  # type: ignore[union-attr]
            # Keep if it has real content (not empty, not placeholder prompts)
            if text and "click to" not in text.lower() and "master text" not in text.lower():
                continue

        # Remove this placeholder
        shapes_to_remove.append(shape)

    # Remove shapes by deleting their XML elements
    for shape in shapes_to_remove:
        sp = shape.element
        parent = sp.getparent()
        if parent is not None:
            parent.remove(sp)


def _hint_init() -> None:
    """Print a one-time hint about `pptx init` when no project config exists."""
    proj = _discover_project_dir()
    if proj is None:
        print(
            "Hint: run `pptx init` to set up a project template and config "
            f"in {PROJECT_DIR_NAME}/",
            file=sys.stderr,
        )


def _resolve_generate_template_path(template_path: str | None) -> str | None:
    if template_path is not None:
        return template_path

    discovered_tpl = _discover_template()
    if discovered_tpl is not None:
        return str(discovered_tpl)

    return None


def _init_generate_presentation(template_path: str | None) -> PresentationObj:
    if template_path:
        prs = Presentation(template_path)
        # Remove template's content slides - keep only layouts/masters
        _delete_all_slides(prs)
        return prs

    _hint_init()
    prs = Presentation()
    prs.slide_width = Emu(int(Layout.SLIDE_WIDTH))
    prs.slide_height = Emu(int(Layout.SLIDE_HEIGHT))
    return prs


def _load_validated_generate_data(path: Path) -> tuple[YamlDict, bool] | None:
    data = load_yaml(str(path))
    errors, warnings = validate_spec(data)
    if errors:
        print(f"{path}:\n  - " + "\n  - ".join(errors), file=sys.stderr)
        return None
    if warnings:
        print(f"{path}:\n  - " + "\n  - ".join(warnings))

    has_table = _is_str_object_dict(data.get("table")) and bool(data.get("table"))
    return data, has_table


def _resolve_generate_slide(
    prs: PresentationObj,
    data: YamlDict,
    slide_index: int | None,
) -> tuple[Slide, SlideLayout | None, str | None] | None:
    layout_override: str | None = None

    if slide_index is not None:
        if slide_index < 0 or slide_index >= len(prs.slides):
            print("slide_index out of range", file=sys.stderr)
            return None

        slide = prs.slides[slide_index]
        if "content_layout" not in data and "layout" not in data:
            layout_override = _infer_layout_from_slide(slide)

        return slide, None, layout_override

    slide_layout_name = str(data.get("slide_layout") or "Default")
    slide_layout_obj = _find_layout(prs, slide_layout_name, fallback=True)
    slide = prs.slides.add_slide(slide_layout_obj)
    return slide, slide_layout_obj, None


def _render_chart_cells_for_spec(
    slide: Slide,
    spec: TableSpec,
    layout: TableLayout,
    area: ContentArea,
) -> None:
    if not spec.chart_defs:
        return

    from .chart_render import render_chart_cells
    from .charts import load_charts_module

    charts_mod = load_charts_module()
    body_pt = layout.body_font_size // 100
    render_chart_cells(
        slide,
        spec,
        layout,
        area,
        charts_mod,
        label_font_size_pt=body_pt,
    )


def _render_sidebar_content(
    data: YamlDict,
    slide_layout_obj: SlideLayout | None,
    renderer: TableRenderer,
    metrics: TextMetrics,
) -> None:
    sidebar_raw = data.get("sidebar")
    if sidebar_raw is None or slide_layout_obj is None:
        return

    sidebar_area = _sidebar_content_area(slide_layout_obj)
    if sidebar_area is None:
        print("  WARNING: sidebar content specified but layout has no secondary content area")
        return

    from .content import Paragraph, normalize_cell

    default_para = Paragraph(text="", lvl=0)
    sidebar_paras = normalize_cell(sidebar_raw, default_para, parse_bullets=True)
    if not sidebar_paras:
        return

    if data.get("sidebar_shrink"):
        _shrink_sidebar_to_fit(sidebar_paras, sidebar_area, metrics)
    else:
        _warn_sidebar_overflow(sidebar_paras, sidebar_area, metrics)

    renderer.render_sidebar(sidebar_paras, sidebar_area)


def _render_table_for_input(
    path: Path,
    slide: Slide,
    slide_layout_obj: SlideLayout | None,
    data: YamlDict,
    layout_override: str | None,
    solver: ConstraintSolver,
    metrics: TextMetrics,
    *,
    detail: bool,
    target_slide_index: int | None,
    keep_existing: bool,
) -> None:
    spec, area, options, placeholders = parse_spec(data, layout_override=layout_override)
    if placeholders:
        spec = fill_placeholders(spec)

    # Derive content area from the layout's primary content placeholder
    # unless the YAML explicitly overrides via content_area / content_layout.
    if slide_layout_obj is not None and "content_area" not in data:
        layout_area = _content_area_from_layout(slide_layout_obj)
        if layout_area is not None:
            area = layout_area

    if target_slide_index is not None and not keep_existing:
        _clear_content_area(slide, area)

    layout, report = solver.solve(spec, area, options)
    print(f"{path}: {report.to_text(detail=detail)}")

    sp_tree = slide.shapes.element

    shape_id: int = TableDefaults.SHAPE_ID_START

    def next_shape_id() -> int:
        nonlocal shape_id
        shape_id += 1
        return shape_id

    renderer = TableRenderer(sp_tree, next_shape_id, slide_part=slide.part)
    renderer.render(spec, layout, area)

    _render_chart_cells_for_spec(slide, spec, layout, area)
    _render_sidebar_content(data, slide_layout_obj, renderer, metrics)


def _finalize_generated_slide(slide: Slide, data: YamlDict) -> None:
    fill_slide_metadata(slide, data)

    # Agent-friendly: warn when title/subtitle likely wrap beyond configured limits.
    for shape in slide.shapes:
        if isinstance(shape, Shape) and shape.has_text_frame:
            _warn_placeholder_text_limits(slide, shape)

    _clear_body_placeholders(slide)


def _process_generate_input(
    path: Path,
    prs: PresentationObj,
    args: _GenerateArgs,
    solver: ConstraintSolver,
    metrics: TextMetrics,
) -> bool:
    loaded = _load_validated_generate_data(path)
    if loaded is None:
        return False

    data, has_table = loaded

    resolved = _resolve_generate_slide(prs, data, args.slide_index)
    if resolved is None:
        return False

    slide, slide_layout_obj, layout_override = resolved

    if has_table:
        _render_table_for_input(
            path,
            slide,
            slide_layout_obj,
            data,
            layout_override,
            solver,
            metrics,
            detail=args.detail,
            target_slide_index=args.slide_index,
            keep_existing=args.keep_existing,
        )
    else:
        print(f"{path}: metadata-only slide (no table)")

    _finalize_generated_slide(slide, data)
    return True


def cmd_generate(args: _GenerateArgs) -> int:
    """Generate PPTX for text-only tables."""
    _apply_config(args.config)

    input_files = _expand_inputs(args.input)
    if not input_files:
        print("No input files found", file=sys.stderr)
        return 1

    template_path = _resolve_generate_template_path(args.template)
    prs = _init_generate_presentation(template_path)

    metrics = TextMetrics()
    solver = ConstraintSolver(metrics)

    for path in input_files:
        ok = _process_generate_input(path, prs, args, solver, metrics)
        if not ok:
            return 1

    output_path = args.output or "output.pptx"
    prs.save(output_path)
    print(f"Saved: {output_path}")

    return 0


def _example_template_dir() -> Path:
    """Return path to the bundled example-template directory."""
    return Path(__file__).resolve().parent / "example-template"


def cmd_init(args: _InitArgs) -> int:
    """Initialise a .clean-slides/ project directory.

    Without --template: copies the bundled example template and config.
    With --template <file.pptx>: runs init-config to generate a config,
    then copies the template + generated config into .clean-slides/.
    """
    import shutil

    target = Path(args.output) if args.output else Path.cwd()
    project_dir = target / PROJECT_DIR_NAME

    if project_dir.exists():
        print(f"Already initialised: {project_dir}", file=sys.stderr)
        return 1

    project_dir.mkdir(parents=True)

    if args.template:
        # Copy the user's template
        src_tpl = Path(args.template)
        if not src_tpl.is_file():
            print(f"Template not found: {src_tpl}", file=sys.stderr)
            return 1
        dst_tpl = project_dir / _TEMPLATE_NAME
        shutil.copy2(src_tpl, dst_tpl)

        # Generate config via init-config
        dst_cfg = project_dir / _CONFIG_NAME

        class _FakeArgs:
            file = str(src_tpl)
            output = str(dst_cfg)

        rc = cmd_init_config(_FakeArgs())  # type: ignore[arg-type]
        if rc != 0:
            return rc
    else:
        # Copy bundled example
        example_dir = _example_template_dir()
        src_tpl = example_dir / "example-template.pptx"
        src_cfg = example_dir / "example-config.yaml"
        if not src_tpl.is_file():
            print(f"Bundled example not found: {src_tpl}", file=sys.stderr)
            return 1
        shutil.copy2(src_tpl, project_dir / _TEMPLATE_NAME)
        shutil.copy2(src_cfg, project_dir / _CONFIG_NAME)

    print(f"Initialised {project_dir}/")
    print(f"  {_TEMPLATE_NAME}  — slide template")
    print(f"  {_CONFIG_NAME}    — colours, fonts, layout config")
    print()
    print("Generate slides:  pptx generate spec.yaml -o output.pptx")
    print(f"Edit config:      $EDITOR {project_dir / _CONFIG_NAME}")
    return 0


def cmd_init_config(args: _InitConfigArgs) -> int:
    """Generate a starter template-config.yaml by introspecting a PPTX template."""
    from .template_init_config import build_init_config_output

    output = build_init_config_output(args.file)

    if args.output:
        Path(args.output).write_text(output)
        print(f"Saved → {args.output}")
        print("Review the config, especially:")
        print("  - colors: map your template's theme colors to semantic names")
        print("  - bullets: inspect slide master lstStyle for accurate margins/chars")
        print("  - placeholders: verify indices match your template")
        print("  - font_sizes: adjust to match your template's type scale")
    else:
        print(output)

    return 0


def cmd_validate(args: _ValidateArgs) -> int:
    """Validate schema for YAML files."""
    _apply_config(args.config)

    input_files = _expand_inputs(args.input)
    if not input_files:
        print("No input files found", file=sys.stderr)
        return 1

    all_valid = True
    for path in input_files:
        data = load_yaml(str(path))
        errors, warnings = validate_spec(data)
        if errors:
            all_valid = False
            print(f"{path}:\n  - " + "\n  - ".join(errors))
        elif warnings:
            print(f"{path}:\n  - " + "\n  - ".join(warnings))
        else:
            print(f"{path}: OK")
    return 0 if all_valid else 1


def cmd_preview(args: _InputArgs) -> int:
    """Preview table structure without generating PPTX."""
    input_files = _expand_inputs(args.input)
    if not input_files:
        print("No input files found", file=sys.stderr)
        return 1

    for path in input_files:
        data = load_yaml(str(path))
        print(preview_spec(data))
    return 0


def cmd_verify(args: _VerifyArgs) -> int:
    """Run layout solver + report without generating PPTX."""
    _apply_config(args.config)

    input_files = _expand_inputs(args.input)
    if not input_files:
        print("No input files found", file=sys.stderr)
        return 1

    metrics = TextMetrics()
    solver = ConstraintSolver(metrics)

    for path in input_files:
        data = load_yaml(str(path))
        errors, warnings = validate_spec(data)
        if errors:
            print(f"{path}:\n  - " + "\n  - ".join(errors), file=sys.stderr)
            return 1
        if warnings:
            print(f"{path}:\n  - " + "\n  - ".join(warnings))

        spec, area, options, placeholders = parse_spec(data)
        if placeholders:
            spec = fill_placeholders(spec)

        _, report = solver.solve(spec, area, options)
        print(f"{path}: {report.to_text(detail=args.detail)}")

        if args.json:
            output_path = Path(args.json)
            output_path.write_text(json.dumps(report.to_json(), indent=2))

    return 0


def cmd_screenshot(args: _ScreenshotArgs) -> int:
    """Generate PNG screenshots from PPTX files."""
    input_files = _expand_inputs(args.input)
    if not input_files:
        print("No input files found", file=sys.stderr)
        return 1

    output_dir = Path(args.output_dir or "outputs/images")
    generator = ScreenshotGenerator(soffice_path=args.soffice)

    for path in input_files:
        png_path = generator.capture(Path(path), output_dir, slide_index=args.slide)
        print(f"{path}: {png_path}")

    return 0


# ============================================================================
# ARGPARSE SETUP
# ============================================================================


def _build_parser() -> argparse.ArgumentParser:
    from .cli_parser import build_parser

    return build_parser(
        cmd_show=cmd_show,
        cmd_list=cmd_list,
        cmd_summary=cmd_summary,
        cmd_slide=cmd_slide,
        cmd_shape=cmd_shape,
        cmd_chart=cmd_chart,
        cmd_layout=cmd_layout,
        cmd_layouts=cmd_layouts,
        cmd_theme=cmd_theme,
        cmd_color=cmd_color,
        cmd_xml=cmd_xml,
        cmd_edit=cmd_edit,
        cmd_batch=cmd_batch,
        cmd_add_slide=cmd_add_slide,
        cmd_delete_slide=cmd_delete_slide,
        cmd_delete_shape=cmd_delete_shape,
        cmd_insert=cmd_insert,
        cmd_render=cmd_render,
        cmd_crop=cmd_crop,
        cmd_charts=cmd_charts,
        cmd_generate=cmd_generate,
        cmd_validate=cmd_validate,
        cmd_preview=cmd_preview,
        cmd_verify=cmd_verify,
        cmd_init=cmd_init,
        cmd_init_config=cmd_init_config,
        cmd_screenshot=cmd_screenshot,
    )


def main() -> int:
    parser = _build_parser()
    args = parser.parse_args()

    if not args.command:
        parser.print_help()
        return 1

    return args.func(args)


if __name__ == "__main__":
    sys.exit(main())
