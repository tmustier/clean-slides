"""
Slide metadata — fill title, subtitle, tracker placeholders.

Handles multi-line text via ``<a:br>`` (line break within single
paragraph) to avoid extra paragraph spacing.
"""

from __future__ import annotations

import contextlib
from typing import Any

from lxml import etree
from pptx.shapes.autoshape import Shape
from pptx.slide import Slide

from .template_config import TEMPLATE_CONFIG, TemplateConfig

NS_A = "http://schemas.openxmlformats.org/drawingml/2006/main"

# Default placeholder indices (overridable via template-config.yaml).
_DEFAULT_PLACEHOLDERS: dict[str, int] = {
    "title": 0,
    "subtitle": 1,
    "tracker": 17,
}

# Active placeholder map — refreshed by _reload().
PLACEHOLDER_MAP: dict[str, int] = dict(_DEFAULT_PLACEHOLDERS)


def reload_placeholders(config: TemplateConfig | None = None) -> None:
    """Refresh :data:`PLACEHOLDER_MAP` from *config*.

    When a ``placeholders`` section is present in the config, **only** its
    keys are used — no defaults are injected.  This lets a config omit
    ``subtitle`` or ``tracker`` when the template lacks those placeholders.
    When no ``placeholders`` section exists at all, the built-in defaults
    are used for backward compatibility.
    """
    if config is None:
        config = TEMPLATE_CONFIG

    PLACEHOLDER_MAP.clear()

    raw = config.section("placeholders") if config.has_section("placeholders") else None
    if raw is not None:
        # Explicit config — use only what's declared.
        for key, val in raw.items():
            with contextlib.suppress(TypeError, ValueError):
                PLACEHOLDER_MAP[str(key)] = int(val)
    else:
        # No placeholders section — use built-in defaults.
        PLACEHOLDER_MAP.update(_DEFAULT_PLACEHOLDERS)


# Initialize at import time.
reload_placeholders()


def fill_slide_metadata(slide: Slide, data: dict[str, Any]) -> None:
    """Write title / subtitle / tracker into existing placeholders.

    Newlines (``\\n``) within a value produce ``<a:br>`` elements inside
    a single paragraph so there is no extra paragraph spacing.
    """
    for key, idx in PLACEHOLDER_MAP.items():
        text = data.get(key)
        if not text:
            continue
        text = str(text)
        for ph in slide.placeholders:
            if not isinstance(ph, Shape):
                continue
            if ph.placeholder_format.idx != idx:
                continue

            if "\n" not in text:
                ph.text = text
                break

            _set_multiline(ph, text)
            break


def _set_multiline(ph: Shape, text: str) -> None:
    """Rewrite placeholder XML to use ``<a:br>`` within a single paragraph."""
    txBody = ph.element.find(f".//{{{NS_A}}}txBody")
    if txBody is None:
        ph.text = text
        return

    # Preserve first paragraph's pPr (inherits placeholder style)
    first_p = txBody.find(f"{{{NS_A}}}p")
    pPr_bytes: bytes | None = None
    if first_p is not None:
        pPr = first_p.find(f"{{{NS_A}}}pPr")
        if pPr is not None:
            pPr_bytes = etree.tostring(pPr)

    # Remove all existing paragraphs
    for old_p in txBody.findall(f"{{{NS_A}}}p"):
        txBody.remove(old_p)

    # Build single paragraph with <a:br> between lines
    new_p = etree.SubElement(txBody, f"{{{NS_A}}}p")
    if pPr_bytes is not None:
        new_p.insert(0, etree.fromstring(pPr_bytes))

    for li, line in enumerate(text.split("\n")):
        if li > 0:
            br = etree.SubElement(new_p, f"{{{NS_A}}}br")
            etree.SubElement(br, f"{{{NS_A}}}rPr", lang="en-US")
        r = etree.SubElement(new_p, f"{{{NS_A}}}r")
        etree.SubElement(r, f"{{{NS_A}}}rPr", lang="en-US")
        t = etree.SubElement(r, f"{{{NS_A}}}t")
        t.text = line
