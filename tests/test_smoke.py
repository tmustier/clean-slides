from __future__ import annotations

from dataclasses import dataclass
from pathlib import Path

from pptx import Presentation

from clean_slides.cli import cmd_generate, cmd_verify


@dataclass
class GenerateArgs:
    input: list[str]
    template: str | None
    slide_index: int | None
    keep_existing: bool
    detail: bool
    output: str | None
    config: str | None = None


@dataclass
class VerifyArgs:
    input: list[str]
    detail: bool
    json: str | None
    config: str | None = None


def _write_yaml(path: Path) -> None:
    path.write_text(
        """
layout: default

table:
  rows: 2
  cols: 2
  has_col_header: true
  col_headers: ["Header A", "Header B"]
  cells:
    - ["Alpha", "Beta"]
    - ["Gamma", "Delta"]
""".strip()
    )


def test_generate_smoke(tmp_path: Path) -> None:
    yaml_path = tmp_path / "table.yaml"
    _write_yaml(yaml_path)

    output_path = tmp_path / "output.pptx"
    args = GenerateArgs(
        input=[str(yaml_path)],
        template=None,
        slide_index=None,
        keep_existing=False,
        detail=False,
        output=str(output_path),
    )

    result = cmd_generate(args)

    assert result == 0
    assert output_path.exists()
    assert output_path.stat().st_size > 0

    prs = Presentation(str(output_path))
    assert len(prs.slides) == 1
    assert len(prs.slides[0].shapes) > 0


def test_generate_metadata_only(tmp_path: Path) -> None:
    """Generate a slide with no table section — metadata-only."""
    yaml_path = tmp_path / "section.yaml"
    yaml_path.write_text(
        """
title: "Section Divider"
subtitle: "No table here"
slide_layout: "Default"
""".strip()
    )

    output_path = tmp_path / "output.pptx"
    args = GenerateArgs(
        input=[str(yaml_path)],
        template=None,
        slide_index=None,
        keep_existing=False,
        detail=False,
        output=str(output_path),
    )

    result = cmd_generate(args)

    assert result == 0
    assert output_path.exists()

    prs = Presentation(str(output_path))
    assert len(prs.slides) == 1


def test_generate_mixed_table_and_metadata(tmp_path: Path) -> None:
    """Generate a deck with both table and metadata-only slides."""
    meta_path = tmp_path / "01-title.yaml"
    meta_path.write_text(
        """
title: "Deck Title"
subtitle: "A subtitle"
slide_layout: "Default"
""".strip()
    )

    table_path = tmp_path / "02-table.yaml"
    table_path.write_text(
        """
title: "Data Slide"
table:
  rows: 2
  cols: 2
  has_col_header: true
  col_headers: ["A", "B"]
  cells:
    - ["One", "Two"]
    - ["Three", "Four"]
""".strip()
    )

    output_path = tmp_path / "output.pptx"
    args = GenerateArgs(
        input=[str(meta_path), str(table_path)],
        template=None,
        slide_index=None,
        keep_existing=False,
        detail=False,
        output=str(output_path),
    )

    result = cmd_generate(args)

    assert result == 0

    prs = Presentation(str(output_path))
    assert len(prs.slides) == 2
    # First slide: metadata-only (no table shapes)
    # Second slide: has table content
    assert len(prs.slides[1].shapes) > len(prs.slides[0].shapes)


def test_verify_smoke(tmp_path: Path) -> None:
    yaml_path = tmp_path / "table.yaml"
    _write_yaml(yaml_path)

    args = VerifyArgs(
        input=[str(yaml_path)],
        detail=False,
        json=None,
    )

    result = cmd_verify(args)

    assert result == 0
