"""Text body template cache and application helpers."""

from __future__ import annotations

import copy
from collections.abc import Iterable
from pathlib import Path
from typing import Protocol, cast

from pptx import Presentation
from pptx.oxml.ns import qn

from ..pptx_access import shape_has_text_frame, shape_text, shape_xml_element


class _XmlElementLike(Protocol):
    text: str | None

    def find(self, path: str) -> object | None: ...

    def iter(self, tag: str) -> Iterable[object]: ...

    def append(self, element: object) -> None: ...

    def getparent(self) -> object | None: ...


class _XmlParentLike(Protocol):
    def replace(self, old: object, new: object) -> None: ...


TxBodyTemplate = object

_TEXT_STYLE_CACHE: dict[tuple[str, str], TxBodyTemplate] = {}


def load_txbody_template(template_path: Path, sample_text: str) -> TxBodyTemplate | None:
    key = (str(template_path), sample_text)
    cached = _TEXT_STYLE_CACHE.get(key)
    if cached is not None:
        return cached

    prs = Presentation(str(template_path))
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape_has_text_frame(shape):
                continue

            text_value = shape_text(shape)
            if text_value is None:
                continue
            if text_value.strip() != sample_text:
                continue

            shape_element = shape_xml_element(shape)
            if shape_element is None:
                continue

            shape_element_obj = cast(_XmlElementLike, shape_element)
            tx_body = shape_element_obj.find(qn("p:txBody"))
            if tx_body is None:
                continue

            cached = copy.deepcopy(tx_body)
            _TEXT_STYLE_CACHE[key] = cached
            return cached

    return None


def apply_txbody_template(box: object, template: TxBodyTemplate | None, text: str | None) -> None:
    if template is None:
        return

    tx_body = copy.deepcopy(template)
    tx_body_element = cast(_XmlElementLike, tx_body)

    if text is not None:
        for t_elem_obj in tx_body_element.iter(qn("a:t")):
            t_elem = cast(_XmlElementLike, t_elem_obj)
            t_elem.text = text

    box_element = shape_xml_element(box)
    if box_element is None:
        return

    box_element_obj = cast(_XmlElementLike, box_element)
    existing = box_element_obj.find(qn("p:txBody"))
    if existing is not None:
        existing_el = cast(_XmlElementLike, existing)
        parent = existing_el.getparent()
        if parent is not None:
            parent_el = cast(_XmlParentLike, parent)
            parent_el.replace(existing_el, tx_body)
    else:
        box_element_obj.append(tx_body)


def resolve_txbody_template(
    template_path: Path | None,
    text: str,
    fallback: TxBodyTemplate | None,
) -> TxBodyTemplate | None:
    if template_path is None:
        return fallback

    template = load_txbody_template(template_path, text)
    return template if template is not None else fallback
