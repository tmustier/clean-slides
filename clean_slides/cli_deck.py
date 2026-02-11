"""Deck/slide management CLI commands (add/delete/insert)."""

from __future__ import annotations

import sys
from copy import deepcopy
from typing import Protocol

from pptx.presentation import Presentation as PresentationObj
from pptx.slide import Slide, SlideLayout

from .cli_common import (
    find_layout,
    find_shape,
    get_slide,
    is_text_shape,
    open_presentation,
    try_find_layout,
)


class FileArgs(Protocol):
    file: str


class AddSlideArgs(FileArgs, Protocol):
    layout: str
    at: int | None
    out: str | None


class DeleteSlideArgs(FileArgs, Protocol):
    slide: str
    confirm: bool
    out: str | None


class DeleteShapeArgs(FileArgs, Protocol):
    slide: str
    shape: str
    out: str | None


class InsertArgs(FileArgs, Protocol):
    source: str
    at: int | None
    slides: str | None
    out: str | None


def parse_slide_selection(selection: str | None, total: int) -> list[int]:
    """Parse a slide selection string (e.g. "1,3-5") into 1-indexed slide numbers."""
    if not selection:
        return list(range(1, total + 1))

    result: list[int] = []
    seen: set[int] = set()

    for part in selection.split(","):
        part = part.strip()
        if not part:
            continue
        if "-" in part:
            a, b = part.split("-", 1)
            start = int(a)
            end = int(b)
            if start > end:
                start, end = end, start
            nums = range(start, end + 1)
        else:
            nums = [int(part)]

        for n in nums:
            if n < 1 or n > total:
                raise ValueError(f"slide {n} out of range (1-{total})")
            if n not in seen:
                seen.add(n)
                result.append(n)

    return result


def _move_last_slide_to(prs: PresentationObj, at_pos: int) -> None:
    """Move the last slide in *prs* to position *at_pos* (1-indexed)."""
    sld_id_list = prs.slides.element
    sld_ids = sld_id_list.sldId_lst
    if not sld_ids:
        raise ValueError("presentation has no slide ID list")

    last = sld_ids[-1]
    sld_id_list.remove(last)

    insert_idx = at_pos - 1
    if insert_idx < len(sld_id_list.sldId_lst):
        sld_id_list.insert(insert_idx, last)
    else:
        sld_id_list.append(last)


_NS_R = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
_HYPERLINK_RELTYPE = "http://schemas.openxmlformats.org/officeDocument/2006/relationships/hyperlink"


def _collect_relationship_ids(slide: Slide) -> set[str]:
    """Return all r:id values referenced by shapes on this slide."""
    ids: set[str] = set()
    sp_tree = slide.shapes.element
    for el in sp_tree.iter():
        for key, val in el.attrib.items():
            if str(key).startswith(f"{{{_NS_R}}}"):
                ids.add(str(val))
    return ids


def _check_and_copy_relationships(src: Slide, dst: Slide) -> None:
    """Copy external hyperlink rels from src to dst, reject embedded content.

    Hyperlinks are external relationships (just a URL string) and can be
    safely re-created on the destination slide. Embedded content (images,
    charts, media) requires OPC-level blob copying which is not yet supported.
    """
    ref_ids = _collect_relationship_ids(src)
    if not ref_ids:
        return

    # Build a rId remapping: src rId → dst rId
    remap: dict[str, str] = {}

    for rid in ref_ids:
        try:
            rel = src.part.rels[rid]
        except KeyError:
            continue

        if rel.is_external and rel.reltype == _HYPERLINK_RELTYPE:
            # Add the hyperlink as a new external relationship on dst
            new_rid = dst.part.rels.get_or_add_ext_rel(_HYPERLINK_RELTYPE, rel.target_ref)
            remap[rid] = new_rid
        else:
            raise ValueError(
                f"source slide contains embedded relationship rId={rid} "
                f"(type={rel.reltype}); pptx insert supports text and hyperlinks only"
            )

    # Remap r:id attributes in the destination shapes
    if remap:
        dst_sp_tree = dst.shapes.element
        for el in dst_sp_tree.iter():
            for key in list(el.attrib):
                if str(key).startswith(f"{{{_NS_R}}}"):
                    old_id = str(el.attrib[key])
                    if old_id in remap:
                        el.attrib[key] = remap[old_id]


def _replace_slide_shapes(dst: Slide, src: Slide) -> None:
    """Replace dst slide shapes with a deep-copy of src slide shapes."""
    dst_sp_tree = dst.shapes.element
    src_sp_tree = src.shapes.element

    # Keep spTree group headers (nvGrpSpPr + grpSpPr). Remove everything else.
    for child in list(dst_sp_tree)[2:]:
        dst_sp_tree.remove(child)

    for child in list(src_sp_tree)[2:]:
        dst_sp_tree.append(deepcopy(child))


def cmd_add_slide(args: AddSlideArgs) -> int:
    prs = open_presentation(args.file)
    layout = find_layout(prs, args.layout)
    total_before = len(prs.slides)
    out_path = args.out or args.file

    slide = prs.slides.add_slide(layout)

    if args.at is not None:
        at_pos = args.at
        if at_pos < 1 or at_pos > total_before + 1:
            print(f"Error: --at {at_pos} out of range (1-{total_before + 1})", file=sys.stderr)
            return 1

        # Reorder slide by moving its <p:sldId> entry in the presentation XML.
        sld_id_list = prs.slides.element
        sld_ids = sld_id_list.sldId_lst
        if not sld_ids:
            print("Error: presentation has no slide ID list", file=sys.stderr)
            return 1

        new_sld_id = sld_ids[-1]
        sld_id_list.remove(new_sld_id)

        insert_idx = at_pos - 1
        if insert_idx < len(sld_id_list.sldId_lst):
            sld_id_list.insert(insert_idx, new_sld_id)
        else:
            sld_id_list.append(new_sld_id)

        final_pos = at_pos
    else:
        final_pos = total_before + 1

    phs: list[str] = [f"{ph.placeholder_format.idx}:{ph.name}" for ph in slide.placeholders]
    print(f'Added slide {final_pos} from layout "{layout.name}"')
    print(f"  Placeholders: {', '.join(phs) if phs else '(none)'}")
    print(f"  Total slides: {total_before + 1}")

    prs.save(out_path)
    print(f"Saved → {out_path}")
    return 0


def cmd_delete_slide(args: DeleteSlideArgs) -> int:
    if not args.confirm:
        print("Error: delete-slide requires --confirm flag", file=sys.stderr)
        return 1

    prs = open_presentation(args.file)
    total = len(prs.slides)
    out_path = args.out or args.file

    if total <= 1:
        print("Error: cannot delete the only slide in the presentation", file=sys.stderr)
        return 1

    slide = get_slide(prs, args.slide)

    title = "(no title)"
    title_shape = slide.shapes.title
    if title_shape is not None:
        title = title_shape.text.replace("\x0b", " | ").replace("\n", " | ")[:60]

    print(f'Deleting slide {args.slide}: "{title}" ({len(slide.shapes)} shapes)')

    # Find relationship id (rId) for this slide part
    slide_part = slide.part
    r_id: str | None = None
    for rel_key in prs.part.rels:
        rel = prs.part.rels[rel_key]
        if rel.target_part is slide_part:
            r_id = rel_key
            break

    if r_id is None:
        print("Error: could not find slide relationship", file=sys.stderr)
        return 1

    # Remove the <p:sldId> element referencing this slide
    sld_id_list = prs.slides.element
    for sld_id in list(sld_id_list.sldId_lst):
        if sld_id.rId == r_id:
            sld_id_list.remove(sld_id)
            break

    prs.part.rels.pop(r_id)

    print(f"  Remaining slides: {total - 1}")
    prs.save(out_path)
    print(f"Saved → {out_path}")
    return 0


def cmd_delete_shape(args: DeleteShapeArgs) -> int:
    prs = open_presentation(args.file)
    slide = get_slide(prs, args.slide)
    shape = find_shape(slide, args.shape)
    out_path = args.out or args.file

    text_preview = ""
    if is_text_shape(shape):
        text_preview = shape.text_frame.text.replace("\x0b", " ").replace("\n", " ")[:50]

    print(f'Deleting shape: "{shape.name}"')
    if text_preview:
        print(f"  Text: {text_preview}")

    slide.shapes.element.remove(shape.element)

    prs.save(out_path)
    print(f"Saved → {out_path}")
    return 0


def cmd_insert(args: InsertArgs) -> int:
    """Insert slides from another PPTX into this presentation."""
    prs = open_presentation(args.file)
    src_prs = open_presentation(args.source)
    out_path = args.out or args.file

    total_before = len(prs.slides)
    src_total = len(src_prs.slides)

    try:
        selected = parse_slide_selection(args.slides, src_total)
    except ValueError as e:
        print(f"Error: {e}", file=sys.stderr)
        return 1

    if args.at is None:
        at_pos = total_before + 1
    else:
        at_pos = args.at
        if at_pos < 1 or at_pos > total_before + 1:
            print(f"Error: --at {at_pos} out of range (1-{total_before + 1})", file=sys.stderr)
            return 1

    insert_pos = at_pos
    for slide_num in selected:
        src_slide = get_slide(src_prs, slide_num)

        dst_layout: SlideLayout | None = None
        src_layout_name = getattr(src_slide.slide_layout, "name", "")
        if src_layout_name:
            dst_layout = try_find_layout(prs, src_layout_name)
        if dst_layout is None and src_layout_name:
            print(
                f"  WARNING: slide {slide_num} layout '{src_layout_name}' not found "
                f"in destination; falling back to Default. "
                f"(Native PowerPoint would import the source layout instead.)"
            )
        if dst_layout is None:
            dst_layout = try_find_layout(prs, "default")
        if dst_layout is None:
            dst_layout = find_layout(prs, "default", fallback=True)
            print(
                f"  WARNING: slide {slide_num} — 'Default' layout not found either; "
                f"using '{dst_layout.name}'."
            )

        dst_slide = prs.slides.add_slide(dst_layout)
        _replace_slide_shapes(dst_slide, src_slide)

        # Copy hyperlink relationships from source to destination (remapping rIds).
        # Raises ValueError for embedded content (images/charts) which isn't supported yet.
        try:
            _check_and_copy_relationships(src_slide, dst_slide)
        except ValueError as e:
            print(f"Error: slide {slide_num}: {e}", file=sys.stderr)
            return 1

        # Reorder to requested insertion position.
        _move_last_slide_to(prs, insert_pos)
        insert_pos += 1

    prs.save(out_path)
    print(
        f"Inserted {len(selected)} slide(s) from {args.source} into {args.file} at position {at_pos} → {out_path}"
    )
    return 0
