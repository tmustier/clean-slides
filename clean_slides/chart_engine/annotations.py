"""Chart annotation and text-shape helpers."""

from __future__ import annotations

from collections.abc import Mapping
from pathlib import Path
from typing import Any, Union, cast

from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_AUTO_SIZE, MSO_VERTICAL_ANCHOR, PP_ALIGN
from pptx.oxml import ns as _oxml_ns
from pptx.oxml.ns import qn
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Pt

from .colors import apply_color, normalize_theme_color, resolve_color
from .defaults import (
    DEFAULT_BAR_SEGMENT_LABEL_FONT_SIZE,
    DEFAULT_WATERFALL_LABEL_FONT_SIZE,
    DEFAULT_WATERFALL_TITLE_FONT_SIZE,
    DEFAULT_WATERFALL_TITLE_HEIGHT,
)
from .text_style import normalize_alignment, normalize_vertical_anchor
from .text_templates import apply_txbody_template, resolve_txbody_template
from .units import coerce_emu, coerce_line_width, resolve_path

Number = Union[int, float]
ColorValue = Union[RGBColor, str, None]
SpecMap = Mapping[str, object]
BaseDirValue = Union[str, Path, None]

_NS_A14 = "http://schemas.microsoft.com/office/drawing/2010/main"
_ns_map = cast(dict[str, str], _oxml_ns.__dict__.get("_nsmap", {}))
if "a14" not in _ns_map:
    _ns_map["a14"] = _NS_A14

_DEFAULT_TEXT_COLOR = RGBColor(0, 0, 0)


def _mapping(value: object) -> dict[str, object]:
    if not isinstance(value, dict):
        return {}

    typed_mapping = cast(dict[object, object], value)
    mapped: dict[str, object] = {}
    for key, item in typed_mapping.items():
        if isinstance(key, str):
            mapped[key] = item
    return mapped


def _str(value: object, default: str = "") -> str:
    return value if isinstance(value, str) else default


def _str_or_none(value: object) -> str | None:
    return value if isinstance(value, str) else None


def _bool(value: object, default: bool = False) -> bool:
    return value if isinstance(value, bool) else default


def _number(value: object, default: Number = 0) -> Number:
    return value if isinstance(value, (int, float)) else default


def _color(value: object) -> ColorValue:
    if isinstance(value, (RGBColor, str)):
        return value
    return None


def _base_dir(value: object) -> BaseDirValue:
    if isinstance(value, (str, Path)):
        return value
    return None


def _theme_token(value: str) -> str | None:
    if normalize_theme_color(value) is None:
        return None
    return value.strip().lower().replace("_", "")


def _rgb_hex(rgb: RGBColor) -> str:
    return bytes(rgb).hex().upper()


def add_text_label(
    slide: Any,
    text: str,
    x: float,
    y: float,
    width: float,
    height: float,
    align: PP_ALIGN = PP_ALIGN.CENTER,
    color: ColorValue = _DEFAULT_TEXT_COLOR,
    font_size: Pt | int | float = DEFAULT_WATERFALL_LABEL_FONT_SIZE,
    fill_color: ColorValue = None,
    shape_type: MSO_SHAPE | None = None,
    margin_left: object | None = None,
    margin_right: object | None = None,
    margin_top: object | None = None,
    margin_bottom: object | None = None,
    vertical_anchor: MSO_VERTICAL_ANCHOR | str | None = None,
    bold: bool | None = None,
    line_color: ColorValue = None,
    line_width: object | None = None,
    bw_mode: str | None = None,
    txbody_template: object | None = None,
) -> Any:
    x_int = round(x)
    y_int = round(y)
    width_int = round(width)
    height_int = round(height)

    if shape_type is None:
        box = slide.shapes.add_textbox(x_int, y_int, width_int, height_int)
    else:
        box = slide.shapes.add_shape(shape_type, x_int, y_int, width_int, height_int)

    box.text_frame.text = text
    box.text_frame.word_wrap = False
    box.text_frame.auto_size = MSO_AUTO_SIZE.NONE

    if vertical_anchor is not None:
        anchor = (
            normalize_vertical_anchor(vertical_anchor)
            if isinstance(vertical_anchor, str)
            else vertical_anchor
        )
        if anchor is not None:
            box.text_frame.vertical_anchor = anchor

    if margin_left is not None:
        box.text_frame.margin_left = coerce_emu(margin_left) or 0
    if margin_right is not None:
        box.text_frame.margin_right = coerce_emu(margin_right) or 0
    if margin_top is not None:
        box.text_frame.margin_top = coerce_emu(margin_top) or 0
    if margin_bottom is not None:
        box.text_frame.margin_bottom = coerce_emu(margin_bottom) or 0

    paragraph = box.text_frame.paragraphs[0]
    paragraph.alignment = align
    if isinstance(font_size, Pt):
        paragraph.font.size = font_size
    else:
        paragraph.font.size = Pt(
            float(_number(font_size, float(DEFAULT_WATERFALL_LABEL_FONT_SIZE)))
        )

    if bold is not None:
        paragraph.font.bold = bool(bold)
    if color is not None:
        apply_color(paragraph.font.color, color)

    if fill_color is None:
        box.fill.background()
    else:
        box.fill.solid()
        if not apply_color(box.fill.fore_color, fill_color):
            box.fill.background()

    if line_color is None:
        box.line.fill.background()
    else:
        box.line.fill.solid()
        apply_color(box.line.color, line_color)
        if line_width is not None:
            line_width_emu = coerce_line_width(line_width)
            if line_width_emu is not None:
                box.line.width = line_width_emu

    resolved_bw_mode = bw_mode
    if resolved_bw_mode is None and shape_type is None:
        resolved_bw_mode = "gray" if fill_color is not None else "auto"

    box_element = box._element

    if resolved_bw_mode:
        shape_props = box_element.find(qn("p:spPr"))
        if shape_props is not None:
            shape_props.set("bwMode", resolved_bw_mode)

    if shape_type is None:
        shape_props = box_element.find(qn("p:spPr"))
        if shape_props is not None:
            preset_geometry = shape_props.find(qn("a:prstGeom"))
            if preset_geometry is not None:
                preset_geometry.set("prst", "rect")

            ext_uri = "{909E8E84-426E-40DD-AFC4-6F175D3DCCD1}"
            ext_list = shape_props.find(qn("a:extLst"))
            if fill_color is None:
                if ext_list is None:
                    ext_list = OxmlElement("a:extLst")
                    shape_props.append(ext_list)

                existing_ext = None
                for child in ext_list:
                    if child.tag == qn("a:ext") and child.get("uri") == ext_uri:
                        existing_ext = child
                        break

                if existing_ext is None:
                    ext = OxmlElement("a:ext")
                    ext.set("uri", ext_uri)
                    hidden_fill = OxmlElement("a14:hiddenFill")
                    solid_fill = OxmlElement("a:solidFill")
                    scheme = OxmlElement("a:schemeClr")
                    scheme.set("val", "accent1")
                    solid_fill.append(scheme)
                    hidden_fill.append(solid_fill)
                    ext.append(hidden_fill)
                    ext_list.append(ext)

            effect_list = shape_props.find(qn("a:effectLst"))
            if effect_list is None:
                effect_list = OxmlElement("a:effectLst")
                if ext_list is not None:
                    shape_props.insert(shape_props.index(ext_list), effect_list)
                else:
                    shape_props.append(effect_list)

    if txbody_template is not None:
        apply_txbody_template(box, txbody_template, text)

    return box


def set_line_endings(
    line: Any,
    head: Mapping[str, object] | None = None,
    tail: Mapping[str, object] | None = None,
) -> None:
    line_props = line._get_or_add_ln()

    def apply_end(tag: str, data: Mapping[str, object]) -> None:
        element = line_props.find(qn(tag))
        if element is None:
            element = OxmlElement(tag)
            line_props.append(element)

        element.attrib.clear()
        element.set("type", _str(data.get("type"), "none"))
        element.set("w", _str(data.get("w"), "med"))
        element.set("len", _str(data.get("len"), "med"))

    if head is not None:
        apply_end("a:headEnd", head)
    if tail is not None:
        apply_end("a:tailEnd", tail)


def add_line_annotation(slide: Any, spec: SpecMap) -> None:
    spec_map = _mapping(spec)

    x = coerce_emu(spec_map.get("x"))
    y = coerce_emu(spec_map.get("y"))

    raw_width = spec_map.get("w") if spec_map.get("w") is not None else spec_map.get("width")
    raw_height = spec_map.get("h") if spec_map.get("h") is not None else spec_map.get("height")
    width = coerce_emu(raw_width)
    height = coerce_emu(raw_height)

    if x is None or y is None or width is None or height is None:
        return

    x_int = round(x)
    y_int = round(y)
    width_int = round(width)
    height_int = round(height)

    line = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        x_int,
        y_int,
        x_int + (width_int or 1),
        y_int + (height_int or 1),
    )

    shape_props = line._element.find(qn("p:spPr"))
    if shape_props is not None:
        shape_props.set("bwMode", "auto")
        transform = shape_props.find(qn("a:xfrm"))
        if transform is not None:
            offset = transform.find(qn("a:off"))
            extent = transform.find(qn("a:ext"))
            if offset is not None:
                offset.set("x", str(x_int))
                offset.set("y", str(y_int))
            if extent is not None:
                extent.set("cx", str(width_int))
                extent.set("cy", str(height_int))

    line_width = spec_map.get("line_width")
    if line_width is not None:
        line_width_emu = coerce_line_width(line_width)
        if line_width_emu is not None:
            line.line.width = line_width_emu

    line_color = _color(spec_map.get("line_color"))
    if line_color is not None:
        line.line.fill.solid()
        apply_color(line.line.color, line_color)

    dash_value = _str_or_none(spec_map.get("dash_style"))
    if dash_value:
        dash_map = {
            "solid": MSO_LINE_DASH_STYLE.SOLID,
            "dash": MSO_LINE_DASH_STYLE.DASH,
            "long_dash": MSO_LINE_DASH_STYLE.LONG_DASH,
            "dot": MSO_LINE_DASH_STYLE.ROUND_DOT,
        }
        dash_key = dash_value.lower()
        dash_style = dash_map.get(dash_key)
        if dash_style is not None:
            line.line.dash_style = dash_style

        if dash_key in {"solid", "dash", "long_dash", "dot"}:
            line_props = line.line._get_or_add_ln()
            preset_dash = line_props.find(qn("a:prstDash"))
            if preset_dash is None:
                preset_dash = OxmlElement("a:prstDash")
                line_props.append(preset_dash)
            preset_dash.set("val", "solid" if dash_key == "solid" else dash_key)

    if _bool(spec_map.get("round"), False):
        line_props = line.line._get_or_add_ln()
        if line_props.find(qn("a:round")) is None:
            line_props.append(OxmlElement("a:round"))

    head_end = _mapping(spec_map.get("head_end"))
    tail_end = _mapping(spec_map.get("tail_end"))
    if head_end or tail_end:
        set_line_endings(line.line, head=head_end or None, tail=tail_end or None)

    cap = _str_or_none(spec_map.get("cap"))
    if cap:
        line_props = line.line._get_or_add_ln()
        line_props.set("cap", cap)

    cmpd = _str_or_none(spec_map.get("cmpd"))
    if cmpd:
        line_props = line.line._get_or_add_ln()
        line_props.set("cmpd", cmpd)

    align = _str_or_none(spec_map.get("algn"))
    if align:
        line_props = line.line._get_or_add_ln()
        line_props.set("algn", align)

    if _bool(spec_map.get("flip_v"), False):
        transform = line._element.find(qn("p:spPr"))
        if transform is not None:
            transform = transform.find(qn("a:xfrm"))
        if transform is not None:
            transform.set("flipV", "1")

    if _bool(spec_map.get("flip_h"), False):
        transform = line._element.find(qn("p:spPr"))
        if transform is not None:
            transform = transform.find(qn("a:xfrm"))
        if transform is not None:
            transform.set("flipH", "1")

    shape_props = line._element.find(qn("p:spPr"))
    if shape_props is not None and shape_props.find(qn("a:effectLst")) is None:
        shape_props.append(OxmlElement("a:effectLst"))


def add_shape_annotation(slide: Any, spec: SpecMap) -> None:
    spec_map = _mapping(spec)

    shape_type = _str_or_none(spec_map.get("shape"))
    if shape_type == "ellipse":
        shape_enum = MSO_SHAPE.OVAL
    elif shape_type == "rectangle":
        shape_enum = MSO_SHAPE.RECTANGLE
    else:
        shape_enum = None

    x = coerce_emu(spec_map.get("x"))
    y = coerce_emu(spec_map.get("y"))

    raw_width = spec_map.get("w") if spec_map.get("w") is not None else spec_map.get("width")
    raw_height = spec_map.get("h") if spec_map.get("h") is not None else spec_map.get("height")
    width = coerce_emu(raw_width)
    height = coerce_emu(raw_height)

    if x is None or y is None or width is None or height is None:
        return

    use_text_box = shape_type == "ellipse"

    font_size_value = spec_map.get("font_size", DEFAULT_BAR_SEGMENT_LABEL_FONT_SIZE)
    if isinstance(font_size_value, Pt):
        font_size = font_size_value
    else:
        font_size = Pt(float(_number(font_size_value, float(DEFAULT_BAR_SEGMENT_LABEL_FONT_SIZE))))

    txbody_template: object | None = None
    base_dir = _base_dir(spec_map.get("_base_dir"))
    text_style_template = _str_or_none(spec_map.get("text_style_template"))
    text_value = _str(spec_map.get("text"), "")

    if text_style_template and text_value:
        template_path = resolve_path(text_style_template, base_dir)
        txbody_template = resolve_txbody_template(template_path, text_value, None)

    shape = add_text_label(
        slide,
        text_value,
        float(x),
        float(y),
        float(width),
        float(height),
        align=normalize_alignment(_str_or_none(spec_map.get("align"))) or PP_ALIGN.CENTER,
        color=_color(spec_map.get("text_color")),
        font_size=font_size,
        fill_color=_color(spec_map.get("fill_color")),
        shape_type=None if use_text_box else shape_enum,
        margin_left=spec_map.get("margin_left"),
        margin_right=spec_map.get("margin_right"),
        margin_top=spec_map.get("margin_top"),
        margin_bottom=spec_map.get("margin_bottom"),
        vertical_anchor=_str_or_none(spec_map.get("vertical_anchor")),
        bold=cast(
            bool | None, spec_map.get("bold") if isinstance(spec_map.get("bold"), bool) else None
        ),
        line_color=_color(spec_map.get("line_color")),
        line_width=spec_map.get("line_width"),
    )

    if shape is not None and txbody_template is not None:
        template_text = "".join(
            _str(getattr(t_elem, "text", ""), "")
            for t_elem in cast(Any, txbody_template).iter(qn("a:t"))
        )
        override_text = None if template_text == text_value else text_value
        apply_txbody_template(shape, txbody_template, override_text)

    if not use_text_box or shape is None:
        return

    shape_props = shape._element.find(qn("p:spPr"))
    if shape_props is None:
        return

    shape_props.set("bwMode", "auto")
    preset_geometry = shape_props.find(qn("a:prstGeom"))
    if preset_geometry is not None:
        preset_geometry.set("prst", "ellipse")

    line_color = _color(spec_map.get("line_color"))
    if line_color is None:
        return

    line_props = shape_props.find(qn("a:ln"))
    if line_props is None:
        line_props = OxmlElement("a:ln")
        shape_props.append(line_props)

    line_width_value = spec_map.get("line_width")
    if line_width_value is not None:
        line_width_emu = coerce_line_width(line_width_value)
        if line_width_emu is not None:
            line_props.set("w", str(line_width_emu))

    cmpd = _str_or_none(spec_map.get("cmpd"))
    if cmpd:
        line_props.set("cmpd", cmpd)

    for child in list(line_props):
        if child.tag == qn("a:solidFill"):
            line_props.remove(child)

    if line_props.find(qn("a:noFill")) is None:
        line_props.append(OxmlElement("a:noFill"))

    ext_uri = "{91240B29-F687-4F45-9708-019B960494DF}"
    ext_list = shape_props.find(qn("a:extLst"))
    if ext_list is None:
        ext_list = OxmlElement("a:extLst")
        shape_props.append(ext_list)

    ext = None
    for child in ext_list:
        if child.tag == qn("a:ext") and child.get("uri") == ext_uri:
            ext = child
            break

    if ext is None:
        ext = OxmlElement("a:ext")
        ext.set("uri", ext_uri)
        ext_list.append(ext)

    hidden_line = ext.find(qn("a14:hiddenLine"))
    if hidden_line is None:
        hidden_line = OxmlElement("a14:hiddenLine")
        ext.append(hidden_line)

    if line_width_value is not None:
        line_width_emu = coerce_line_width(line_width_value)
        if line_width_emu is not None:
            hidden_line.set("w", str(line_width_emu))

    if cmpd:
        hidden_line.set("cmpd", cmpd)

    solid_fill = hidden_line.find(qn("a:solidFill"))
    if solid_fill is None:
        solid_fill = OxmlElement("a:solidFill")
        hidden_line.append(solid_fill)

    if isinstance(line_color, str):
        theme_token = _theme_token(line_color)
        if theme_token is not None:
            scheme = solid_fill.find(qn("a:schemeClr"))
            if scheme is None:
                scheme = OxmlElement("a:schemeClr")
                solid_fill.append(scheme)
            scheme.set("val", theme_token)
        else:
            rgb, _theme = resolve_color(line_color)
            if rgb is not None:
                srgb = solid_fill.find(qn("a:srgbClr"))
                if srgb is None:
                    srgb = OxmlElement("a:srgbClr")
                    solid_fill.append(srgb)
                srgb.set("val", _rgb_hex(rgb))
    else:
        srgb = solid_fill.find(qn("a:srgbClr"))
        if srgb is None:
            srgb = OxmlElement("a:srgbClr")
            solid_fill.append(srgb)
        srgb.set("val", _rgb_hex(line_color))


def add_waterfall_title(
    slide: Any, chart_box: tuple[int, int, int, int], title: str, offset: Number
) -> None:
    x, y, width, _ = chart_box
    title_y = y - offset
    add_text_label(
        slide,
        title,
        x,
        title_y,
        width,
        int(DEFAULT_WATERFALL_TITLE_HEIGHT),
        align=PP_ALIGN.CENTER,
        font_size=DEFAULT_WATERFALL_TITLE_FONT_SIZE,
    )
