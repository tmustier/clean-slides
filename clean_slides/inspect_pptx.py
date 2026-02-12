"""
Read-only introspection for PowerPoint slides.

All positions in inches (4 decimal places). Colors as hex RGB + theme name.
Font sizes in points. Returns dataclasses with .to_dict() for serialization.
"""

from __future__ import annotations

from dataclasses import asdict, dataclass, field
from typing import Any, Protocol, TypedDict, cast

from lxml import etree
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT
from pptx.oxml.chart.chart import CT_ChartSpace
from pptx.oxml.ns import qn
from pptx.presentation import Presentation
from pptx.shapes.autoshape import Shape as PptxShape
from pptx.shapes.base import BaseShape
from pptx.shapes.graphfrm import GraphicFrame
from pptx.slide import Slide, SlideLayout
from pptx.text.text import TextFrame
from typing_extensions import TypeGuard

from .pptx_access import chart_xml_space, paragraph_xml_element

# ── Helpers ────────────────────────────────────────────────────────────


XmlElement = Any


def _is_text_shape(shape: BaseShape) -> TypeGuard[PptxShape]:
    return isinstance(shape, PptxShape) and shape.has_text_frame


def _is_chart_frame(shape: BaseShape) -> TypeGuard[_ChartFrameLike]:
    return isinstance(shape, GraphicFrame) and shape.has_chart


def _emu_to_inches(emu: int | None) -> float | None:
    """Convert EMU to inches, rounded to 4 decimal places."""
    if emu is None:
        return None
    return round(int(emu) / 914400, 4)


def _pt_to_float(pt_val: int | None) -> float | None:
    """Convert EMU (or Length) to float points."""
    if pt_val is None:
        return None
    return round(int(pt_val) / 12700, 1)


def _resolve_rPr_color(rPr: XmlElement | None) -> tuple[str | None, str | None]:
    """Extract color from an a:rPr XML element directly."""
    if rPr is None:
        return None, None

    # Check solidFill
    solidFill = rPr.find(qn("a:solidFill"))
    if solidFill is not None:
        schemeClr = solidFill.find(qn("a:schemeClr"))
        if schemeClr is not None:
            return None, schemeClr.get("val")
        srgbClr = solidFill.find(qn("a:srgbClr"))
        if srgbClr is not None:
            return f"#{srgbClr.get('val')}", None

    return None, None


def _get_fill_info(element: XmlElement) -> FillInfo:
    """Extract fill information from an spPr-containing element."""
    for ns in ["p", "a"]:
        spPr = element.find(qn(f"{ns}:spPr"))
        if spPr is not None:
            return _parse_spPr_fill(spPr)
    return FillInfo()


def _parse_spPr_fill(spPr: XmlElement | None) -> FillInfo:
    """Parse fill from an spPr element."""
    if spPr is None:
        return FillInfo()

    if spPr.find(qn("a:noFill")) is not None:
        return FillInfo(type="none")

    solidFill = spPr.find(qn("a:solidFill"))
    if solidFill is not None:
        schemeClr = solidFill.find(qn("a:schemeClr"))
        if schemeClr is not None:
            return FillInfo(type="theme", theme_color=schemeClr.get("val"))
        srgbClr = solidFill.find(qn("a:srgbClr"))
        if srgbClr is not None:
            return FillInfo(type="solid", color_rgb=f"#{srgbClr.get('val')}")

    if spPr.find(qn("a:gradFill")) is not None:
        return FillInfo(type="gradient")
    if spPr.find(qn("a:pattFill")) is not None:
        return FillInfo(type="pattern")
    if spPr.find(qn("a:blipFill")) is not None:
        return FillInfo(type="picture")

    return FillInfo(type="inherited")


def _get_line_info(element: XmlElement) -> LineInfo:
    """Extract line/outline info from an spPr-containing element."""
    for ns in ["p", "a"]:
        spPr = element.find(qn(f"{ns}:spPr"))
        if spPr is not None:
            ln = spPr.find(qn("a:ln"))
            if ln is not None:
                return _parse_ln(ln)
    return LineInfo()


def _parse_ln(ln: XmlElement) -> LineInfo:
    """Parse line properties from an a:ln element."""
    width = ln.get("w")
    width_pt = round(int(width) / 12700, 1) if width else None

    color_rgb = None
    theme_color = None
    dash = None

    solidFill = ln.find(qn("a:solidFill"))
    if solidFill is not None:
        schemeClr = solidFill.find(qn("a:schemeClr"))
        if schemeClr is not None:
            theme_color = schemeClr.get("val")
        srgbClr = solidFill.find(qn("a:srgbClr"))
        if srgbClr is not None:
            color_rgb = f"#{srgbClr.get('val')}"

    if ln.find(qn("a:noFill")) is not None:
        return LineInfo(width=width_pt, dash_style="none")

    prstDash = ln.find(qn("a:prstDash"))
    if prstDash is not None:
        dash = prstDash.get("val")

    return LineInfo(width=width_pt, color_rgb=color_rgb, theme_color=theme_color, dash_style=dash)


def _classify_shape(shape: BaseShape) -> str:
    """Return a human-readable shape type string."""
    if shape.is_placeholder:
        return "placeholder"
    st = shape.shape_type
    type_map = {
        MSO_SHAPE_TYPE.TEXT_BOX: "text_box",
        MSO_SHAPE_TYPE.AUTO_SHAPE: "auto_shape",
        MSO_SHAPE_TYPE.CHART: "chart",
        MSO_SHAPE_TYPE.TABLE: "table",
        MSO_SHAPE_TYPE.PICTURE: "picture",
        MSO_SHAPE_TYPE.GROUP: "group",
        MSO_SHAPE_TYPE.FREEFORM: "freeform",
    }
    # Connectors
    if hasattr(shape, "begin_x"):
        return "connector"
    return type_map.get(st, str(st))


def _get_text_preview(shape: BaseShape, max_len: int = 80) -> str:
    """Get truncated text content from a shape."""
    if not _is_text_shape(shape):
        return ""

    text = shape.text_frame.text or ""
    if len(text) > max_len:
        return text[:max_len] + "…"
    return text


class SlideListEntry(TypedDict):
    slide: int
    title: str


class TypographyEntry(TypedDict):
    font: str
    size: float
    bold_mix: bool
    chars: int
    shapes: int


class SlideSummary(TypedDict):
    title: str
    subtitle: str
    tracker: str | None
    text_blocks: list[str]
    shape_count: int
    chart_count: int | None
    typography: list[TypographyEntry]


class _ParagraphLike(Protocol):
    @property
    def alignment(self) -> PP_PARAGRAPH_ALIGNMENT | None: ...

    @property
    def level(self) -> int: ...


class _ChartLike(Protocol):
    @property
    def has_legend(self) -> bool: ...


class _ChartFrameLike(Protocol):
    @property
    def has_chart(self) -> bool: ...

    @property
    def chart(self) -> _ChartLike: ...


def _chart_space(chart: _ChartLike) -> CT_ChartSpace:
    chart_space = chart_xml_space(chart)
    if chart_space is None:
        raise ValueError("Chart object missing _chartSpace")
    return cast(CT_ChartSpace, chart_space)


def _paragraph_xml(paragraph: _ParagraphLike) -> XmlElement:
    paragraph_el = paragraph_xml_element(paragraph)
    if paragraph_el is None:
        raise ValueError("Paragraph object missing _element")
    return paragraph_el


# ── Data classes ───────────────────────────────────────────────────────


@dataclass
class FillInfo:
    type: str = (
        "inherited"  # "solid", "theme", "gradient", "pattern", "picture", "none", "inherited"
    )
    color_rgb: str | None = None
    theme_color: str | None = None

    def to_dict(self):
        return {k: v for k, v in asdict(self).items() if v is not None and v != "inherited"}

    def __str__(self):
        if self.type == "none":
            return "none"
        if self.type == "theme":
            return f"theme:{self.theme_color}"
        if self.type == "solid":
            return f"solid:{self.color_rgb}"
        return self.type


@dataclass
class LineInfo:
    width: float | None = None  # points
    color_rgb: str | None = None
    theme_color: str | None = None
    dash_style: str | None = None

    def to_dict(self):
        return {k: v for k, v in asdict(self).items() if v is not None}

    def __str__(self):
        parts: list[str] = []
        if self.width:
            parts.append(f"{self.width}pt")
        if self.theme_color:
            parts.append(self.theme_color)
        elif self.color_rgb:
            parts.append(self.color_rgb)
        if self.dash_style and self.dash_style not in ("solid", "none", None):
            parts.append(self.dash_style)
        if self.dash_style == "none":
            parts.append("hidden")
        return " ".join(parts) if parts else "none"


@dataclass
class RunInfo:
    """
    Run formatting info — field names match edit command for copy-paste workflow.

    Fields match `pptx edit` overrides:
        font, size, bold, italic, underline, color, superscript, subscript
    """

    text: str = ""
    font: str | None = None  # font name
    size: float | None = None  # points
    bold: bool | None = None
    italic: bool | None = None
    underline: bool | None = None
    color: str | None = None  # theme name (e.g. "accent1") or hex (e.g. "#FF0000")
    superscript: bool | None = None
    subscript: bool | None = None

    def to_dict(self):
        d = {"text": self.text}
        # Output in consistent order matching edit command
        for k in (
            "font",
            "size",
            "bold",
            "italic",
            "underline",
            "color",
            "superscript",
            "subscript",
        ):
            v = getattr(self, k)
            if v is not None:
                d[k] = v
        return d

    def __str__(self):
        props: list[str] = []
        if self.bold:
            props.append("B")
        if self.italic:
            props.append("I")
        if self.underline:
            props.append("U")
        if self.superscript:
            props.append("^")
        if self.subscript:
            props.append("_")
        if self.size:
            props.append(f"{self.size}pt")
        if self.color:
            props.append(self.color)
        suffix = f" [{', '.join(props)}]" if props else ""
        return f'"{self.text}"{suffix}'


def _runinfo_list() -> list[RunInfo]:
    return []


@dataclass
class ParagraphInfo:
    runs: list[RunInfo] = field(default_factory=_runinfo_list)
    alignment: str | None = None
    level: int = 0
    bullet_char: str | None = None
    spacing_before: float | None = None  # points
    spacing_after: float | None = None
    line_spacing: float | None = None  # percentage

    def to_dict(self):
        d = {"runs": [r.to_dict() for r in self.runs]}
        for k, v in asdict(self).items():
            if k != "runs" and v is not None and v != 0:
                d[k] = v
        return d

    def __str__(self):
        prefix = ""
        if self.bullet_char:
            prefix = f"{self.bullet_char} "
        text = "".join(r.text for r in self.runs)
        return f"{'  ' * self.level}{prefix}{text}"


def _paragraphinfo_list() -> list[ParagraphInfo]:
    return []


@dataclass
class TextFrameInfo:
    paragraphs: list[ParagraphInfo] = field(default_factory=_paragraphinfo_list)
    anchor: str | None = None
    word_wrap: bool | None = None
    auto_size: str | None = None
    margins: dict[str, float] | None = None  # inches

    def to_dict(self):
        return {
            "paragraphs": [p.to_dict() for p in self.paragraphs],
            "anchor": self.anchor,
            "word_wrap": self.word_wrap,
            "auto_size": self.auto_size,
            "margins": self.margins,
        }


@dataclass
class ShapeInfo:
    """Summary of a shape — returned by inspect_slide()."""

    index: int = 0
    name: str = ""
    shape_type: str = ""
    left: float = 0.0
    top: float = 0.0
    width: float = 0.0
    height: float = 0.0
    rotation: float = 0.0
    text_preview: str = ""
    placeholder_idx: int | None = None
    fill: FillInfo = field(default_factory=FillInfo)
    line: LineInfo = field(default_factory=LineInfo)

    def to_dict(self):
        d = asdict(self)
        d["fill"] = self.fill.to_dict()
        d["line"] = self.line.to_dict()
        # Drop defaults
        if not d["text_preview"]:
            del d["text_preview"]
        if d["rotation"] == 0.0:
            del d["rotation"]
        if d["placeholder_idx"] is None:
            del d["placeholder_idx"]
        return d

    def __str__(self):
        parts = [
            f"[{self.index:2d}] {self.shape_type:<12s} {self.name:<25s}",
            f"({self.left:.2f}, {self.top:.2f}) {self.width:.2f}×{self.height:.2f}",
        ]
        if self.placeholder_idx is not None:
            parts.append(f"ph={self.placeholder_idx}")
        if self.fill.type not in ("inherited",):
            parts.append(f"fill={self.fill}")
        if self.text_preview:
            preview = self.text_preview[:50]
            parts.append(f'"{preview}"')
        return "  ".join(parts)


@dataclass
class ShapeDetail(ShapeInfo):
    """Full detail for a shape — returned by inspect_shape()."""

    text_frame: TextFrameInfo | None = None

    def to_dict(self):
        d = super().to_dict()
        if self.text_frame:
            d["text_frame"] = self.text_frame.to_dict()
        return d


def _object_list() -> list[object]:
    return []


@dataclass
class SeriesInfo:
    index: int = 0
    name: str = ""
    values: list[object] = field(default_factory=_object_list)
    fill: FillInfo = field(default_factory=FillInfo)

    def to_dict(self):
        return {
            "index": self.index,
            "name": self.name,
            "values": self.values,
            "fill": self.fill.to_dict(),
        }


@dataclass
class AxisInfo:
    type: str = ""  # "value" or "category"
    min_val: float | None = None
    max_val: float | None = None
    visible: bool = True
    gridlines: bool = False

    def to_dict(self):
        return {k: v for k, v in asdict(self).items() if v is not None}


def _seriesinfo_list() -> list[SeriesInfo]:
    return []


def _str_list() -> list[str]:
    return []


def _axisinfo_list() -> list[AxisInfo]:
    return []


@dataclass
class ChartInfo:
    chart_type: str = ""
    series: list[SeriesInfo] = field(default_factory=_seriesinfo_list)
    categories: list[str] = field(default_factory=_str_list)
    axes: list[AxisInfo] = field(default_factory=_axisinfo_list)
    gap_width: int | None = None
    plot_area: dict[str, float] | None = None
    has_legend: bool = False

    def to_dict(self):
        d = {
            "chart_type": self.chart_type,
            "series": [s.to_dict() for s in self.series],
            "categories": self.categories,
            "axes": [a.to_dict() for a in self.axes],
            "gap_width": self.gap_width,
            "plot_area": self.plot_area,
            "has_legend": self.has_legend,
        }
        return {k: v for k, v in d.items() if v}


@dataclass
class PlaceholderInfo:
    idx: int = 0
    name: str = ""
    type: str = ""
    left: float = 0.0
    top: float = 0.0
    width: float = 0.0
    height: float = 0.0
    default_font: str | None = None
    default_size: float | None = None
    default_bold: bool | None = None
    default_alignment: str | None = None

    def to_dict(self):
        return {k: v for k, v in asdict(self).items() if v is not None}

    def __str__(self):
        return (
            f"  ph[{self.idx:2d}] {self.type:<12s} {self.name:<25s} "
            f"({self.left:.2f}, {self.top:.2f}) {self.width:.2f}×{self.height:.2f}"
            f"  font={self.default_font} {self.default_size}pt"
        )


# ── Public API ─────────────────────────────────────────────────────────


@dataclass
class CommentInfo:
    """A single slide comment."""

    author: str
    text: str
    date: str = ""

    def to_dict(self) -> dict[str, object]:
        return asdict(self)


def get_slide_comments(slide: Slide) -> list[CommentInfo]:
    """Extract native PowerPoint comments from a slide.

    Comments are stored in a separate part (``ppt/comments/commentN.xml``)
    linked via a relationship.  python-pptx has no high-level API for them
    so we parse the XML directly.
    """
    NS_P = "http://schemas.openxmlformats.org/presentationml/2006/main"
    REL_COMMENTS = "http://schemas.openxmlformats.org/officeDocument/2006/relationships/comments"
    REL_AUTHORS = (
        "http://schemas.openxmlformats.org/officeDocument/2006/relationships/commentAuthors"
    )

    # Build author map from presentation-level commentAuthors.xml
    authors: dict[str, str] = {}  # authorId → name
    try:
        pres_part = slide.part.package.presentation_part  # type: ignore[union-attr]
        for rel in pres_part.rels.values():
            if rel.reltype == REL_AUTHORS:
                root = etree.fromstring(rel.target_part.blob)
                for cm_author in root.findall(f"{{{NS_P}}}cmAuthor"):
                    aid = cm_author.get("id", "")
                    name = cm_author.get("name", "Unknown")
                    authors[aid] = name
                break
    except Exception:
        pass

    # Find comment parts on this slide
    results: list[CommentInfo] = []
    for rel in slide.part.rels.values():
        if rel.reltype != REL_COMMENTS:
            continue
        try:
            root = etree.fromstring(rel.target_part.blob)
        except Exception:
            continue
        for cm in root.findall(f"{{{NS_P}}}cm"):
            author_id = cm.get("authorId", "")
            author = authors.get(author_id, f"Author {author_id}")
            text_el = cm.find(f"{{{NS_P}}}text")
            text = text_el.text if text_el is not None and text_el.text else ""
            dt = cm.get("dt", "")
            results.append(CommentInfo(author=author, text=text, date=dt))

    return results


def list_slides(prs: Presentation) -> list[SlideListEntry]:
    """
    Table of contents: slide number + title for every slide.
    Minimal context — use this first to orient, then drill into specific slides.

    Returns list of dicts: [{"slide": 1, "title": "..."}, ...]
    """
    result: list[SlideListEntry] = []
    for i, slide in enumerate(prs.slides):
        title = ""

        title_shape = slide.shapes.title
        if title_shape is not None:
            raw = title_shape.text
            title = raw.replace("\x0b", " | ").replace("\n", " | ").strip()
        else:
            # No title placeholder — scan for any text
            for shape in slide.shapes:
                if _is_text_shape(shape) and shape.text.strip():
                    title = shape.text.replace("\x0b", " | ").replace("\n", " | ").strip()[:80]
                    break

        result.append({"slide": i + 1, "title": title})

    return result


def _resolve_theme_fonts(slide: Slide) -> tuple[str | None, str | None]:
    """
    Get major (heading) and minor (body) font names from the slide's theme.
    Returns (major, minor) — typically both strings, or None if unresolvable.
    """
    try:
        master = slide.slide_layout.slide_master
        for rel in master.part.rels.values():
            if "theme" in rel.reltype:
                theme_xml = rel.target_part.blob
                root = etree.fromstring(theme_xml)
                ns = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}
                major = root.find(".//a:majorFont/a:latin", ns)
                minor = root.find(".//a:minorFont/a:latin", ns)
                return (
                    major.get("typeface") if major is not None else None,
                    minor.get("typeface") if minor is not None else None,
                )
    except Exception:
        pass
    return None, None


def _profile_typography(slide: Slide) -> list[TypographyEntry]:
    """
    Profile font usage across a slide, excluding title/subtitle placeholders.

    Returns list of dicts sorted by char count (descending):
        [{"font": "Arial", "size": 8, "bold_mix": True, "chars": 1200, "shapes": 5}, ...]

    bold_mix is True if runs in that font+size combo include both bold and non-bold,
    indicating a "bold label: regular text" pattern.

    Resolves inherited fonts from theme (runs with no explicit font name).
    """
    from collections import defaultdict

    major_font, minor_font = _resolve_theme_fonts(slide)

    def _int_set() -> set[int]:
        return set()

    @dataclass
    class _Combo:
        chars: int = 0
        shapes: set[int] = field(default_factory=_int_set)
        has_bold: bool = False
        has_nonbold: bool = False

    combos: defaultdict[tuple[str, float], _Combo] = defaultdict(_Combo)

    for shape in slide.shapes:
        # Skip title/subtitle placeholders — those have their own sizes
        if shape.is_placeholder:
            idx = shape.placeholder_format.idx
            if idx in (0, 1, 17):  # title, subtitle, tracker
                continue

        if not _is_text_shape(shape):
            continue

        shape_fonts: set[tuple[str, float]] = set()
        for p in shape.text_frame.paragraphs:
            for run in p.runs:
                if not run.text.strip():
                    continue

                font_name = run.font.name
                font_size = _pt_to_float(run.font.size) if run.font.size is not None else None
                if font_name is None and font_size is None:
                    continue

                # Resolve inherited font from theme
                if font_name is None:
                    font_name = minor_font or major_font

                size = font_size if font_size is not None else 0.0
                key = (font_name or "?", size)

                entry = combos[key]
                entry.chars += len(run.text)
                shape_fonts.add(key)

                if run.font.bold:
                    entry.has_bold = True
                else:
                    entry.has_nonbold = True

        for key in shape_fonts:
            combos[key].shapes.add(id(shape))

    result: list[TypographyEntry] = []
    for (font, size), data in combos.items():
        result.append(
            {
                "font": font,
                "size": size,
                "bold_mix": data.has_bold and data.has_nonbold,
                "chars": data.chars,
                "shapes": len(data.shapes),
            }
        )

    result.sort(key=lambda x: -x["chars"])
    return result


def summarize_slide(slide: Slide) -> SlideSummary:
    """
    Lightweight text summary of a slide — title, subtitle, and text content.
    No formatting, no positions, no XML. Just what's on the page.

    Returns dict with keys: title, subtitle, tracker, text_blocks, shape_count,
    chart_count, typography.
    """
    title = ""
    subtitle = ""
    tracker = ""
    text_blocks: list[str] = []
    chart_count = 0

    for shape in slide.shapes:
        # Placeholders
        if shape.is_placeholder:
            idx = shape.placeholder_format.idx
            text = ""
            if isinstance(shape, PptxShape):
                text = shape.text.strip()

            if idx == 0:
                title = text.replace("\x0b", " | ").replace("\n", " | ")
            elif idx == 1:
                subtitle = text.replace("\x0b", " | ").replace("\n", " | ")
            elif idx == 17:
                tracker = text
            elif text:
                text_blocks.append(text)
            continue

        # Charts
        if shape.shape_type == MSO_SHAPE_TYPE.CHART:
            chart_count += 1
            continue

        # Text content from non-placeholder shapes
        if _is_text_shape(shape):
            text = shape.text.replace("\x0b", "\n").strip()
            if text and len(text) > 2:  # skip trivial shapes
                text_blocks.append(text)

    return {
        "title": title,
        "subtitle": subtitle,
        "tracker": tracker or None,
        "text_blocks": text_blocks,
        "shape_count": len(slide.shapes),
        "chart_count": chart_count or None,
        "typography": _profile_typography(slide),
    }


def inspect_slide(slide: Slide) -> list[ShapeInfo]:
    """
    Summary of every shape on a slide, sorted by position (top, then left).

    Returns list[ShapeInfo].
    """
    shapes: list[ShapeInfo] = []
    for i, shape in enumerate(slide.shapes):
        info = ShapeInfo(
            index=i,
            name=shape.name or "",
            shape_type=_classify_shape(shape),
            left=_emu_to_inches(shape.left) or 0.0,
            top=_emu_to_inches(shape.top) or 0.0,
            width=_emu_to_inches(shape.width) or 0.0,
            height=_emu_to_inches(shape.height) or 0.0,
            rotation=shape.rotation or 0.0,
            text_preview=_get_text_preview(shape),
            placeholder_idx=shape.placeholder_format.idx if shape.is_placeholder else None,
            fill=_get_fill_info(shape.element),
            line=_get_line_info(shape.element),
        )
        shapes.append(info)

    shapes.sort(key=lambda s: (s.top, s.left))
    return shapes


def inspect_shape(shape: BaseShape) -> ShapeDetail:
    """
    Full detail for a single shape including all text formatting.

    Returns ShapeDetail.
    """
    detail = ShapeDetail(
        index=0,
        name=shape.name or "",
        shape_type=_classify_shape(shape),
        left=_emu_to_inches(shape.left) or 0.0,
        top=_emu_to_inches(shape.top) or 0.0,
        width=_emu_to_inches(shape.width) or 0.0,
        height=_emu_to_inches(shape.height) or 0.0,
        rotation=shape.rotation or 0.0,
        text_preview=_get_text_preview(shape),
        placeholder_idx=shape.placeholder_format.idx if shape.is_placeholder else None,
        fill=_get_fill_info(shape.element),
        line=_get_line_info(shape.element),
    )

    if _is_text_shape(shape):
        detail.text_frame = _parse_text_frame(shape.text_frame)

    return detail


def inspect_chart(shape: BaseShape) -> ChartInfo:
    """
    Extract chart data and styling from a chart shape.

    Returns ChartInfo.
    """
    if not _is_chart_frame(shape):
        raise ValueError("Shape is not a chart")

    chart = shape.chart
    cs = _chart_space(chart)

    chart_el = cs.find(qn("c:chart"))
    if chart_el is None:
        raise ValueError("Chart XML missing c:chart")

    plotArea = chart_el.find(qn("c:plotArea"))
    if plotArea is None:
        raise ValueError("Chart XML missing c:plotArea")

    # Determine chart type
    chart_type_map = {
        "barChart": "bar",
        "bar3DChart": "bar_3d",
        "lineChart": "line",
        "pieChart": "pie",
        "areaChart": "area",
        "scatterChart": "scatter",
        "doughnutChart": "doughnut",
    }
    chart_type = "unknown"
    chart_type_el: XmlElement | None = None
    for local_name, friendly in chart_type_map.items():
        el = plotArea.find(qn(f"c:{local_name}"))
        if el is not None:
            chart_type = friendly
            chart_type_el = el
            # Check for stacked/clustered
            grouping = el.find(qn("c:grouping"))
            if grouping is not None:
                g = grouping.get("val", "")
                if "stacked" in g.lower():
                    chart_type += "_stacked"
                elif "percent" in g.lower():
                    chart_type += "_100"
            # Check direction for bar charts
            if local_name == "barChart":
                dir_el = el.find(qn("c:barDir"))
                if dir_el is not None and dir_el.get("val") == "col":
                    chart_type = chart_type.replace("bar", "column")
            break

    # Series
    series_list: list[SeriesInfo] = []
    if chart_type_el is not None:
        for ser_el in chart_type_el.findall(qn("c:ser")):
            idx_el = ser_el.find(qn("c:idx"))
            idx = int(idx_el.get("val") or 0) if idx_el is not None else 0

            # Name
            name = ""
            tx = ser_el.find(qn("c:tx"))
            if tx is not None:
                v = tx.find(".//" + qn("c:v"))
                if v is not None and v.text:
                    name = v.text

            # Values
            values: list[object] = []
            numRef = ser_el.find(".//" + qn("c:numRef"))
            if numRef is not None:
                numCache = numRef.find(qn("c:numCache"))
                if numCache is not None:
                    for pt in numCache.findall(qn("c:pt")):
                        v = pt.find(qn("c:v"))
                        if v is not None and v.text:
                            try:
                                values.append(float(v.text))
                            except ValueError:
                                values.append(v.text)

            # Fill
            spPr = ser_el.find(qn("c:spPr"))
            fill = _parse_spPr_fill(spPr) if spPr is not None else FillInfo()

            series_list.append(SeriesInfo(index=idx, name=name, values=values, fill=fill))

    # Categories
    categories: list[str] = []
    if chart_type_el is not None:
        first_ser = chart_type_el.find(qn("c:ser"))
        if first_ser is not None:
            cat = first_ser.find(qn("c:cat"))
            if cat is not None:
                strRef = cat.find(qn("c:strRef"))
                if strRef is not None:
                    strCache = strRef.find(qn("c:strCache"))
                    if strCache is not None:
                        for pt in strCache.findall(qn("c:pt")):
                            v = pt.find(qn("c:v"))
                            if v is None:
                                categories.append("")
                            else:
                                categories.append(v.text or "")

    # Axes
    axes: list[AxisInfo] = []
    for ax_tag, ax_type in [(qn("c:valAx"), "value"), (qn("c:catAx"), "category")]:
        for ax in plotArea.findall(ax_tag):
            scaling = ax.find(qn("c:scaling"))
            min_val = max_val = None
            if scaling is not None:
                min_el = scaling.find(qn("c:min"))
                max_el = scaling.find(qn("c:max"))
                if min_el is not None:
                    min_raw = min_el.get("val")
                    if min_raw is not None:
                        min_val = float(min_raw)
                if max_el is not None:
                    max_raw = max_el.get("val")
                    if max_raw is not None:
                        max_val = float(max_raw)

            tick_lbl = ax.find(qn("c:tickLblPos"))
            visible = tick_lbl is None or tick_lbl.get("val") != "none"

            gridlines = ax.find(qn("c:majorGridlines")) is not None

            axes.append(
                AxisInfo(
                    type=ax_type,
                    min_val=min_val,
                    max_val=max_val,
                    visible=visible,
                    gridlines=gridlines,
                )
            )

    # Gap width
    gap_width = None
    if chart_type_el is not None:
        gw = chart_type_el.find(qn("c:gapWidth"))
        if gw is not None:
            gw_val = gw.get("val")
            if gw_val is not None:
                gap_width = int(gw_val)

    # Plot area layout
    plot_layout: dict[str, float] | None = None
    layout = plotArea.find(qn("c:layout"))
    if layout is not None:
        ml = layout.find(qn("c:manualLayout"))
        if ml is not None:
            layout_dict: dict[str, float] = {}
            for tag in ["x", "y", "w", "h"]:
                el = ml.find(qn(f"c:{tag}"))
                if el is not None:
                    raw_val = el.get("val")
                    if raw_val is not None:
                        layout_dict[tag] = float(raw_val)

            if layout_dict:
                plot_layout = layout_dict

    return ChartInfo(
        chart_type=chart_type,
        series=series_list,
        categories=categories,
        axes=axes,
        gap_width=gap_width,
        plot_area=plot_layout,
        has_legend=chart.has_legend,
    )


def inspect_layout(layout: SlideLayout) -> list[PlaceholderInfo]:
    """
    What placeholders does this layout define?

    Returns list[PlaceholderInfo].
    """
    placeholders: list[PlaceholderInfo] = []
    for ph in layout.placeholders:
        # Get default formatting from the placeholder's first paragraph
        default_font = None
        default_size = None
        default_bold = None
        default_alignment = None

        try:
            tf = ph.text_frame
            if tf.paragraphs:
                p = tf.paragraphs[0]
                if p.alignment is not None:
                    align_map = {0: "left", 1: "center", 2: "right", 3: "justify"}
                    default_alignment = align_map.get(int(p.alignment), str(p.alignment))

                # Check defRPr on the paragraph
                pPr = p._element.find(qn("a:pPr"))
                if pPr is not None:
                    defRPr = pPr.find(qn("a:defRPr"))
                    if defRPr is not None:
                        sz = defRPr.get("sz")
                        if sz:
                            default_size = round(int(sz) / 100, 1)
                        default_bold = defRPr.get("b") == "1" if defRPr.get("b") else None
                        latin = defRPr.find(qn("a:latin"))
                        if latin is not None:
                            default_font = latin.get("typeface")
        except Exception:
            pass

        # Placeholder type
        ph_type = "custom"
        try:
            idx = ph.placeholder_format.idx
            type_map = {
                0: "title",
                1: "subtitle",
                2: "body",
                10: "body",
                12: "slide_number",
                13: "date",
                14: "footer",
            }
            ph_type = type_map.get(idx, f"custom_{idx}")
        except Exception:
            idx = -1

        placeholders.append(
            PlaceholderInfo(
                idx=idx,
                name=ph.name or "",
                type=ph_type,
                left=_emu_to_inches(ph.left) or 0.0,
                top=_emu_to_inches(ph.top) or 0.0,
                width=_emu_to_inches(ph.width) or 0.0,
                height=_emu_to_inches(ph.height) or 0.0,
                default_font=default_font,
                default_size=default_size,
                default_bold=default_bold,
                default_alignment=default_alignment,
            )
        )

    placeholders.sort(key=lambda p: p.idx)
    return placeholders


def resolve_theme_colors(prs: Presentation) -> dict[str, str]:
    """
    Map scheme color names to hex RGB values.

    Returns dict like {"accent1": "#1D4ED8", "tx1": "#0D193B", ...}
    """
    colors: dict[str, str] = {}

    try:
        # Theme is a separate part linked via relationship from the slide master
        master = prs.slide_masters[0]
        clrScheme = None

        for rel in master.part.rels.values():
            if "theme" in str(rel.reltype).lower():
                theme_el = etree.fromstring(rel.target_part.blob)
                clrScheme = theme_el.find(".//" + qn("a:clrScheme"))
                if clrScheme is not None:
                    break

        if clrScheme is None:
            return colors

        for child in clrScheme:
            local_tag = child.tag.split("}")[-1]
            srgb = child.find(qn("a:srgbClr"))
            sys_clr = child.find(qn("a:sysClr"))

            hex_val = None
            if srgb is not None:
                hex_val = f"#{srgb.get('val')}"
            elif sys_clr is not None:
                hex_val = f"#{sys_clr.get('lastClr', sys_clr.get('val', '000000'))}"

            if hex_val:
                colors[local_tag] = hex_val

        # Add standard aliases
        alias_map = {"tx1": "dk1", "tx2": "dk2", "bg1": "lt1", "bg2": "lt2"}
        for alias, source in alias_map.items():
            if source in colors and alias not in colors:
                colors[alias] = colors[source]

    except Exception as e:
        colors["_error"] = str(e)

    return colors


def identify_color(prs: Presentation, rgb_hex: str) -> str | None:
    """
    Reverse lookup: given "#1B3D6F", find if it matches a theme color.

    Returns scheme name ("accent1") or None.
    """
    if not rgb_hex.startswith("#"):
        rgb_hex = f"#{rgb_hex}"
    rgb_hex = rgb_hex.upper()

    theme_colors = resolve_theme_colors(prs)
    for name, hex_val in theme_colors.items():
        if hex_val.upper() == rgb_hex and not name.startswith("_"):
            return name
    return None


# ── Text parsing ───────────────────────────────────────────────────────


def _parse_text_frame(tf: TextFrame) -> TextFrameInfo:
    """Parse a text frame into TextFrameInfo."""
    bodyPr = tf._element.find(qn("a:bodyPr"))
    anchor = bodyPr.get("anchor") if bodyPr is not None else None

    margins = None
    if bodyPr is not None:
        m: dict[str, float] = {}
        for attr, key in [("lIns", "left"), ("rIns", "right"), ("tIns", "top"), ("bIns", "bottom")]:
            val = bodyPr.get(attr)
            if val is not None:
                m[key] = round(int(val) / 914400, 4)
        if m:
            margins = m

    auto_size = None
    try:
        if tf.auto_size is not None:
            auto_size = str(tf.auto_size).split(".")[-1].split("(")[0].lower()
    except Exception:
        pass

    paragraphs: list[ParagraphInfo] = [_parse_paragraph(p) for p in tf.paragraphs]

    return TextFrameInfo(
        paragraphs=paragraphs,
        anchor=anchor,
        word_wrap=tf.word_wrap,
        auto_size=auto_size,
        margins=margins,
    )


def _parse_paragraph(p: _ParagraphLike) -> ParagraphInfo:
    """Parse a paragraph into ParagraphInfo."""
    alignment = None
    if p.alignment is not None:
        align_map = {0: "left", 1: "center", 2: "right", 3: "justify"}
        alignment = align_map.get(int(p.alignment), str(p.alignment))

    level = p.level or 0

    # Bullet character
    bullet_char = None
    paragraph_el = _paragraph_xml(p)
    pPr = paragraph_el.find(qn("a:pPr"))
    if pPr is not None:
        buChar = pPr.find(qn("a:buChar"))
        if buChar is not None:
            bullet_char = buChar.get("char")
        buNone = pPr.find(qn("a:buNone"))
        if buNone is not None:
            bullet_char = None

    # Spacing
    spacing_before = spacing_after = line_spacing = None
    if pPr is not None:
        spcBef = pPr.find(qn("a:spcBef"))
        if spcBef is not None:
            pts = spcBef.find(qn("a:spcPts"))
            if pts is not None:
                raw = pts.get("val")
                if raw is not None:
                    spacing_before = round(int(raw) / 100, 1)
        spcAft = pPr.find(qn("a:spcAft"))
        if spcAft is not None:
            pts = spcAft.find(qn("a:spcPts"))
            if pts is not None:
                raw = pts.get("val")
                if raw is not None:
                    spacing_after = round(int(raw) / 100, 1)
        lnSpc = pPr.find(qn("a:lnSpc"))
        if lnSpc is not None:
            pct = lnSpc.find(qn("a:spcPct"))
            if pct is not None:
                raw = pct.get("val")
                if raw is not None:
                    line_spacing = round(int(raw) / 1000, 1)

    # Parse runs (including line breaks as runs with text="\n")
    runs: list[RunInfo] = []
    for child in paragraph_el:
        tag = child.tag.split("}")[-1]
        if tag == "r":
            runs.append(_parse_run_element(child))
        elif tag == "br":
            runs.append(RunInfo(text="\n"))

    return ParagraphInfo(
        runs=runs,
        alignment=alignment,
        level=level,
        bullet_char=bullet_char,
        spacing_before=spacing_before,
        spacing_after=spacing_after,
        line_spacing=line_spacing,
    )


def _parse_run_element(r_el: XmlElement) -> RunInfo:
    """Parse an a:r element into RunInfo."""
    # Text
    t = r_el.find(qn("a:t"))
    text = (t.text or "") if t is not None else ""

    # Run properties
    rPr = r_el.find(qn("a:rPr"))
    font = None
    size = None
    bold = None
    italic = None
    underline = None
    superscript = None
    subscript = None

    if rPr is not None:
        sz = rPr.get("sz")
        if sz:
            size = round(int(sz) / 100, 1)
        # None = inherit, True = explicit on, False = explicit off
        b_attr = rPr.get("b")
        bold = None if b_attr is None else (b_attr == "1")
        i_attr = rPr.get("i")
        italic = None if i_attr is None else (i_attr == "1")
        u_attr = rPr.get("u")
        underline = None if u_attr is None else (u_attr not in ("none",))

        # Derive superscript/subscript from baseline
        bl = rPr.get("baseline")
        if bl:
            baseline = int(bl)
            if baseline > 0:
                superscript = True
            elif baseline < 0:
                subscript = True

        latin = rPr.find(qn("a:latin"))
        if latin is not None:
            typeface = latin.get("typeface")
            if typeface and not typeface.startswith("+"):
                font = typeface

    # Merge color: prefer theme name, fallback to hex
    color_rgb, theme_color = _resolve_rPr_color(rPr) if rPr is not None else (None, None)
    color = theme_color or color_rgb

    return RunInfo(
        text=text,
        font=font,
        size=size,
        bold=bold,
        italic=italic,
        underline=underline,
        color=color,
        superscript=superscript,
        subscript=subscript,
    )
