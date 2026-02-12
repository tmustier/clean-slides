"""Chart slide construction orchestration."""

from __future__ import annotations

import copy
from collections.abc import Mapping, Sequence
from pathlib import Path
from typing import Any, Callable, Union, cast

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, PP_PLACEHOLDER
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt

from ..pptx_access import chart_series_names, chart_xml_space, shape_has_text_frame, slide_charts
from . import annotations as _annotations
from . import payloads as _payloads
from .colors import apply_color
from .defaults import (
    CLEAN_SLIDES_CONTENT_BOX,
    CLEAN_SLIDES_TEMPLATE_PATH,
    DEFAULT_BAR_DATA_LABEL_FONT_SIZE,
    DEFAULT_BAR_DATA_LABEL_FORMAT,
    DEFAULT_BAR_LEGEND_LABEL_HEIGHT,
    DEFAULT_BAR_LEGEND_LABEL_OFFSET,
    DEFAULT_BAR_LEGEND_MARKER_HEIGHT,
    DEFAULT_BAR_LEGEND_MARKER_Y_OFFSET,
    DEFAULT_BAR_OVERLAY_BAND_EXTRA,
    DEFAULT_WATERFALL_CHART_BOX,
    DEFAULT_WATERFALL_TITLE_OFFSET,
    WATERFALL_TYPES,
)
from .overlays import (
    add_bar_overlays,
    add_waterfall_overlays,
    apply_waterfall_data_label_layout,
    get_chart_series,
)
from .spec_utils import normalize_list
from .style import (
    apply_bar_chart_style,
    apply_series_colors,
    apply_waterfall_chart_style,
    apply_waterfall_data_labels,
    apply_waterfall_style,
)
from .template_ops import ChartTemplateReplacement, apply_chart_template_replacements
from .text_style import normalize_label_position
from .units import coerce_emu

Number = Union[int, float]
ChartBox = tuple[int, int, int, int]
SpecMap = Mapping[str, object]

AddWaterfallTitleFn = Callable[[Any, ChartBox, str, object], None]
BuildPayloadFn = Callable[[dict[str, object]], tuple[Any, Any, dict[str, object]]]


def _require_attr(module: object, name: str) -> object:
    value = getattr(module, name, None)
    if value is None:
        raise AttributeError(f"{module!r} does not expose {name}")
    return value


add_waterfall_title = cast(
    AddWaterfallTitleFn,
    _require_attr(_annotations, "add_waterfall_title"),
)
build_bar_payload = cast(BuildPayloadFn, _require_attr(_payloads, "build_bar_payload"))
build_waterfall_payload = cast(
    BuildPayloadFn,
    _require_attr(_payloads, "build_waterfall_payload"),
)


def _mapping(value: object) -> dict[str, object]:
    if not isinstance(value, dict):
        return {}

    typed_mapping = cast(dict[object, object], value)
    mapped: dict[str, object] = {}
    for key, item in typed_mapping.items():
        if isinstance(key, str):
            mapped[key] = item
    return mapped


def _list(value: object) -> list[object]:
    if isinstance(value, list):
        return cast(list[object], value)
    if isinstance(value, tuple):
        return list(cast(tuple[object, ...], value))
    return []


def _str(value: object, default: str = "") -> str:
    return value if isinstance(value, str) else default


def _str_or_none(value: object) -> str | None:
    return value if isinstance(value, str) else None


def _bool(value: object, default: bool) -> bool:
    return value if isinstance(value, bool) else default


def _int(value: object, default: int = 0) -> int:
    if isinstance(value, bool):
        return int(value)
    if isinstance(value, int):
        return value
    if isinstance(value, float):
        return int(value)
    return default


def _parse_int(value: object) -> int:
    if isinstance(value, bool):
        return int(value)
    if isinstance(value, int):
        return value
    if isinstance(value, float):
        return int(value)
    if isinstance(value, str):
        return int(value.strip())
    raise ValueError(f"Expected integer-compatible value, got {value!r}")


def _to_float(value: object, default: float = 0.0) -> float:
    if isinstance(value, (int, float)):
        return float(value)
    if isinstance(value, str):
        try:
            return float(value.strip())
        except ValueError:
            return default
    return default


def _string_or_none_list(value: object) -> list[str | None]:
    result: list[str | None] = []
    for item in _list(value):
        if isinstance(item, str):
            result.append(item)
        elif item is None:
            result.append(None)
    return result


def _chart_box(value: tuple[object, object, object, object]) -> ChartBox | None:
    emu_values: list[int] = []
    for item in value:
        emu_value = coerce_emu(item)
        if emu_value is None:
            return None
        emu_values.append(int(emu_value))

    return (emu_values[0], emu_values[1], emu_values[2], emu_values[3])


def apply_chart_template_dlbls(
    target_chart: object,
    template_path: Path,
    slide_index: int = 0,
    chart_index: int = 0,
    series_index: int = 0,
) -> None:
    prs = Presentation(str(template_path))
    if slide_index < 0 or slide_index >= len(prs.slides):
        return

    slide = prs.slides[slide_index]
    template_charts = slide_charts(slide)
    if chart_index < 0 or chart_index >= len(template_charts):
        return

    template_chart = template_charts[chart_index]
    template_chart_space = chart_xml_space(template_chart)
    target_chart_space = chart_xml_space(target_chart)
    if template_chart_space is None or target_chart_space is None:
        return

    template_series = get_chart_series(template_chart_space)
    target_series = get_chart_series(target_chart_space)
    if not template_series or not target_series:
        return
    if series_index < 0 or series_index >= len(template_series):
        return
    if series_index >= len(target_series):
        return

    template_series_element = template_series[series_index]
    target_series_element = target_series[series_index]

    template_dlbls = template_series_element.find(qn("c:dLbls"))
    if template_dlbls is None:
        return

    existing = target_series_element.find(qn("c:dLbls"))
    if existing is not None:
        target_series_element.remove(existing)

    target_series_element.insert(len(target_series_element), copy.deepcopy(template_dlbls))


def resolve_series_indices(series_spec: object, series_names: Sequence[str]) -> list[int]:
    indices: list[int] = []
    for item in normalize_list(series_spec):
        if isinstance(item, int):
            indices.append(item)
        elif isinstance(item, str) and item in series_names:
            indices.append(series_names.index(item))
    return list(dict.fromkeys(indices))


def apply_data_label_style(labels: Any, data_cfg: SpecMap) -> None:
    labels.number_format = _str(data_cfg.get("format"), DEFAULT_BAR_DATA_LABEL_FORMAT)
    labels.number_format_is_linked = False

    font_size = data_cfg.get("font_size", DEFAULT_BAR_DATA_LABEL_FONT_SIZE)
    if isinstance(font_size, Pt):
        labels.font.size = font_size
    else:
        labels.font.size = Pt(float(_int(font_size, int(DEFAULT_BAR_DATA_LABEL_FONT_SIZE))))

    label_position = normalize_label_position(_str_or_none(data_cfg.get("position")))
    if label_position is not None:
        labels.position = label_position
    if hasattr(labels, "show_value"):
        labels.show_value = True

    color_value = data_cfg.get("color")
    if isinstance(color_value, (RGBColor, str)):
        apply_color(labels.font.color, color_value)


def chart_box_from_spec(raw: object) -> ChartBox | None:
    if not raw:
        return None

    if isinstance(raw, dict):
        typed = cast(dict[object, object], raw)
        values: list[object] = [
            typed.get("x"),
            typed.get("y"),
            typed.get("cx"),
            typed.get("cy"),
        ]
    else:
        values = _list(raw)

    if len(values) != 4:
        return None

    box = _chart_box((values[0], values[1], values[2], values[3]))
    return box


def template_content_box(slide: Any, template_path: Path | None) -> ChartBox | None:
    if template_path is None:
        return None

    try:
        if template_path.resolve() == CLEAN_SLIDES_TEMPLATE_PATH.resolve():
            return (
                int(CLEAN_SLIDES_CONTENT_BOX[0]),
                int(CLEAN_SLIDES_CONTENT_BOX[1]),
                int(CLEAN_SLIDES_CONTENT_BOX[2]),
                int(CLEAN_SLIDES_CONTENT_BOX[3]),
            )
    except FileNotFoundError:
        if template_path == CLEAN_SLIDES_TEMPLATE_PATH:
            return (
                int(CLEAN_SLIDES_CONTENT_BOX[0]),
                int(CLEAN_SLIDES_CONTENT_BOX[1]),
                int(CLEAN_SLIDES_CONTENT_BOX[2]),
                int(CLEAN_SLIDES_CONTENT_BOX[3]),
            )

    placeholder = find_content_placeholder(slide)
    if placeholder is None:
        return None

    return (
        int(placeholder.left),
        int(placeholder.top),
        int(placeholder.width),
        int(placeholder.height),
    )


def find_content_placeholder(slide: Any) -> Any | None:
    candidates: list[Any] = []
    for placeholder in slide.placeholders:
        ph_type = placeholder.placeholder_format.type
        if ph_type in (PP_PLACEHOLDER.BODY, PP_PLACEHOLDER.OBJECT):
            candidates.append(placeholder)

    if not candidates:
        return None

    return max(candidates, key=lambda ph: int(ph.width) * int(ph.height))


def remove_shape(shape: Any) -> None:
    element = shape._element
    parent = element.getparent()
    if parent is not None:
        parent.remove(element)


def adjust_bar_chart_box_for_overlays(
    chart_box: ChartBox,
    bar_meta: Mapping[str, object] | None,
) -> ChartBox:
    if not bar_meta:
        return chart_box

    overlay = _mapping(bar_meta.get("overlay"))
    if not overlay:
        return chart_box

    legend_offset = _int(overlay.get("legend_label_offset"), int(DEFAULT_BAR_LEGEND_LABEL_OFFSET))
    legend_height = int(DEFAULT_BAR_LEGEND_LABEL_HEIGHT)
    marker_height = _int(
        overlay.get("legend_marker_height"),
        int(DEFAULT_BAR_LEGEND_MARKER_HEIGHT),
    )
    marker_y_offset = _int(
        overlay.get("legend_marker_y_offset"),
        int(DEFAULT_BAR_LEGEND_MARKER_Y_OFFSET),
    )
    overlay_extra = _int(overlay.get("overlay_band_extra"), int(DEFAULT_BAR_OVERLAY_BAND_EXTRA))

    bottom_band = (
        max(
            0,
            legend_offset + legend_height,
            legend_offset + marker_y_offset + marker_height,
        )
        + overlay_extra
    )
    if bottom_band <= 0:
        return chart_box

    x, y, cx, cy = chart_box
    new_cy = max(int(Emu(100000)), int(cy) - bottom_band)
    return (int(x), int(y), int(cx), int(new_cy))


def find_layout(prs: Any, name: str | None) -> Any | None:
    if not name:
        return None

    target = name.strip().lower()
    for layout in prs.slide_layouts:
        layout_name = str(getattr(layout, "name", "")).strip().lower()
        if layout_name == target:
            return layout
    return None


def apply_template_placeholders(slide: Any, title: str | None, subtitle: str | None) -> None:
    for placeholder in slide.placeholders:
        if not shape_has_text_frame(placeholder):
            continue

        ph_type = placeholder.placeholder_format.type
        if ph_type == PP_PLACEHOLDER.TITLE:
            placeholder.text = title or ""
        elif ph_type == PP_PLACEHOLDER.SUBTITLE:
            placeholder.text = subtitle or ""
        else:
            placeholder.text = ""


def add_hidden_anchor(slide: Any) -> None:
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.02),
        Inches(0.02),
        Inches(0.02),
        Inches(0.02),
    )
    shape.name = "chart data - do not delete"
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(255, 255, 255)
    shape.line.fill.background()


def add_overlay_labels(slide: Any, categories: Sequence[object], chart_box: ChartBox) -> None:
    x, y, cx, cy = chart_box
    label_y = int(y + cy + Inches(0.1))
    if not categories:
        return

    label_width = int(cx / len(categories))
    for idx, label in enumerate(categories):
        text_box = slide.shapes.add_textbox(
            int(x + (label_width * idx)),
            label_y,
            label_width,
            int(Inches(0.3)),
        )
        text_box.name = f"tc_category_label_{idx}"
        text_frame = text_box.text_frame
        text_frame.text = str(label)
        paragraph = text_frame.paragraphs[0]
        paragraph.font.size = Pt(10)


def select_slide(
    prs: Any,
    slide_layout: Any,
    use_template: bool,
    spec: SpecMap,
) -> Any:
    template_slide_index = spec.get("template_slide_index")
    append_slide = _bool(spec.get("append_slide", False), False)

    if use_template and template_slide_index is not None:
        try:
            index = _parse_int(template_slide_index) - 1
        except ValueError as exc:
            raise ValueError("template_slide_index must be an integer (1-based)") from exc

        if index < 0 or index >= len(prs.slides):
            raise ValueError(
                f"template_slide_index {template_slide_index} is out of range (1-{len(prs.slides)})"
            )
        return prs.slides[index]

    if use_template and not append_slide and len(prs.slides) == 1:
        return prs.slides[0]

    if use_template and append_slide and len(prs.slides) == 1:
        slide = prs.slides[0]
        if not slide_charts(slide):
            return slide

    return prs.slides.add_slide(slide_layout)


def build_chart(
    prs: Any,
    spec: SpecMap,
    output_path: Path,
    template_path: Path | None = None,
    layout_name: str | None = None,
    save: bool = True,
    defer_template_copy: bool = False,
) -> list[ChartTemplateReplacement]:
    spec_map = _mapping(spec)

    slide_layout = find_layout(prs, layout_name)
    if slide_layout is None:
        slide_layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[0]

    use_template = template_path is not None
    slide = select_slide(prs, slide_layout, use_template, spec_map)

    add_hidden_anchor(slide)

    chart_type_key = _str(spec_map.get("type"), "clustered")
    is_waterfall = chart_type_key in WATERFALL_TYPES

    if is_waterfall:
        chart_type, chart_data, style_raw = build_waterfall_payload(dict(spec_map))
    else:
        chart_type, chart_data, style_raw = build_bar_payload(dict(spec_map))

    style = _mapping(style_raw)
    bar_style = _mapping(style.get("bar"))
    waterfall_style = _mapping(style.get("waterfall"))

    title = _str_or_none(spec_map.get("title"))
    subtitle = _str_or_none(spec_map.get("subtitle"))

    content_placeholder: Any | None = None
    remove_placeholder_value = spec_map.get("remove_content_placeholder")
    if remove_placeholder_value is None:
        remove_placeholder = use_template
    else:
        remove_placeholder = bool(remove_placeholder_value)

    if use_template:
        apply_template_placeholders(slide, title, subtitle)
        content_placeholder = find_content_placeholder(slide)

    if is_waterfall:
        default_box = (
            int(DEFAULT_WATERFALL_CHART_BOX[0]),
            int(DEFAULT_WATERFALL_CHART_BOX[1]),
            int(DEFAULT_WATERFALL_CHART_BOX[2]),
            int(DEFAULT_WATERFALL_CHART_BOX[3]),
        )
    else:
        default_box = (
            int(Inches(0.5)),
            int(Inches(1.0)),
            int(Inches(9.0)),
            int(Inches(4.5)),
        )

    chart_box = chart_box_from_spec(spec_map.get("chart_box")) or default_box
    waterfall_spec = _mapping(spec_map.get("waterfall"))
    waterfall_chart_box_raw = waterfall_spec.get("chart_box") if is_waterfall else None
    if is_waterfall:
        chart_box = chart_box_from_spec(waterfall_chart_box_raw) or chart_box

    template_box = template_content_box(slide, template_path)

    if content_placeholder is not None and remove_placeholder:
        remove_shape(content_placeholder)

    if template_box and waterfall_chart_box_raw is None and spec_map.get("chart_box") is None:
        chart_box = template_box
        if bar_style and _bool(spec_map.get("add_overlay_labels", False), False):
            chart_box = adjust_bar_chart_box_for_overlays(chart_box, bar_style)

    x, y, cx, cy = chart_box
    chart = slide.shapes.add_chart(chart_type, x, y, cx, cy, chart_data).chart
    chart_part = str(chart.part.partname).lstrip("/")

    if title and not is_waterfall and not use_template:
        chart.has_title = True
        chart.chart_title.text_frame.text = title

    show_legend_value = spec_map.get("show_legend")
    if show_legend_value is None:
        show_legend = not (is_waterfall or bool(bar_style))
    else:
        show_legend = bool(show_legend_value)

    chart.has_legend = show_legend
    if chart.has_legend:
        chart.legend.include_in_layout = False

    if is_waterfall:
        show_segment_labels = _bool(spec_map.get("show_data_labels", True), True)
        if show_segment_labels:
            apply_waterfall_data_labels(chart, waterfall_style)
    elif _bool(spec_map.get("show_data_labels", False), False):
        data_cfg = _mapping(spec_map.get("data_labels"))
        series_selector = (
            data_cfg.get("series_indices") or data_cfg.get("series") or data_cfg.get("series_names")
        )

        if series_selector is not None:
            series_names = chart_series_names(chart)
            indices = resolve_series_indices(series_selector, series_names)
            for idx in indices:
                if 0 <= idx < len(chart.series):
                    series = chart.series[idx]
                    series.has_data_labels = True
                    apply_data_label_style(series.data_labels, data_cfg)
        else:
            plot = chart.plots[0]
            plot.has_data_labels = True
            apply_data_label_style(plot.data_labels, data_cfg)

    apply_series_colors(chart, _string_or_none_list(style.get("series_colors")))

    if waterfall_style:
        apply_waterfall_style(chart, waterfall_style)
        apply_waterfall_chart_style(chart, waterfall_style)
        if is_waterfall and _bool(spec_map.get("show_data_labels", True), True):
            apply_waterfall_data_label_layout(chart, (x, y, cx, cy), waterfall_style)

    if bar_style:
        apply_bar_chart_style(chart, bar_style)
        bar_template = _str_or_none(bar_style.get("chart_template"))
        if bar_template:
            apply_chart_template_dlbls(
                chart,
                Path(bar_template),
                slide_index=_int(bar_style.get("chart_template_slide"), 1) - 1,
                chart_index=_int(bar_style.get("chart_template_chart_index"), 0),
                series_index=_int(bar_style.get("chart_template_series_index"), 0),
            )

    overlay_title_value = waterfall_spec.get("overlay_title")
    if overlay_title_value is None:
        overlay_title = not use_template
    else:
        overlay_title = bool(overlay_title_value)

    if is_waterfall and title and overlay_title:
        title_offset_value = waterfall_spec.get("title_offset")
        if title_offset_value is None:
            title_offset = DEFAULT_WATERFALL_TITLE_OFFSET
        else:
            title_offset = Inches(_to_float(title_offset_value))
        add_waterfall_title(slide, (x, y, cx, cy), title, title_offset)

    if _bool(spec_map.get("add_overlay_labels", False), False):
        chart_box_emu = (x, y, cx, cy)
        if is_waterfall:
            add_waterfall_overlays(
                slide,
                chart_box_emu,
                waterfall_style,
                slide_size=(int(prs.slide_width), int(prs.slide_height)),
            )
        elif bar_style:
            add_bar_overlays(slide, chart_box_emu, bar_style)
        else:
            add_overlay_labels(slide, _list(spec_map.get("categories")), chart_box_emu)

    replacements: list[ChartTemplateReplacement] = []
    if _bool(bar_style.get("chart_template_copy"), False):
        chart_template = _str_or_none(bar_style.get("chart_template"))
        if chart_template:
            replacements.append(
                ChartTemplateReplacement(
                    chart_part=chart_part,
                    template_path=Path(chart_template),
                    template_slide_index=_int(bar_style.get("chart_template_slide"), 1) - 1,
                    template_chart_index=_int(bar_style.get("chart_template_chart_index"), 0),
                )
            )

    if save:
        prs.save(output_path)
        if replacements and not defer_template_copy:
            apply_chart_template_replacements(output_path, replacements)
            replacements = []

    return replacements
