from __future__ import annotations

import json
import zipfile
from dataclasses import dataclass
from pathlib import Path

from pptx import Presentation
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.shapes import MSO_SHAPE_TYPE

from clean_slides.cli import cmd_generate
from clean_slides.cli_render import cmd_charts
from clean_slides.pptx_access import (
    iter_shapes,
    presentation_chart_types,
    shape_has_connector_endpoints,
    shape_text_frame_text,
)


@dataclass
class _GenerateArgs:
    input: list[str]
    template: str | None
    slide_index: int | None
    keep_existing: bool
    detail: bool
    output: str | None
    config: str | None = None


@dataclass
class _ChartsArgs:
    input: str
    output: str
    template: str | None = None
    layout: str | None = None
    expected_template: str | None = None


def _run_generate(
    tmp_path: Path,
    name: str,
    slides: list[tuple[str, str]],
) -> Path:
    inputs: list[str] = []
    for filename, yaml_text in slides:
        yaml_path = tmp_path / filename
        yaml_path.write_text(yaml_text.strip(), encoding="utf-8")
        inputs.append(str(yaml_path))

    output_path = tmp_path / f"{name}.pptx"
    result = cmd_generate(
        _GenerateArgs(
            input=inputs,
            template=None,
            slide_index=None,
            keep_existing=False,
            detail=False,
            output=str(output_path),
        )
    )

    assert result == 0
    assert output_path.exists()
    return output_path


def _run_charts(tmp_path: Path, name: str, spec: dict[str, object]) -> Path:
    input_path = tmp_path / f"{name}.json"
    input_path.write_text(json.dumps(spec), encoding="utf-8")

    output_path = tmp_path / f"{name}.pptx"
    result = cmd_charts(_ChartsArgs(input=str(input_path), output=str(output_path)))

    assert result == 0
    assert output_path.exists()
    return output_path


def _chart_types(path: Path) -> list[int]:
    prs = Presentation(str(path))
    return presentation_chart_types(prs)


def _slide_texts(path: Path, slide_index: int) -> list[str]:
    prs = Presentation(str(path))
    texts: list[str] = []
    for shape in iter_shapes(prs.slides[slide_index]):
        raw_text = shape_text_frame_text(shape)
        if raw_text is None:
            continue

        text = raw_text.strip()
        if text:
            texts.append(text)
    return texts


def _connector_count(path: Path, slide_index: int) -> int:
    prs = Presentation(str(path))
    count = 0
    for shape in iter_shapes(prs.slides[slide_index]):
        if shape_has_connector_endpoints(shape):
            count += 1
    return count


def _blank_auto_shape_count(path: Path, slide_index: int) -> int:
    prs = Presentation(str(path))
    count = 0
    for shape in iter_shapes(prs.slides[slide_index]):
        shape_type = getattr(shape, "shape_type", None)
        if shape_type != MSO_SHAPE_TYPE.AUTO_SHAPE:
            continue

        text = shape_text_frame_text(shape)
        if text is not None and text.strip():
            continue

        count += 1
    return count


def _pptx_part(path: Path, part_name: str) -> bytes:
    with zipfile.ZipFile(path) as archive:
        return archive.read(part_name)


def test_canary_generate_mixed_metadata_and_table_deck(tmp_path: Path) -> None:
    output_path = _run_generate(
        tmp_path,
        "canary-mixed-generate",
        [
            (
                "01-title.yaml",
                """
title: Release Canary
subtitle: Stability gate
""",
            ),
            (
                "02-table.yaml",
                """
title: KPI snapshot
table:
  rows: 3
  cols: 2
  has_col_header: true
  col_headers:
    -
      text: CAGR
      sub: "% FY26-33E"
    - Notes
  cells:
    - [North, Stable]
    - [South, Upside]
""",
            ),
        ],
    )

    prs = Presentation(str(output_path))
    assert len(prs.slides) == 2
    assert _chart_types(output_path) == []

    slide2_text = "\n".join(_slide_texts(output_path, 1))
    assert "North" in slide2_text
    assert "South" in slide2_text
    assert "CAGR" in slide2_text


def test_canary_generate_waterfall_chart_cells_with_connectors(tmp_path: Path) -> None:
    output_path = _run_generate(
        tmp_path,
        "canary-waterfall-cells",
        [
            (
                "waterfall-cells.yaml",
                """
title: Waterfall canary
charts:
  wf:
    type: waterfall
    dir: horizontal
    values: [954, 13, -45, 1209]
    totals: [1, 4]
    decreases: [3]
    format: "{:,.0f}"
    connector: true

table:
  rows: 4
  cols: 1
  has_col_header: false
  cells:
    - [wf-1]
    - [wf-2]
    - [wf-3]
    - [wf-4]
""",
            ),
        ],
    )

    assert _chart_types(output_path) == [int(XL_CHART_TYPE.BAR_STACKED)]
    assert _connector_count(output_path, 0) >= 3

    texts = set(_slide_texts(output_path, 0))
    assert {"954", "13", "-45", "1,209"}.issubset(texts)


def test_canary_charts_cli_multi_chart_overlay_deck(tmp_path: Path) -> None:
    output_path = _run_charts(
        tmp_path,
        "canary-charts-multi",
        {
            "layout": "Default",
            "charts": [
                {
                    "type": "clustered",
                    "categories": ["A", "B"],
                    "series": [{"name": "S1", "values": [10, 20], "color": "accent1"}],
                    "show_data_labels": True,
                    "bar": {"orientation": "horizontal"},
                },
                {
                    "type": "waterfall",
                    "categories": ["Start", "Growth", "Costs", "End"],
                    "series": [
                        {
                            "name": "Values",
                            "values": [100, 40, -20, 120],
                            "color": "accent1",
                        }
                    ],
                    "show_data_labels": True,
                    "add_overlay_labels": True,
                    "waterfall": {
                        "orientation": "horizontal",
                        "total_categories": ["Start", "End"],
                        "decrease_categories": ["Costs"],
                        "connector_style": "step",
                        "connector_dash_style": "dot",
                    },
                },
            ],
        },
    )

    prs = Presentation(str(output_path))
    assert len(prs.slides) == 2
    assert _chart_types(output_path) == [
        int(XL_CHART_TYPE.BAR_CLUSTERED),
        int(XL_CHART_TYPE.BAR_STACKED),
    ]

    slide2_text = set(_slide_texts(output_path, 1))
    assert {"Start", "End"}.issubset(slide2_text)

    # Dot connector mode should render many small line-segment rectangles.
    assert _blank_auto_shape_count(output_path, 1) >= 20


def test_canary_charts_cli_template_copy_roundtrip(tmp_path: Path) -> None:
    template_path = _run_charts(
        tmp_path,
        "canary-template-source",
        {
            "type": "clustered",
            "categories": ["A", "B"],
            "series": [{"name": "S1", "values": [1, 2], "color": "accent1"}],
            "show_data_labels": True,
            "data_labels": {"format": "0.00", "font_size": 18},
        },
    )

    target_path = _run_charts(
        tmp_path,
        "canary-template-target",
        {
            "type": "clustered",
            "categories": ["A", "B"],
            "series": [{"name": "S1", "values": [9, 8], "color": "accent2"}],
            "show_data_labels": True,
            "data_labels": {"format": "0", "font_size": 10},
            "bar": {
                "chart_template": str(template_path),
                "chart_template_copy": True,
            },
        },
    )

    assert _pptx_part(target_path, "ppt/charts/chart1.xml") == _pptx_part(
        template_path,
        "ppt/charts/chart1.xml",
    )
    assert _pptx_part(target_path, "ppt/charts/_rels/chart1.xml.rels") == _pptx_part(
        template_path,
        "ppt/charts/_rels/chart1.xml.rels",
    )
